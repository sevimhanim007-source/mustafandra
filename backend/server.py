from fastapi import FastAPI, APIRouter, HTTPException, Depends, BackgroundTasks, status, UploadFile, File, Form, Query, Response
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse, StreamingResponse
from fastapi.encoders import jsonable_encoder
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
import os
import logging
from pathlib import Path
from collections import defaultdict
from pydantic import BaseModel, Field, EmailStr
from typing import List, Optional, Dict, Any, Set, Iterable, Literal, Union
import uuid
from datetime import datetime, timezone, timedelta
import hashlib
import jwt
import smtplib
import ssl
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
import aiofiles
import shutil
import re
from io import BytesIO
import pandas as pd
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import mammoth
from pptx import Presentation
import mimetypes
import math

ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')
DEPLOYMENT_GUIDE_PATH = ROOT_DIR / "DEPLOYMENT.md"

# Create uploads directory
UPLOAD_DIR = ROOT_DIR / "uploads"
UPLOAD_DIR.mkdir(exist_ok=True)

# MongoDB connection
mongo_url = os.environ['MONGO_URL']
client = AsyncIOMotorClient(mongo_url)
db = client[os.environ['DB_NAME']]

# Email configuration
SMTP_SERVER = os.getenv("SMTP_SERVER", "mail.calista.com.tr")
SMTP_PORT = int(os.getenv("SMTP_PORT", 587))
SMTP_USER = os.getenv("SMTP_USER", "mozdemir@calista.com.tr")
SMTP_PASSWORD = os.getenv("SMTP_PASSWORD", "")
SMTP_FROM_NAME = os.getenv("SMTP_FROM_NAME", "QDMS Portal")

# JWT Settings
SECRET_KEY = os.getenv("JWT_SECRET", "your-secret-key-change-this")
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 30
DISABLE_AUTH = os.getenv("DISABLE_AUTH", "false").lower() == "true"
GUEST_USERNAME = os.getenv("GUEST_USERNAME", "guest")
GUEST_FULL_NAME = os.getenv("GUEST_FULL_NAME", "Guest User")
GUEST_EMAIL = os.getenv("GUEST_EMAIL", "guest@example.com")
GUEST_DEPARTMENT = os.getenv("GUEST_DEPARTMENT", "Genel")

# Create the main app
app = FastAPI(title="QDMS Portal", version="1.0.0")
api_router = APIRouter(prefix="/api")

# Static files
app.mount("/uploads", StaticFiles(directory=str(UPLOAD_DIR)), name="uploads")

# Security
security = HTTPBearer(auto_error=not DISABLE_AUTH)

# User Models
class UserRole(BaseModel):
    name: str
    description: Optional[str] = None
    permissions: List[str] = Field(default_factory=list)
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    updated_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

class User(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    username: str
    email: EmailStr
    full_name: str
    role: str
    department: str
    roles: List[str] = Field(default_factory=list)
    groups: List[str] = Field(default_factory=list)
    permissions: List[str] = Field(default_factory=list)
    is_active: bool = True
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    
class UserCreate(BaseModel):
    username: str
    email: EmailStr
    password: str
    full_name: str
    role: str
    department: str
    roles: Optional[List[str]] = None
    groups: Optional[List[str]] = None

class UserLogin(BaseModel):
    username: str
    password: str

class Token(BaseModel):
    access_token: str
    token_type: str
    user: User


class RoleCreate(BaseModel):
    name: str
    description: Optional[str] = None
    permissions: List[str] = Field(default_factory=list)


class RoleUpdate(BaseModel):
    description: Optional[str] = None
    permissions: Optional[List[str]] = None


class UserListItem(BaseModel):
    id: str
    username: str
    email: EmailStr
    full_name: str
    role: str
    roles: List[str]
    department: str
    groups: List[str]
    permissions: List[str]
    is_active: bool
    created_at: datetime


class UserUpdate(BaseModel):
    role: Optional[str] = None
    roles: Optional[List[str]] = None
    department: Optional[str] = None
    groups: Optional[List[str]] = None
    permissions: Optional[List[str]] = None
    is_active: Optional[bool] = None


ADMIN_ROLE_KEYS = {"admin", "system_admin", "systemadministrator"}
ALLOWED_FOLDER_CAPABILITIES = ["read", "download", "create", "revise", "approve", "cancel", "manage"]


def _collect_role_names(data: Dict[str, Any]) -> List[str]:
    roles: Set[str] = set()
    primary_role = data.get("role")
    if primary_role:
        roles.add(str(primary_role))
    for name in data.get("roles") or []:
        if name:
            roles.add(str(name))
    return sorted(roles)


def _is_admin_role(role_name: Optional[str]) -> bool:
    if not role_name:
        return False
    return role_name.lower() in ADMIN_ROLE_KEYS


def user_is_admin(user: User) -> bool:
    if _is_admin_role(user.role):
        return True
    return any(_is_admin_role(role) for role in user.roles)


def user_has_permission(user: User, permission: str) -> bool:
    if user_is_admin(user):
        return True
    if "*" in user.permissions:
        return True
    return permission in user.permissions


def ensure_permission(user: User, permission: str) -> None:
    if not user_has_permission(user, permission):
        raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail="Permission denied")


def _normalize_token(value: Optional[str]) -> str:
    return (value or "").strip().lower()


def _user_identifier_tokens(user: User) -> Set[str]:
    tokens: Set[str] = set()
    tokens.add(_normalize_token(user.id))
    tokens.add(_normalize_token(user.username))
    tokens.add(_normalize_token(user.role))
    if user.email:
        tokens.add(_normalize_token(user.email))
    tokens.update(_normalize_token(role) for role in user.roles)
    if user.department:
        tokens.add(_normalize_token(user.department))
    for group in user.groups or []:
        tokens.add(_normalize_token(group))
    return {token for token in tokens if token}


def user_matches_approver(user: User, approver_token: str) -> bool:
    token = _normalize_token(approver_token)
    if not token:
        return False

    def strip_prefix(value: str) -> str:
        return value.split(":", 1)[1] if ":" in value else ""

    if token.startswith("role:"):
        role_value = strip_prefix(token)
        return any(_normalize_token(role) == role_value for role in _collect_role_names(user.dict()))
    if token.startswith("department:"):
        dept_value = strip_prefix(token)
        return _normalize_token(user.department) == dept_value
    if token.startswith("group:"):
        group_value = strip_prefix(token)
        return any(_normalize_token(group) == group_value for group in user.groups or [])
    if token.startswith("user:"):
        target = strip_prefix(token)
        return target in _user_identifier_tokens(user)

    identifiers = _user_identifier_tokens(user)
    return token in identifiers

# File Models
class FileUpload(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    filename: str
    original_filename: str
    file_path: str
    file_size: int
    mime_type: str
    uploaded_by: str
    uploaded_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    module_type: str  # document, complaint, audit, etc.
    module_id: Optional[str] = None  # related record id (can be linked later)


class FileMetadata(BaseModel):
    id: str
    original_filename: str
    mime_type: str
    file_size: int
    download_url: str


class DocumentPreview(BaseModel):
    kind: str  # html, slides, text, binary
    content: Optional[str] = None
    slides: Optional[List[str]] = None
    mime_type: Optional[str] = None
    download_url: Optional[str] = None


def _generate_document_preview(
    document: Dict[str, Any],
    file_record: Dict[str, Any],
) -> DocumentPreview:
    file_path = Path(file_record["file_path"])
    if not file_path.exists():
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Dosya bulunamadi")

    original_name = file_record.get("original_filename", file_path.name)
    mime_type = file_record.get("mime_type") or mimetypes.guess_type(original_name)[0] or "application/octet-stream"
    extension = file_path.suffix.lower()

    download_url = app.url_path_for("download_file", file_id=file_record["id"])

    try:
        if extension in {".docx"}:
            with open(file_path, "rb") as doc_file:
                result = mammoth.convert_to_html(doc_file)
            html_content = result.value or "<p>Onizleme olusturulamadi.</p>"
            return DocumentPreview(
                kind="html",
                content=html_content,
                mime_type="text/html",
                download_url=download_url,
            )

        if extension in {".xlsx", ".xls"}:
            try:
                sheets = pd.read_excel(file_path, sheet_name=None)
            except Exception:
                sheets = {}
            if not sheets:
                return DocumentPreview(
                    kind="binary",
                    mime_type=mime_type,
                    download_url=download_url,
                )
            parts: List[str] = []
            for sheet_name, df in sheets.items():
                parts.append(f"<h3>{sheet_name}</h3>")
                preview_df = df.head(100)
                parts.append(preview_df.to_html(index=False, escape=False))
            html_content = "".join(parts)
            return DocumentPreview(
                kind="html",
                content=html_content,
                mime_type="text/html",
                download_url=download_url,
            )

        if extension == ".csv":
            try:
                df = pd.read_csv(file_path, nrows=200)
            except Exception:
                df = None
            if df is not None:
                html_content = df.to_html(index=False)
                return DocumentPreview(
                    kind="html",
                    content=html_content,
                    mime_type="text/html",
                    download_url=download_url,
                )

        if extension in {".pptx"}:
            presentation = Presentation(str(file_path))
            slides: List[str] = []
            for slide in presentation.slides:
                slide_text: List[str] = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            slide_text.append(text)
                if slide_text:
                    slides.append("\n".join(slide_text))
            if slides:
                return DocumentPreview(
                    kind="slides",
                    slides=slides,
                    mime_type="text/plain",
                    download_url=download_url,
                )

        if mime_type.startswith("text/"):
            try:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as text_file:
                    content = text_file.read(50000)
            except Exception:
                content = ""
            if content:
                return DocumentPreview(
                    kind="text",
                    content=content,
                    mime_type=mime_type,
                    download_url=download_url,
                )

    except HTTPException:
        raise
    except Exception as exc:
        logging.error("Preview generation failed for %s: %s", original_name, exc)

    return DocumentPreview(
        kind="binary",
        mime_type=mime_type,
        download_url=download_url,
    )


def _build_markdown_blocks(lines: List[str]) -> List[Dict[str, Any]]:
    blocks: List[Dict[str, Any]] = []
    current_block: Optional[Dict[str, Any]] = None
    in_code_block = False

    def flush_block() -> None:
        nonlocal current_block
        if not current_block:
            return
        block_type = current_block["type"]
        if block_type == "paragraph":
            text = " ".join(current_block.get("lines", [])).strip()
            if text:
                blocks.append({"type": "paragraph", "text": text})
        elif block_type in {"unordered_list", "ordered_list"}:
            items = [item.strip() for item in current_block.get("items", []) if item.strip()]
            if items:
                blocks.append({"type": block_type, "items": items})
        elif block_type == "code":
            blocks.append(
                {
                    "type": "code",
                    "lines": current_block.get("lines", []),
                    "language": current_block.get("language"),
                }
            )
        current_block = None

    for raw_line in lines:
        line = raw_line.rstrip("\r\n")
        stripped = line.strip()

        if in_code_block:
            if stripped.startswith("```"):
                flush_block()
                in_code_block = False
            else:
                if current_block:
                    current_block.setdefault("lines", []).append(line)
            continue

        if stripped.startswith("```"):
            flush_block()
            language = stripped[3:].strip() or None
            current_block = {"type": "code", "lines": [], "language": language}
            in_code_block = True
            continue

        if not stripped:
            flush_block()
            continue

        if stripped.startswith("- "):
            if not current_block or current_block.get("type") != "unordered_list":
                flush_block()
                current_block = {"type": "unordered_list", "items": []}
            current_block.setdefault("items", []).append(stripped[2:].strip())
            continue

        ordered_match = re.match(r"(\d+)\.\s+(.*)", stripped)
        if ordered_match:
            if not current_block or current_block.get("type") != "ordered_list":
                flush_block()
                current_block = {"type": "ordered_list", "items": []}
            current_block.setdefault("items", []).append(ordered_match.group(2).strip())
            continue

        if not current_block or current_block.get("type") != "paragraph":
            flush_block()
            current_block = {"type": "paragraph", "lines": []}
        current_block.setdefault("lines", []).append(stripped)

    flush_block()
    return blocks


def _parse_deployment_markdown(content: str) -> Dict[str, Any]:
    parts = content.split("\n## ")
    title = ""
    sections: List[Dict[str, Any]] = []

    if parts:
        raw_title = parts[0].strip()
        if raw_title.startswith("# "):
            title = raw_title[2:].strip()
        else:
            title = raw_title

    for section_raw in parts[1:]:
        segment = section_raw.strip()
        if not segment:
            continue
        lines = segment.splitlines()
        section_title = lines[0].strip()
        body_lines: List[str] = []
        subsections: List[Dict[str, Any]] = []
        current_sub: Optional[Dict[str, Any]] = None

        for line in lines[1:]:
            if line.startswith("### "):
                if current_sub:
                    sub_lines = current_sub.get("body_lines", [])
                    current_sub["body"] = "\n".join(sub_lines).strip() or None
                    current_sub["blocks"] = _build_markdown_blocks(sub_lines)
                    del current_sub["body_lines"]
                    subsections.append(current_sub)
                current_sub = {
                    "title": line[4:].strip(),
                    "body_lines": [],
                }
            else:
                if current_sub is not None:
                    current_sub.setdefault("body_lines", []).append(line)
                else:
                    body_lines.append(line)

        if current_sub:
            sub_lines = current_sub.get("body_lines", [])
            current_sub["body"] = "\n".join(sub_lines).strip() or None
            current_sub["blocks"] = _build_markdown_blocks(sub_lines)
            del current_sub["body_lines"]
            subsections.append(current_sub)

        section_body = "\n".join(body_lines).strip()
        sections.append(
            {
                "title": section_title,
                "body": section_body or None,
                "blocks": _build_markdown_blocks(body_lines),
                "subsections": subsections,
            }
        )

    return {
        "title": title,
        "sections": sections,
    }


class DeploymentGuideBlock(BaseModel):
    type: str
    text: Optional[str] = None
    items: Optional[List[str]] = None
    lines: Optional[List[str]] = None
    language: Optional[str] = None


class DeploymentGuideSubsection(BaseModel):
    title: str
    body: Optional[str] = None
    blocks: List[DeploymentGuideBlock] = Field(default_factory=list)


class DeploymentGuideSection(BaseModel):
    title: str
    body: Optional[str] = None
    blocks: List[DeploymentGuideBlock] = Field(default_factory=list)
    subsections: List[DeploymentGuideSubsection] = Field(default_factory=list)


class DeploymentGuideResponse(BaseModel):
    title: str
    sections: List[DeploymentGuideSection] = Field(default_factory=list)


@api_router.get("/meta/deployment-guide", response_model=DeploymentGuideResponse)
async def get_deployment_guide() -> DeploymentGuideResponse:
    if not DEPLOYMENT_GUIDE_PATH.exists():
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Deployment dokumani bulunamadi")
    try:
        content = DEPLOYMENT_GUIDE_PATH.read_text(encoding="utf-8")
    except Exception as exc:
        logging.error("Deployment dokumani okunamadi: %s", exc)
        raise HTTPException(status_code=status.HTTP_500_INTERNAL_SERVER_ERROR, detail="Deployment dokumani okunamadi")

    parsed = _parse_deployment_markdown(content)
    return DeploymentGuideResponse(**parsed)


# Document Management Models
class FolderPermission(BaseModel):
    principal_type: str = Field(pattern="^(user|role|department|group)$")
    principal_id: str
    capabilities: List[str] = Field(default_factory=list)


class DocumentFolder(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    name: str
    code_prefix: Optional[str] = None
    department: Optional[str] = None
    description: Optional[str] = None
    parent_id: Optional[str] = None
    auto_code_pattern: str = "{PREFIX}-{TYPE}-{SEQ:000}"
    auto_code_seq: int = 0
    permissions: List[FolderPermission] = Field(default_factory=list)
    current_user_capabilities: List[str] = Field(default_factory=list)
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    updated_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))


def _sanitize_capability_name(value: str) -> str:
    return value.lower().strip()


def _normalize_capabilities(capabilities: List[str]) -> List[str]:
    normalized = {_sanitize_capability_name(cap) for cap in capabilities}
    return [cap for cap in ALLOWED_FOLDER_CAPABILITIES if cap in normalized]


def _normalize_string_list(values: Optional[List[str]]) -> List[str]:
    return sorted({value.strip() for value in (values or []) if value and value.strip()})


def _normalize_permissions_list(values: Optional[List[str]]) -> List[str]:
    return sorted({value.strip() for value in (values or []) if value and value.strip()})


def _user_role_names(user: User) -> Set[str]:
    roles: Set[str] = set()
    if user.role:
        roles.add(user.role)
    for role in user.roles:
        if role:
            roles.add(role)
    return roles


def user_has_folder_capability(user: User, folder: DocumentFolder, capability: str) -> bool:
    capability = _sanitize_capability_name(capability)
    if capability not in ALLOWED_FOLDER_CAPABILITIES:
        return False

    if user_is_admin(user):
        return True

    if capability == "manage" and user_has_permission(user, "doc.folder.manage_permissions"):
        return True

    if not folder.permissions:
        # Legacy klasörler için haklar tanımlı değilse tüm kullanıcılar okuyabilir,
        # yönetim hakları ise sadece yetkili kullanıcıya açılsın.
        return capability != "manage"

    user_roles = _user_role_names(user)
    user_departments = {user.department} if user.department else set()
    user_groups = set(user.groups or [])

    for perm in folder.permissions:
        capabilities = _normalize_capabilities(perm.capabilities)
        if capability not in capabilities:
            continue
        principal_id = perm.principal_id
        ptype = perm.principal_type

        if ptype == "role" and principal_id in user_roles:
            return True
        if ptype == "user" and principal_id == user.id:
            return True
        if ptype == "department" and principal_id in user_departments:
            return True
        if ptype == "group" and principal_id in user_groups:
            return True

    if capability in {"read", "download"} and user_has_permission(user, "doc.document.read"):
        return True

    return False


def compute_folder_capabilities_for_user(folder: DocumentFolder, user: User) -> List[str]:
    available = []
    for capability in ALLOWED_FOLDER_CAPABILITIES:
        if user_has_folder_capability(user, folder, capability):
            available.append(capability)
    return available


async def build_user_model(user_data: Dict[str, Any]) -> User:
    data = dict(user_data)
    data.pop("password", None)
    data.setdefault("groups", [])

    roles = _collect_role_names(data)
    data["roles"] = roles

    permissions: Set[str] = set(cap for cap in data.get("permissions") or [] if cap)

    for role_name in roles:
        role_doc = await db.roles.find_one({"name": role_name})
        if role_doc and role_doc.get("permissions"):
            permissions.update(role_doc["permissions"])

    if any(_is_admin_role(role_name) for role_name in roles):
        permissions.add("*")

    data["permissions"] = sorted({perm for perm in permissions if perm})

    if "department" not in data or data["department"] is None:
        data["department"] = ""

    return User(**data)


class DocumentFolderCreate(BaseModel):
    name: str
    code_prefix: Optional[str] = None
    department: Optional[str] = None
    description: Optional[str] = None
    parent_id: Optional[str] = None
    auto_code_pattern: Optional[str] = None
    permissions: List[FolderPermission] = Field(default_factory=list)


class DocumentFolderUpdate(BaseModel):
    name: Optional[str] = None
    code_prefix: Optional[str] = None
    department: Optional[str] = None
    description: Optional[str] = None
    parent_id: Optional[str] = None
    auto_code_pattern: Optional[str] = None


class DocumentFolderPermissionsUpdate(BaseModel):
    permissions: List[FolderPermission] = Field(default_factory=list)


class DocumentDistribution(BaseModel):
    principal_type: str = Field(pattern="^(user|role|department|group)$")
    principal_id: str
    required_to_read: bool = True


class StageDecision(BaseModel):
    user_id: str
    decision: str = Field(pattern="^(approved|rejected)$")
    comment: Optional[str] = None
    decided_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    matched_token: Optional[str] = None


class DocumentApprovalStage(BaseModel):
    stage: int
    approvers: List[str]
    approval_type: str = Field(pattern="^(all|any)$", default="all")
    deadline: Optional[datetime] = None
    status: str = Field(default="pending", pattern="^(pending|approved|rejected)$")
    decided_by: Optional[str] = None
    decided_at: Optional[datetime] = None
    comment: Optional[str] = None
    decisions: List[StageDecision] = Field(default_factory=list)


class DocumentReadReceipt(BaseModel):
    user_id: str
    required: bool = True
    status: str = Field(default="pending", pattern="^(pending|read|overdue)$")
    read_at: Optional[datetime] = None
    note: Optional[str] = None


class DocumentStatus(BaseModel):
    status: str  # draft, review, approved, archived, retired
    changed_by: str
    changed_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    comment: Optional[str] = None


class DocumentVersion(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    version: str
    changes: Optional[str] = None
    notes: Optional[str] = None
    created_by: str
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    file_id: Optional[str] = None
    status: str = Field(default="draft", pattern="^(draft|pending_approval|published|retired)$")


class Document(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    folder_id: str
    code: str
    title: str
    description: Optional[str] = None
    document_type: str  # SOP, Specification, Procedure, Policy
    department: Optional[str] = None
    status: str = Field(default="draft", pattern="^(draft|review|approved|archived|retired)$")
    author_id: str
    version: str = "1.0"
    tags: List[str] = Field(default_factory=list)
    distribution_list: List[DocumentDistribution] = Field(default_factory=list)
    approval_matrix: List[DocumentApprovalStage] = Field(default_factory=list)
    read_receipts: List[DocumentReadReceipt] = Field(default_factory=list)
    status_history: List[DocumentStatus] = Field(default_factory=list)
    version_history: List[DocumentVersion] = Field(default_factory=list)
    current_version_id: Optional[str] = None
    review_date: Optional[datetime] = None
    expiry_date: Optional[datetime] = None
    published_at: Optional[datetime] = None
    archived_at: Optional[datetime] = None
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    updated_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))


class DocumentCreate(BaseModel):
    folder_id: str
    title: str
    document_type: str
    department: Optional[str] = None
    description: Optional[str] = None
    tags: List[str] = Field(default_factory=list)
    distribution_list: List[DocumentDistribution] = Field(default_factory=list)
    approval_matrix: List[DocumentApprovalStage] = Field(default_factory=list)
    file_id: Optional[str] = None
    version_notes: Optional[str] = None
    review_date: Optional[datetime] = None
    expiry_date: Optional[datetime] = None


class DocumentUpdate(BaseModel):
    title: Optional[str] = None
    description: Optional[str] = None
    department: Optional[str] = None
    document_type: Optional[str] = None
    tags: Optional[List[str]] = None
    distribution_list: Optional[List[DocumentDistribution]] = None
    approval_matrix: Optional[List[DocumentApprovalStage]] = None
    review_date: Optional[datetime] = None
    expiry_date: Optional[datetime] = None


class DocumentVersionCreate(BaseModel):
    changes: Optional[str] = None
    notes: Optional[str] = None
    file_id: Optional[str] = None
    mark_as_published: bool = False


class DocumentStatusUpdate(BaseModel):
    status: str
    comment: Optional[str] = None


class DocumentApprovalDecision(BaseModel):
    stage: int
    decision: str = Field(pattern="^(approved|rejected)$")
    comment: Optional[str] = None


class DocumentReadAcknowledge(BaseModel):
    user_id: Optional[str] = None
    note: Optional[str] = None


class DocumentReportFilters(BaseModel):
    department: Optional[str] = None
    status: Optional[str] = None
    document_type: Optional[str] = None
    start_date: Optional[datetime] = None
    end_date: Optional[datetime] = None


class DocumentStatusReport(BaseModel):
    total: int
    status_counts: Dict[str, int]
    department_counts: Dict[str, int]
    type_counts: Dict[str, int]


async def build_document_status_report(
    current_user: User,
    department: Optional[str],
    document_type: Optional[str],
) -> DocumentStatusReport:
    query: Dict[str, Any] = {}
    if department:
        query["department"] = department
    if document_type:
        query["document_type"] = document_type

    allowed_folders = await get_authorized_folder_ids(current_user, "read")
    if allowed_folders is not None:
        if not allowed_folders:
            return DocumentStatusReport(
                total=0,
                status_counts={},
                department_counts={},
                type_counts={},
            )
        query["folder_id"] = {"$in": list(allowed_folders)}

    total = await db.documents.count_documents(query)

    status_cursor = db.documents.aggregate(
        [
            {"$match": query},
            {"$group": {"_id": "$status", "count": {"$sum": 1}}},
        ]
    )
    status_counts = {
        (item["_id"] or "unknown"): item["count"]
        for item in await status_cursor.to_list(length=None)
    }

    department_cursor = db.documents.aggregate(
        [
            {"$match": query},
            {"$group": {"_id": "$department", "count": {"$sum": 1}}},
        ]
    )
    department_counts = {
        (item["_id"] or "Bilinmiyor"): item["count"]
        for item in await department_cursor.to_list(length=None)
    }

    type_cursor = db.documents.aggregate(
        [
            {"$match": query},
            {"$group": {"_id": "$document_type", "count": {"$sum": 1}}},
        ]
    )
    type_counts = {
        (item["_id"] or "Bilinmiyor"): item["count"]
        for item in await type_cursor.to_list(length=None)
    }

    return DocumentStatusReport(
        total=total,
        status_counts=status_counts,
        department_counts=department_counts,
        type_counts=type_counts,
    )


class DocumentApprovalTask(BaseModel):
    document_id: str
    document_code: str
    title: str
    version: str
    stage: int
    approval_type: str
    deadline: Optional[datetime] = None
    folder_id: str
    status: str
    approvers: List[str]


class DocumentReadTask(BaseModel):
    document_id: str
    document_code: str
    title: str
    version: str
    required: bool
    status: str
    folder_id: str


DOCUMENT_CODE_PATTERN = "{PREFIX}-{TYPE}-{SEQ:000}"
DOCUMENT_SEQ_REGEX = re.compile(r"\{SEQ:(\d+)\}")


def _sanitize_code_token(value: Optional[str], fallback: str) -> str:
    token = (value or fallback).upper()
    token = re.sub(r"[^A-Z0-9]", "", token)
    return token or fallback


async def ensure_document_folder(
    folder_id: str,
    user: Optional[User] = None,
    capability: Optional[str] = None,
) -> DocumentFolder:
    folder_data = await db.document_folders.find_one({"id": folder_id})
    if not folder_data:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Document folder not found")
    folder = DocumentFolder(**folder_data)
    if user:
        folder.current_user_capabilities = compute_folder_capabilities_for_user(folder, user)
        if capability and not user_has_folder_capability(user, folder, capability):
            raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail="Folder access denied")
    elif capability:
        # capability istendi ama kullanıcı verilmediyse kontrol yapılamaz
        raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail="Folder access denied")
    return folder


async def get_authorized_folder_ids(user: User, capability: str) -> Optional[Set[str]]:
    capability = _sanitize_capability_name(capability)
    if user_is_admin(user):
        return None

    allowed: Set[str] = set()
    cursor = db.document_folders.find({})
    async for folder_data in cursor:
        folder = DocumentFolder(**folder_data)
        if user_has_folder_capability(user, folder, capability):
            allowed.add(folder.id)
    return allowed


def sanitize_folder_permissions(permissions: List[FolderPermission]) -> List[Dict[str, Any]]:
    sanitized: List[Dict[str, Any]] = []
    for perm in permissions:
        capabilities = _normalize_capabilities(perm.capabilities)
        principal_id = perm.principal_id.strip()
        if not capabilities or not principal_id:
            continue
        sanitized.append(
            {
                "principal_type": perm.principal_type,
                "principal_id": principal_id,
                "capabilities": capabilities,
            }
        )
    return sanitized


async def generate_document_code(folder: DocumentFolder, document_type: str) -> str:
    pattern = folder.auto_code_pattern or DOCUMENT_CODE_PATTERN
    prefix = _sanitize_code_token(folder.code_prefix, "DOC")
    doc_type_token = _sanitize_code_token(document_type[:8], "DOC")
    now = datetime.now(timezone.utc)

    next_seq = folder.auto_code_seq + 1

    def replace_seq(match: re.Match) -> str:
        width = int(match.group(1))
        return str(next_seq).zfill(width)

    code = DOCUMENT_SEQ_REGEX.sub(replace_seq, pattern)
    code = code.replace("{PREFIX}", prefix)
    code = code.replace("{TYPE}", doc_type_token)
    code = code.replace("{YEAR}", str(now.year))
    code = code.replace("{MONTH}", f"{now.month:02d}")

    await db.document_folders.update_one(
        {"id": folder.id},
        {"$set": {"auto_code_seq": next_seq, "updated_at": now}},
    )

    return code


def build_initial_read_receipts(distribution: List[DocumentDistribution]) -> List[DocumentReadReceipt]:
    receipts: List[DocumentReadReceipt] = []
    for item in distribution:
        if item.principal_type == "user":
            receipts.append(
                DocumentReadReceipt(
                    user_id=item.principal_id,
                    required=item.required_to_read,
                )
            )
    return receipts


def normalize_approval_stages(stages: List[DocumentApprovalStage]) -> List[DocumentApprovalStage]:
    normalized: List[DocumentApprovalStage] = []
    for stage in stages:
        if stage.stage is None:
            continue
        stage.approvers = [token for token in stage.approvers if token]
        stage.decisions = [
            StageDecision(**decision.dict()) if not isinstance(decision, StageDecision) else decision
            for decision in stage.decisions or []
        ]
        normalized.append(stage)
    normalized.sort(key=lambda entry: entry.stage)
    return normalized


def reset_approval_stages(stages: List[DocumentApprovalStage]) -> List[DocumentApprovalStage]:
    reset_list: List[DocumentApprovalStage] = []
    for stage in stages:
        reset_list.append(
            DocumentApprovalStage(
                stage=stage.stage,
                approvers=list(stage.approvers),
                approval_type=stage.approval_type,
                deadline=stage.deadline,
                status="pending",
                comment=None,
                decided_at=None,
                decided_by=None,
                decisions=[],
            )
        )
    return reset_list


def find_pending_stage_index(document: Document) -> Optional[int]:
    pending_index: Optional[int] = None
    pending_stage_value: Optional[int] = None
    for idx, stage in enumerate(document.approval_matrix):
        if stage.status != "pending":
            continue
        if pending_stage_value is None or stage.stage < pending_stage_value:
            pending_stage_value = stage.stage
            pending_index = idx
    return pending_index


def approval_stage_is_complete(stage: DocumentApprovalStage) -> bool:
    if stage.status == "approved":
        return True
    if stage.approval_type == "any":
        return any(decision.decision == "approved" for decision in stage.decisions)
    if not stage.approvers:
        return False
    approved_tokens = {
        _normalize_token(decision.matched_token or decision.user_id)
        for decision in stage.decisions
        if decision.decision == "approved"
    }
    required_tokens = {_normalize_token(token) for token in stage.approvers if token}
    return required_tokens.issubset(approved_tokens)


def resolve_matching_token(user: User, stage: DocumentApprovalStage) -> Optional[str]:
    for token in stage.approvers:
        if user_matches_approver(user, token):
            return token
    return None


def _collect_tokens_from_user_doc(user_doc: Dict[str, Any]) -> Set[str]:
    tokens: Set[str] = set()
    identifiers = [
        user_doc.get("id"),
        user_doc.get("_id"),
        user_doc.get("username"),
        user_doc.get("email"),
        user_doc.get("role"),
        user_doc.get("department"),
    ]
    for identifier in identifiers:
        if identifier:
            tokens.add(_normalize_token(str(identifier)))

    for role_name in user_doc.get("roles") or []:
        if role_name:
            tokens.add(_normalize_token(role_name))

    for group_name in user_doc.get("groups") or []:
        if group_name:
            tokens.add(_normalize_token(group_name))

    return {token for token in tokens if token}


def _user_doc_matches_token(user_doc: Dict[str, Any], token: str) -> bool:
    normalized = _normalize_token(token)
    if not normalized:
        return False

    tokens = _collect_tokens_from_user_doc(user_doc)

    def matches_identifier(target: str) -> bool:
        return target in tokens

    if normalized.startswith("user:"):
        target = normalized.split(":", 1)[1]
        return matches_identifier(target)
    if normalized.startswith("role:"):
        target = normalized.split(":", 1)[1]
        role_tokens = {_normalize_token(user_doc.get("role"))}
        role_tokens.update(_normalize_token(role) for role in user_doc.get("roles") or [])
        return target in {token for token in role_tokens if token}
    if normalized.startswith("department:"):
        target = normalized.split(":", 1)[1]
        return target == _normalize_token(user_doc.get("department"))
    if normalized.startswith("group:"):
        target = normalized.split(":", 1)[1]
        return target in {_normalize_token(group) for group in user_doc.get("groups") or [] if group}

    return normalized in tokens


async def collect_user_ids_for_tokens(tokens: List[str]) -> Set[str]:
    normalized_tokens = [_normalize_token(token) for token in tokens or [] if token]
    normalized_tokens = [token for token in normalized_tokens if token]
    if not normalized_tokens:
        return set()

    users: List[Dict[str, Any]] = await db.users.find({"is_active": True}).to_list(length=None)
    matched: Set[str] = set()
    for user_doc in users:
        user_id = user_doc.get("id") or str(user_doc.get("_id") or "")
        if not user_id:
            continue
        for token in normalized_tokens:
            if _user_doc_matches_token(user_doc, token):
                matched.add(str(user_id))
                break
    return matched


async def notify_users(user_ids: Iterable[str], title: str, message: str, notif_type: str = "info") -> None:
    unique_ids = {str(user_id) for user_id in user_ids if user_id}
    if not unique_ids:
        return
    notifications = [
        Notification(user_id=user_id, title=title, message=message, type=notif_type).dict()
        for user_id in unique_ids
    ]
    if notifications:
        await db.notifications.insert_many(notifications)


async def notify_document_approvers(
    document: Document,
    stage: DocumentApprovalStage,
    *,
    exclude_user_ids: Optional[Iterable[str]] = None,
) -> None:
    recipients = await collect_user_ids_for_tokens(stage.approvers)
    if exclude_user_ids:
        recipients -= {str(user_id) for user_id in exclude_user_ids if user_id}
    if not recipients:
        return
    await notify_users(
        recipients,
        title="Dokuman Onayi Bekliyor",
        message=f"{document.code} / {document.title} icin {stage.stage}. asamada onayiniz gerekiyor.",
        notif_type="warning",
    )


async def notify_document_readers(
    document: Document,
    user_ids: Optional[Iterable[str]] = None,
    note: Optional[str] = None,
) -> None:
    recipients = (
        {str(user_id) for user_id in user_ids or [] if user_id}
        if user_ids is not None
        else {
            receipt.user_id
            for receipt in document.read_receipts
            if receipt.user_id and receipt.status != "read"
        }
    )
    recipients = {str(user_id) for user_id in recipients if user_id}
    if not recipients:
        return
    await notify_users(
        recipients,
        title="Yeni Dokuman Okuma Onayi",
        message=note or f"{document.code} / {document.title} icin okuma onayiniz bekleniyor.",
        notif_type="info",
    )


def hydrate_document(document_raw: Dict[str, Any]) -> Document:
    data = dict(document_raw)
    data.setdefault("folder_id", "LEGACY")
    data.setdefault("distribution_list", [])
    data.setdefault("approval_matrix", [])
    data.setdefault("read_receipts", [])
    data.setdefault("status_history", [])
    data.setdefault("version_history", [])
    data.setdefault("current_version_id", None)
    data.setdefault("version", "1.0")
    data.setdefault("tags", [])
    data.setdefault("created_at", datetime.now(timezone.utc))
    data.setdefault("updated_at", datetime.now(timezone.utc))
    return Document(**data)


def determine_next_version(version_history: List[Dict[str, Any]], mark_as_published: bool) -> str:
    if not version_history:
        return "1.0"
    latest = version_history[-1].get("version", "1.0")
    try:
        major_str, minor_str = latest.split(".")
        major = int(major_str)
        minor = int(minor_str)
    except ValueError:
        major = 1
        minor = len(version_history)

    if mark_as_published:
        major += 1
        minor = 0
    else:
        minor += 1
    return f"{major}.{minor}"

# Customer Complaint Models
class ComplaintStatus(BaseModel):
    status: str
    changed_by: str
    changed_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    comment: Optional[str] = None

class ComplaintCategory(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    name: str
    description: Optional[str] = None
    is_active: bool = True
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    updated_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))


class ComplaintCategoryCreate(BaseModel):
    name: str
    description: Optional[str] = None


class ComplaintCategoryUpdate(BaseModel):
    name: Optional[str] = None
    description: Optional[str] = None
    is_active: Optional[bool] = None


class Complaint(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    complaint_no: str
    customer_name: str
    customer_contact: str
    complaint_date: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    complaint_type: str  # product, service, delivery, etc.
    priority: str  # low, medium, high, critical
    status: str = "open"  # open, investigating, resolved, closed
    description: str
    department: str
    assigned_to: Optional[str] = None
    team_leader: Optional[str] = None
    solution_team: List[str] = Field(default_factory=list)
    initial_response: Optional[str] = None
    investigation_report: Optional[str] = None
    investigation_reported_at: Optional[datetime] = None
    investigation_reported_by: Optional[str] = None
    related_task_ids: List[str] = Field(default_factory=list)
    related_capa_ids: List[str] = Field(default_factory=list)
    final_response: Optional[str] = None
    final_report: Optional[str] = None
    final_reported_at: Optional[datetime] = None
    final_reported_by: Optional[str] = None
    resolution_date: Optional[datetime] = None
    file_attachments: List[str] = Field(default_factory=list)
    status_history: List[ComplaintStatus] = Field(default_factory=list)
    created_by: str
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    updated_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    category_id: Optional[str] = None
    category_name: Optional[str] = None

class ComplaintCreate(BaseModel):
    customer_name: str
    customer_contact: str
    complaint_type: str
    priority: str
    description: str
    department: str
    assigned_to: Optional[str] = None
    team_leader: Optional[str] = None
    solution_team: Optional[List[str]] = None
    initial_response: Optional[str] = None
    file_attachments: Optional[List[str]] = None
    related_task_ids: Optional[List[str]] = None
    related_capa_ids: Optional[List[str]] = None
    category_id: Optional[str] = None


class ComplaintAssignmentUpdate(BaseModel):
    assigned_to: Optional[str] = None
    team_leader: Optional[str] = None
    solution_team: Optional[List[str]] = None
    initial_response: Optional[str] = None


class ComplaintInvestigationUpdate(BaseModel):
    investigation_report: str
    investigation_reported_at: Optional[datetime] = None
    related_task_ids: Optional[List[str]] = None
    file_attachments: Optional[List[str]] = None


class ComplaintFinalizationUpdate(BaseModel):
    final_report: str
    final_response: Optional[str] = None
    file_attachments: Optional[List[str]] = None
    mark_resolved: bool = False


class ComplaintMetadataUpdate(BaseModel):
    complaint_type: Optional[str] = None
    priority: Optional[str] = None
    department: Optional[str] = None
    category_id: Optional[str] = None


class ComplaintAttachmentAdd(BaseModel):
    file_ids: List[str]


class ComplaintDofCreate(BaseModel):
    title: Optional[str] = None
    description: Optional[str] = None
    department: Optional[str] = None
    responsible_person: Optional[str] = None
    due_date: Optional[datetime] = None
    team_members: Optional[List[str]] = None
    initial_improvement_report_date: Optional[datetime] = None


class ComplaintCapaCreate(BaseModel):
    title: Optional[str] = None
    source: Optional[str] = None
    department: Optional[str] = None
    team_leader: Optional[str] = None
    target_date: Optional[datetime] = None
    nonconformity_description: Optional[str] = None
    file_attachments: Optional[List[str]] = None
    team_members: Optional[List[str]] = None
    initial_improvement_report_date: Optional[datetime] = None
    linked_risk_ids: Optional[List[str]] = None
    linked_equipment_ids: Optional[List[str]] = None
    linked_audit_finding_ids: Optional[List[str]] = None

# CAPA Models
class CapaAction(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    action_description: str
    responsible_person: str
    due_date: datetime
    status: str = "open"  # open, in_progress, completed, overdue
    completion_date: Optional[datetime] = None
    evidence: Optional[str] = None

ALLOWED_CAPA_STATUSES = {"open", "investigating", "implementing", "pending_closure", "closed", "cancelled"}
ALLOWED_AUDIT_STATUSES = {"planned", "in_progress", "completed", "cancelled"}
ALLOWED_AUDIT_FINDING_TYPES = {"observation", "minor", "major", "critical"}
ALLOWED_AUDIT_FINDING_STATUSES = {"open", "in_progress", "closed"}
ALLOWED_AUDIT_CHECKLIST_STATUSES = {"pending", "compliant", "noncompliant", "not_applicable"}
ALLOWED_RISK_STATUSES = {"identified", "assessed", "mitigating", "monitoring", "closed"}
DEFAULT_RISK_LEVELS = ["low", "medium", "high", "critical"]
DEFAULT_RISK_MATRIX_SIZE = 5
CALIBRATION_DEVICE_STATUSES = {"active", "inactive", "out_of_service"}
CALIBRATION_WORK_ORDER_STATUSES = {"planned", "in_progress", "completed", "cancelled"}
CALIBRATION_DEVICE_TRANSITIONS = {
    "active": {"inactive", "out_of_service"},
    "inactive": {"active", "out_of_service"},
    "out_of_service": {"active", "inactive"},
}
CALIBRATION_WORK_ORDER_TRANSITIONS = {
    "planned": {"in_progress", "cancelled"},
    "in_progress": {"completed", "cancelled"},
    "completed": set(),
    "cancelled": set(),
}
REPORT_ALLOWED_MODULES = {
    "documents",
    "document_versions",
    "calibration_devices",
    "calibration_work_orders",
    "risks",
    "capas",
    "complaints",
    "dof_tasks",
    "audit_findings",
}
REPORT_ALLOWED_OPERATORS = {
    "eq",
    "ne",
    "in",
    "nin",
    "gte",
    "lte",
    "gt",
    "lt",
    "between",
    "contains",
    "regex",
    "exists",
}
SYSTEM_USER_ID = "system"

SEED_REPORT_TEMPLATES: List[Dict[str, Any]] = [
    {
        "id": "calibration-device-overview-html",
        "name": "Calibration Device Overview (HTML)",
        "module": "calibration_devices",
        "template_type": "html",
        "description": "Baseline HTML template for calibration device overview reports.",
        "default_definition_id": "calibration-device-overview",
        "placeholders": [
            {"key": "report_title", "source": "header", "description": "Title for the rendered report."},
            {"key": "generated_at", "source": "header", "description": "Generation timestamp."},
            {"key": "summary.total_rows", "source": "summary", "description": "Total devices in scope."},
            {
                "key": "visualizations.status.series",
                "source": "visualizations",
                "path": "1.series",
                "description": "Series data for status distribution chart.",
            },
            {"key": "rows", "source": "table", "description": "Tabular device dataset."},
        ],
    },
    {
        "id": "calibration-work-order-summary-html",
        "name": "Calibration Work Order Summary (HTML)",
        "module": "calibration_work_orders",
        "template_type": "html",
        "description": "HTML template for calibration work order pipeline reports.",
        "default_definition_id": "calibration-work-order-pipeline",
        "placeholders": [
            {"key": "report_title", "source": "header", "description": "Title for the rendered report."},
            {"key": "summary.total_rows", "source": "summary", "description": "Total work orders returned."},
            {
                "key": "visualizations.status.series",
                "source": "visualizations",
                "path": "1.series",
                "description": "Status distribution for work orders.",
            },
            {
                "key": "visualizations.due.chart",
                "source": "visualizations",
                "path": "2.series",
                "description": "Due date trend data.",
            },
            {"key": "rows", "source": "table", "description": "Work order table rows."},
        ],
    },
    {
        "id": "complaint-status-summary-html",
        "name": "Complaint Status Summary (HTML)",
        "module": "complaints",
        "template_type": "html",
        "description": "HTML template for complaint lifecycle reporting.",
        "default_definition_id": "complaint-status-summary",
        "placeholders": [
            {"key": "report_title", "source": "header", "description": "Report title placeholder."},
            {"key": "summary.total_rows", "source": "summary", "description": "Total complaints in report."},
            {
                "key": "visualizations.status.series",
                "source": "visualizations",
                "path": "1.series",
                "description": "Complaint status distribution.",
            },
            {"key": "rows", "source": "table", "description": "Complaint listing table."},
        ],
    },
]

SEED_REPORT_DEFINITIONS: List[Dict[str, Any]] = [
    {
        "id": "calibration-device-overview",
        "name": "Calibration Device Overview",
        "module": "calibration_devices",
        "description": "Track calibration devices with status distribution and key metrics.",
        "filters": [
            {
                "field": "status",
                "label": "Status",
                "operator": "in",
                "field_type": "string",
                "allow_multiple": True,
            },
            {
                "field": "department",
                "label": "Department",
                "operator": "in",
                "field_type": "string",
                "allow_multiple": True,
            },
            {
                "field": "category",
                "label": "Category",
                "operator": "in",
                "field_type": "string",
                "allow_multiple": True,
            },
        ],
        "visualizations": [
            {
                "type": "kpi",
                "title": "Device KPIs",
                "metrics": [
                    {"operation": "count", "label": "Total Devices"},
                    {
                        "operation": "avg",
                        "field": "calibration_interval_days",
                        "label": "Avg Interval (days)",
                        "precision": 1,
                    },
                ],
            },
            {
                "type": "bar",
                "title": "Devices by Status",
                "group_by": "status",
                "metrics": [{"operation": "count", "label": "Devices"}],
            },
            {
                "type": "table",
                "title": "Device Details",
                "fields": [
                    "device_code",
                    "name",
                    "status",
                    "department",
                    "category",
                    "calibration_interval_days",
                    "next_due_date",
                    "responsible_person",
                ],
            },
        ],
        "fields": [
            "device_code",
            "name",
            "status",
            "department",
            "category",
            "calibration_interval_days",
            "last_calibrated_at",
            "next_due_date",
            "responsible_person",
        ],
        "max_rows": 5000,
        "default_template_id": "calibration-device-overview-html",
    },
    {
        "id": "calibration-work-order-pipeline",
        "name": "Calibration Work Order Pipeline",
        "module": "calibration_work_orders",
        "description": "Understand work order load, due dates, and assignments for calibration activities.",
        "filters": [
            {
                "field": "status",
                "label": "Status",
                "operator": "in",
                "field_type": "string",
                "allow_multiple": True,
            },
            {
                "field": "assigned_to",
                "label": "Assigned To",
                "operator": "in",
                "field_type": "string",
                "allow_multiple": True,
            },
            {
                "field": "due_date",
                "label": "Due Date Between",
                "operator": "between",
                "field_type": "datetime",
            },
        ],
        "visualizations": [
            {
                "type": "kpi",
                "title": "Work Order KPIs",
                "metrics": [
                    {"operation": "count", "label": "Total Work Orders"},
                ],
            },
            {
                "type": "bar",
                "title": "Work Orders by Status",
                "group_by": "status",
                "metrics": [{"operation": "count", "label": "Work Orders"}],
            },
            {
                "type": "line",
                "title": "Work Orders by Due Date",
                "group_by": "due_date",
                "metrics": [{"operation": "count", "label": "Due"}],
            },
        ],
        "fields": [
            "work_order_no",
            "device_id",
            "status",
            "planned_date",
            "due_date",
            "assigned_to",
            "result",
            "linked_dof_task_id",
        ],
        "max_rows": 5000,
        "default_template_id": "calibration-work-order-summary-html",
    },
    {
        "id": "complaint-status-summary",
        "name": "Complaint Status Summary",
        "module": "complaints",
        "description": "Monitor complaint lifecycle with status counts and departmental distribution.",
        "filters": [
            {
                "field": "status",
                "label": "Status",
                "operator": "in",
                "field_type": "string",
                "allow_multiple": True,
            },
            {
                "field": "department",
                "label": "Department",
                "operator": "in",
                "field_type": "string",
                "allow_multiple": True,
            },
            {
                "field": "priority",
                "label": "Priority",
                "operator": "in",
                "field_type": "string",
                "allow_multiple": True,
            },
            {
                "field": "created_at",
                "label": "Created Between",
                "operator": "between",
                "field_type": "datetime",
            },
        ],
        "visualizations": [
            {
                "type": "kpi",
                "title": "Complaint KPIs",
                "metrics": [
                    {"operation": "count", "label": "Total Complaints"},
                ],
            },
            {
                "type": "pie",
                "title": "Complaints by Status",
                "group_by": "status",
                "metrics": [{"operation": "count", "label": "Complaints"}],
            },
            {
                "type": "bar",
                "title": "Complaints by Department",
                "group_by": "department",
                "metrics": [{"operation": "count", "label": "Complaints"}],
            },
        ],
        "fields": [
            "complaint_no",
            "title",
            "status",
            "department",
            "priority",
            "responsible_person",
            "created_at",
            "due_date",
        ],
        "max_rows": 5000,
        "default_template_id": "complaint-status-summary-html",
    },
]


class Capa(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    capa_no: str
    title: str
    source: str  # internal_audit, customer_complaint, management_review, etc.
    department: str
    initiated_by: str
    team_leader: str
    initiated_date: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    target_date: Optional[datetime] = None
    status: str = "open"  # open, investigating, implementing, pending_closure, closed, cancelled
    nonconformity_description: str
    root_cause_analysis: Optional[str] = None
    immediate_action: Optional[str] = None
    corrective_actions: List[CapaAction] = Field(default_factory=list)
    preventive_actions: List[CapaAction] = Field(default_factory=list)
    effectiveness_review: Optional[str] = None
    file_attachments: List[str] = Field(default_factory=list)
    team_members: List[str] = Field(default_factory=list)
    initial_improvement_report_date: Optional[datetime] = None
    linked_risk_ids: List[str] = Field(default_factory=list)
    linked_equipment_ids: List[str] = Field(default_factory=list)
    linked_audit_finding_ids: List[str] = Field(default_factory=list)
    closure_requested_at: Optional[datetime] = None
    closure_requested_by: Optional[str] = None
    closure_request_note: Optional[str] = None
    closure_approved_at: Optional[datetime] = None
    closure_approved_by: Optional[str] = None
    closure_decision_note: Optional[str] = None
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    updated_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

class CapaCreate(BaseModel):
    title: str
    source: str
    department: str
    team_leader: str
    target_date: Optional[datetime] = None
    nonconformity_description: str
    team_members: Optional[List[str]] = None
    initial_improvement_report_date: Optional[datetime] = None
    linked_risk_ids: Optional[List[str]] = None
    linked_equipment_ids: Optional[List[str]] = None
    linked_audit_finding_ids: Optional[List[str]] = None


class CapaUpdate(BaseModel):
    title: Optional[str] = None
    source: Optional[str] = None
    department: Optional[str] = None
    team_leader: Optional[str] = None
    target_date: Optional[datetime] = None
    status: Optional[str] = None
    nonconformity_description: Optional[str] = None
    root_cause_analysis: Optional[str] = None
    immediate_action: Optional[str] = None
    effectiveness_review: Optional[str] = None
    file_attachments: Optional[List[str]] = None
    team_members: Optional[List[str]] = None
    initial_improvement_report_date: Optional[datetime] = None
    linked_risk_ids: Optional[List[str]] = None
    linked_equipment_ids: Optional[List[str]] = None
    linked_audit_finding_ids: Optional[List[str]] = None


class CapaClosureRequest(BaseModel):
    note: Optional[str] = None


class CapaClosureDecision(BaseModel):
    approve: bool
    note: Optional[str] = None


class CapaActionCreate(BaseModel):
    action_description: str
    responsible_person: str
    due_date: datetime
    action_type: str = Field(pattern="^(corrective|preventive)$")
    status: Optional[str] = "open"
    evidence: Optional[str] = None


class CapaActionUpdate(BaseModel):
    action_description: Optional[str] = None
    responsible_person: Optional[str] = None
    due_date: Optional[datetime] = None
    status: Optional[str] = None
    completion_date: Optional[datetime] = None
    evidence: Optional[str] = None


class AuditTeamMember(BaseModel):
    user_id: str
    role: str
    full_name: Optional[str] = None


class AuditChecklistItem(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    question_id: Optional[str] = None
    question: Optional[str] = None
    response: Optional[str] = None
    status: str = "pending"
    note: Optional[str] = None


class AuditStatusHistory(BaseModel):
    status: str
    changed_by: str
    changed_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    note: Optional[str] = None


class AuditFinding(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    finding_type: str
    description: str
    requirement_reference: Optional[str] = None
    related_capa_id: Optional[str] = None
    status: str = "open"
    corrective_action: Optional[str] = None
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    updated_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))


class Audit(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    audit_code: str
    audit_type: str
    scope: str
    department: str
    start_date: datetime
    end_date: datetime
    status: str = "planned"
    lead_auditor: str
    audit_team: List[AuditTeamMember] = Field(default_factory=list)
    auditee_representative: Optional[str] = None
    objectives: Optional[str] = None
    checklist: List[AuditChecklistItem] = Field(default_factory=list)
    findings: List[AuditFinding] = Field(default_factory=list)
    status_history: List[AuditStatusHistory] = Field(default_factory=list)
    created_by: str
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    updated_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))


class AuditCreate(BaseModel):
    audit_code: Optional[str] = None
    audit_type: str
    scope: str
    department: str
    start_date: datetime
    end_date: datetime
    lead_auditor: str
    audit_team: Optional[List[AuditTeamMember]] = None
    auditee_representative: Optional[str] = None
    objectives: Optional[str] = None
    checklist: Optional[List[AuditChecklistItem]] = None


class AuditUpdate(BaseModel):
    audit_type: Optional[str] = None
    scope: Optional[str] = None
    department: Optional[str] = None
    start_date: Optional[datetime] = None
    end_date: Optional[datetime] = None
    lead_auditor: Optional[str] = None
    audit_team: Optional[List[AuditTeamMember]] = None
    auditee_representative: Optional[str] = None
    objectives: Optional[str] = None
    checklist: Optional[List[AuditChecklistItem]] = None
    status: Optional[str] = None


class AuditStatusUpdate(BaseModel):
    status: str
    note: Optional[str] = None


class AuditFindingCreate(BaseModel):
    finding_type: str
    description: str
    requirement_reference: Optional[str] = None
    related_capa_id: Optional[str] = None
    corrective_action: Optional[str] = None


class AuditFindingUpdate(BaseModel):
    finding_type: Optional[str] = None
    description: Optional[str] = None
    requirement_reference: Optional[str] = None
    related_capa_id: Optional[str] = None
    status: Optional[str] = None
    corrective_action: Optional[str] = None


class AuditQuestion(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    question: str
    category: Optional[str] = None
    requirement_reference: Optional[str] = None
    tags: List[str] = Field(default_factory=list)
    is_active: bool = True
    created_by: str
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    updated_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))


class AuditQuestionCreate(BaseModel):
    question: str
    category: Optional[str] = None
    requirement_reference: Optional[str] = None
    tags: Optional[List[str]] = None
    is_active: Optional[bool] = True


class AuditQuestionUpdate(BaseModel):
    question: Optional[str] = None
    category: Optional[str] = None
    requirement_reference: Optional[str] = None
    tags: Optional[List[str]] = None
    is_active: Optional[bool] = None


class RiskFactor(BaseModel):
    name: str
    value: float
    weight: float = 1.0


class RiskScore(BaseModel):
    inherent: float
    residual: float
    inherent_level: str
    residual_level: str
    matrix_row: int
    matrix_col: int


class RiskRevision(BaseModel):
    revision_no: int
    snapshot: Dict[str, Any]
    changed_by: str
    changed_at: datetime
    note: Optional[str] = None


class RiskTrendPoint(BaseModel):
    recorded_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    inherent_score: float
    residual_score: float
    status: str


class RiskAssessment(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    risk_code: str
    title: str
    category: str
    process: Optional[str] = None
    owner: str
    description: Optional[str] = None
    status: str = "identified"
    likelihood: float
    impact: float
    detection: Optional[float] = None
    controls_effectiveness: Optional[float] = None
    custom_factors: List[RiskFactor] = Field(default_factory=list)
    risk_score: RiskScore
    linked_capa_ids: List[str] = Field(default_factory=list)
    linked_audit_finding_ids: List[str] = Field(default_factory=list)
    next_review_date: Optional[datetime] = None
    last_reviewed_at: Optional[datetime] = None
    created_by: str
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    updated_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    revisions: List[RiskRevision] = Field(default_factory=list)
    trend: List[RiskTrendPoint] = Field(default_factory=list)


class RiskMatrixCell(BaseModel):
    row: int
    col: int
    level: str
    color: str
    label: str
    count: Optional[int] = None


class RiskMatrixSummary(BaseModel):
    matrix: List[List[Dict[str, Any]]]
    palette: Dict[str, str]


class RiskTrendResponse(BaseModel):
    points: List[RiskTrendPoint]


class RiskSettings(BaseModel):
    id: str = "default"
    formula: str = "likelihood * impact"
    residual_formula: Optional[str] = "likelihood * impact * (1 - controls_effectiveness)"
    max_scale: int = 25
    thresholds: Dict[str, float] = Field(
        default_factory=lambda: {
            "low": 5,
            "medium": 12,
            "high": 20,
            "critical": 25,
        }
    )
    matrix: List[List[str]] = Field(
        default_factory=lambda: [
            ["low", "low", "medium", "medium", "high"],
            ["low", "medium", "medium", "high", "high"],
            ["medium", "medium", "high", "high", "critical"],
            ["medium", "high", "high", "critical", "critical"],
            ["high", "high", "critical", "critical", "critical"],
        ]
    )
    palette: Dict[str, str] = Field(
        default_factory=lambda: {
            "low": "#90EE90",
            "medium": "#F9E076",
            "high": "#F39C12",
            "critical": "#E74C3C",
        }
    )
    updated_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    updated_by: Optional[str] = None


class RiskAssessmentCreate(BaseModel):
    title: str
    category: str
    process: Optional[str] = None
    owner: str
    description: Optional[str] = None
    status: Optional[str] = "identified"
    likelihood: float = Field(ge=1, le=5)
    impact: float = Field(ge=1, le=5)
    detection: Optional[float] = Field(default=None, ge=0, le=5)
    controls_effectiveness: Optional[float] = Field(default=None, ge=0, le=1)
    custom_factors: Optional[List[RiskFactor]] = None
    linked_capa_ids: Optional[List[str]] = None
    linked_audit_finding_ids: Optional[List[str]] = None
    next_review_date: Optional[datetime] = None
    revision_note: Optional[str] = None


class RiskAssessmentUpdate(BaseModel):
    title: Optional[str] = None
    category: Optional[str] = None
    process: Optional[str] = None
    owner: Optional[str] = None
    description: Optional[str] = None
    status: Optional[str] = None
    likelihood: Optional[float] = Field(default=None, ge=1, le=5)
    impact: Optional[float] = Field(default=None, ge=1, le=5)
    detection: Optional[float] = Field(default=None, ge=0, le=5)
    controls_effectiveness: Optional[float] = Field(default=None, ge=0, le=1)
    custom_factors: Optional[List[RiskFactor]] = None
    linked_capa_ids: Optional[List[str]] = None
    linked_audit_finding_ids: Optional[List[str]] = None
    next_review_date: Optional[datetime] = None
    revision_note: Optional[str] = None


class RiskReportTemplate(BaseModel):
    id: str = "default"
    name: str = "Risk Executive Summary"
    description: Optional[str] = None
    body: str = (
        "<h1>{{title}}</h1><p>Status: {{status}}</p>"
        "<p>Inherent Score: {{inherent_score}}</p>"
        "<p>Residual Score: {{residual_score}}</p>"
        "<p>Owner: {{owner}}</p>"
        "<p>Controls: {{controls}}</p>"
    )
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    updated_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    updated_by: Optional[str] = None


class ReportFilterDefinition(BaseModel):
    field: str
    label: str
    operator: str = "eq"
    field_type: str = "string"  # string, number, boolean, datetime, list
    description: Optional[str] = None
    required: bool = False
    choices: Optional[List[str]] = None
    default: Optional[Any] = None
    allow_multiple: bool = False


class ReportMetricConfig(BaseModel):
    name: Optional[str] = None
    field: Optional[str] = None
    operation: str = "count"  # count, sum, avg, min, max
    precision: Optional[int] = 2
    label: Optional[str] = None


class ReportVisualizationConfig(BaseModel):
    type: str = "table"  # table, bar, line, pie, kpi
    title: Optional[str] = None
    group_by: Optional[str] = None
    metrics: List[ReportMetricConfig] = Field(default_factory=list)
    fields: Optional[List[str]] = None
    options: Dict[str, Any] = Field(default_factory=dict)


class ReportDefinition(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    name: str
    module: str
    description: Optional[str] = None
    filters: List[ReportFilterDefinition] = Field(default_factory=list)
    visualizations: List[ReportVisualizationConfig] = Field(default_factory=list)
    fields: Optional[List[str]] = None
    max_rows: int = 2000
    default_template_id: Optional[str] = None
    created_by: str
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    updated_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))


class ReportDefinitionCreate(BaseModel):
    name: str
    module: str
    description: Optional[str] = None
    filters: Optional[List[ReportFilterDefinition]] = None
    visualizations: Optional[List[ReportVisualizationConfig]] = None
    fields: Optional[List[str]] = None
    max_rows: Optional[int] = None
    default_template_id: Optional[str] = None


class ReportDefinitionUpdate(BaseModel):
    name: Optional[str] = None
    description: Optional[str] = None
    filters: Optional[List[ReportFilterDefinition]] = None
    visualizations: Optional[List[ReportVisualizationConfig]] = None
    fields: Optional[List[str]] = None
    max_rows: Optional[int] = None
    default_template_id: Optional[str] = None


class ReportFilterValue(BaseModel):
    field: str
    operator: Optional[str] = None
    value: Any


class ReportRunRequest(BaseModel):
    definition_id: str
    filters: Optional[List[ReportFilterValue]] = None
    visualization_overrides: Optional[List[ReportVisualizationConfig]] = None


class ReportRunResult(BaseModel):
    definition: ReportDefinition
    filters: List[ReportFilterValue]
    rows: List[Dict[str, Any]]
    total: int
    visualizations: List[Dict[str, Any]]
    generated_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    summary: Dict[str, Any] = Field(default_factory=dict)


class ReportTemplatePlaceholder(BaseModel):
    key: str
    source: str  # header, summary, table, metric
    path: Optional[str] = None
    description: Optional[str] = None


class ReportTemplate(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    name: str
    module: str
    template_type: Literal["docx", "xlsx", "html"]
    description: Optional[str] = None
    file_upload_id: Optional[str] = None
    placeholders: List[ReportTemplatePlaceholder] = Field(default_factory=list)
    default_definition_id: Optional[str] = None
    created_by: str
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    updated_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))


class ReportTemplateCreate(BaseModel):
    name: str
    module: str
    template_type: Literal["docx", "xlsx", "html"]
    description: Optional[str] = None
    file_upload_id: Optional[str] = None
    placeholders: Optional[List[ReportTemplatePlaceholder]] = None
    default_definition_id: Optional[str] = None


class ReportTemplateUpdate(BaseModel):
    name: Optional[str] = None
    description: Optional[str] = None
    file_upload_id: Optional[str] = None
    template_type: Optional[Literal["docx", "xlsx", "html"]] = None
    placeholders: Optional[List[ReportTemplatePlaceholder]] = None
    default_definition_id: Optional[str] = None


class ReportTemplateRenderRequest(BaseModel):
    filters: Optional[List[ReportFilterValue]] = None
    definition_id: Optional[str] = None
    include_rows: bool = True
    include_summary: bool = True
    include_visualizations: bool = True


class ReportTemplateRenderResponse(BaseModel):
    template: ReportTemplate
    definition: Optional[ReportDefinition] = None
    run_result: Optional[ReportRunResult] = None
    placeholders: List[ReportTemplatePlaceholder] = Field(default_factory=list)


class CalibrationDeviceStatusChange(BaseModel):
    status: str
    changed_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    changed_by: str
    note: Optional[str] = None


class CalibrationWorkOrderStatusChange(BaseModel):
    status: str
    changed_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    changed_by: str
    note: Optional[str] = None


class CalibrationMeasurementRecord(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    parameter: str
    nominal: float
    tolerance: Optional[float] = None
    observed: float
    pass_status: bool
    recorded_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    recorded_by: Optional[str] = None
    note: Optional[str] = None


class CalibrationCostEntry(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    description: str
    amount: float
    currency: str = "TRY"
    incurred_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    supplier: Optional[str] = None
    reference: Optional[str] = None
    note: Optional[str] = None


class CalibrationWorkOrder(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    work_order_no: str
    device_id: str
    planned_date: datetime
    due_date: datetime
    status: str = "planned"
    status_history: List[CalibrationWorkOrderStatusChange] = Field(default_factory=list)
    assigned_to: Optional[str] = None
    created_by: str
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    updated_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    completed_at: Optional[datetime] = None
    result: Optional[str] = None
    notes: Optional[str] = None
    measurement_records: List[CalibrationMeasurementRecord] = Field(default_factory=list)
    cost_entries: List[CalibrationCostEntry] = Field(default_factory=list)
    linked_dof_task_id: Optional[str] = None
    linked_capa_ids: List[str] = Field(default_factory=list)


class CalibrationDevice(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    device_code: str
    name: str
    category: str
    location: Optional[str] = None
    manufacturer: Optional[str] = None
    model: Optional[str] = None
    serial_number: Optional[str] = None
    department: Optional[str] = None
    responsible_person: Optional[str] = None
    status: str = "active"
    status_history: List[CalibrationDeviceStatusChange] = Field(default_factory=list)
    calibration_interval_days: int = 365
    last_calibrated_at: Optional[datetime] = None
    next_due_date: Optional[datetime] = None
    notice_days: int = 14
    file_attachments: List[str] = Field(default_factory=list)
    linked_capa_ids: List[str] = Field(default_factory=list)
    linked_dof_task_ids: List[str] = Field(default_factory=list)
    notes: Optional[str] = None
    created_by: str
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    updated_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))


class CalibrationDeviceCreate(BaseModel):
    name: str
    category: str
    location: Optional[str] = None
    manufacturer: Optional[str] = None
    model: Optional[str] = None
    serial_number: Optional[str] = None
    department: Optional[str] = None
    responsible_person: Optional[str] = None
    status: Optional[str] = "active"
    calibration_interval_days: int = 365
    notice_days: Optional[int] = 14
    last_calibrated_at: Optional[datetime] = None
    next_due_date: Optional[datetime] = None
    linked_capa_ids: Optional[List[str]] = None
    linked_dof_task_ids: Optional[List[str]] = None
    file_attachments: Optional[List[str]] = None
    notes: Optional[str] = None


class CalibrationDeviceUpdate(BaseModel):
    name: Optional[str] = None
    category: Optional[str] = None
    location: Optional[str] = None
    manufacturer: Optional[str] = None
    model: Optional[str] = None
    serial_number: Optional[str] = None
    department: Optional[str] = None
    responsible_person: Optional[str] = None
    status: Optional[str] = None
    status_note: Optional[str] = None
    calibration_interval_days: Optional[int] = None
    notice_days: Optional[int] = None
    last_calibrated_at: Optional[datetime] = None
    next_due_date: Optional[datetime] = None
    linked_capa_ids: Optional[List[str]] = None
    linked_dof_task_ids: Optional[List[str]] = None
    file_attachments: Optional[List[str]] = None
    notes: Optional[str] = None


class CalibrationWorkOrderCreate(BaseModel):
    planned_date: datetime
    due_date: datetime
    assigned_to: Optional[str] = None
    notes: Optional[str] = None


class CalibrationWorkOrderUpdate(BaseModel):
    planned_date: Optional[datetime] = None
    due_date: Optional[datetime] = None
    assigned_to: Optional[str] = None
    status: Optional[str] = None
    status_note: Optional[str] = None
    notes: Optional[str] = None
    linked_dof_task_id: Optional[str] = None


class CalibrationMeasurementCreate(BaseModel):
    parameter: str
    nominal: float
    tolerance: Optional[float] = None
    observed: float
    pass_status: bool
    recorded_at: Optional[datetime] = None
    recorded_by: Optional[str] = None
    note: Optional[str] = None


class CalibrationCostCreate(BaseModel):
    description: str
    amount: float
    currency: Optional[str] = "TRY"
    incurred_at: Optional[datetime] = None
    supplier: Optional[str] = None
    reference: Optional[str] = None
    note: Optional[str] = None


class CalibrationReportSummary(BaseModel):
    total_devices: int
    active_devices: int
    overdue_devices: int
    upcoming_devices: int
    total_work_orders: int
    open_work_orders: int
    cost_total: float
    currency: str = "TRY"


class CalibrationDofRequest(BaseModel):
    title: Optional[str] = None
    department: Optional[str] = None
    responsible_person: Optional[str] = None
    due_date: Optional[datetime] = None
    note: Optional[str] = None


class DofStatusHistory(BaseModel):
    status: str
    changed_by: str
    changed_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    note: Optional[str] = None

class DofTask(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    dof_no: str
    title: str
    description: Optional[str] = None
    department: str
    responsible_person: str
    due_date: datetime
    status: str = "open"  # open, in_progress, pending_closure, closed, cancelled
    team_members: List[str] = Field(default_factory=list)
    initial_improvement_report_date: Optional[datetime] = None
    created_by: str
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    updated_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    closed_at: Optional[datetime] = None
    status_history: List[DofStatusHistory] = Field(default_factory=list)
    closure_requested_at: Optional[datetime] = None
    closure_requested_by: Optional[str] = None
    closure_request_note: Optional[str] = None
    closure_approved_at: Optional[datetime] = None
    closure_approved_by: Optional[str] = None
    closure_decision_note: Optional[str] = None

class DofTaskCreate(BaseModel):
    title: str
    description: Optional[str] = None
    department: str
    responsible_person: str
    due_date: datetime
    team_members: Optional[List[str]] = None
    initial_improvement_report_date: Optional[datetime] = None

class DofTaskUpdate(BaseModel):
    title: Optional[str] = None
    description: Optional[str] = None
    department: Optional[str] = None
    responsible_person: Optional[str] = None
    due_date: Optional[datetime] = None
    team_members: Optional[List[str]] = None
    initial_improvement_report_date: Optional[datetime] = None


class DofClosureRequest(BaseModel):
    note: Optional[str] = None


class DofClosureDecision(BaseModel):
    approve: bool
    note: Optional[str] = None

class DofTaskListResponse(BaseModel):
    items: List[DofTask]
    total: int
    page: int
    page_size: int

class DofStatusUpdate(BaseModel):
    status: str
    note: Optional[str] = None

# Notification Models
class Notification(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    user_id: str
    title: str
    message: str
    type: str  # info, warning, error, success
    is_read: bool = False
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

# Dashboard Stats Model
class DashboardStats(BaseModel):
    total_documents: int
    pending_approvals: int
    total_complaints: int
    open_complaints: int
    total_capas: int
    open_capas: int
    recent_activities: List[Dict[str, Any]]
    notifications_count: int


class DashboardWorkItem(BaseModel):
    id: str
    module: str
    title: str
    code: Optional[str] = None
    status: Optional[str] = None
    due_date: Optional[datetime] = None
    description: Optional[str] = None


class DashboardSection(BaseModel):
    total: int
    open: int
    pending: int
    items: List[DashboardWorkItem] = Field(default_factory=list)


class DashboardOverview(BaseModel):
    documents: DashboardSection
    complaints: DashboardSection
    capas: DashboardSection
    dof: DashboardSection
    notifications: int

# Utility Functions
def hash_password(password: str) -> str:
    return hashlib.sha256(password.encode()).hexdigest()

def verify_password(password: str, hashed_password: str) -> bool:
    return hash_password(password) == hashed_password

def create_access_token(data: dict, expires_delta: Optional[timedelta] = None):
    to_encode = data.copy()
    if expires_delta:
        expire = datetime.now(timezone.utc) + expires_delta
    else:
        expire = datetime.now(timezone.utc) + timedelta(minutes=15)
    to_encode.update({"exp": expire})
    encoded_jwt = jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)
    return encoded_jwt

async def get_current_user(credentials: Optional[HTTPAuthorizationCredentials] = Depends(security)) -> User:
    if DISABLE_AUTH:
        guest = await db.users.find_one({"username": GUEST_USERNAME})
        if guest:
            return await build_user_model(guest)

        guest_id = str(uuid.uuid4())
        guest_record = {
            "id": guest_id,
            "username": GUEST_USERNAME,
            "email": GUEST_EMAIL,
            "password": hash_password(os.getenv("GUEST_PASSWORD", "guest123")),
            "full_name": GUEST_FULL_NAME,
            "role": "Viewer",
            "roles": ["Viewer"],
            "department": GUEST_DEPARTMENT,
            "is_active": True,
            "created_at": datetime.now(timezone.utc),
        }
        await db.users.insert_one(guest_record)
        return await build_user_model(guest_record)

    if credentials is None:
        raise HTTPException(status_code=401, detail="Authentication required")

    try:
        payload = jwt.decode(credentials.credentials, SECRET_KEY, algorithms=[ALGORITHM])
        username: str = payload.get("sub")
        if username is None:
            raise HTTPException(status_code=401, detail="Invalid authentication credentials")
    except jwt.PyJWTError:
        raise HTTPException(status_code=401, detail="Invalid authentication credentials")
    
    user = await db.users.find_one({"username": username})
    if user is None:
        raise HTTPException(status_code=401, detail="User not found")

    return await build_user_model(user)

# Generate unique codes
async def generate_complaint_no() -> str:
    count = await db.complaints.count_documents({})
    return f"COMP-{datetime.now().year}-{count + 1:04d}"

async def generate_capa_no() -> str:
    count = await db.capas.count_documents({})
    return f"CAPA-{datetime.now().year}-{count + 1:04d}"

async def generate_dof_no() -> str:
    count = await db.dof_tasks.count_documents({})
    return f"DOF-{datetime.now().year}-{count + 1:04d}"

async def generate_audit_code() -> str:
    count = await db.audits.count_documents({})
    return f"AUD-{datetime.now().year}-{count + 1:04d}"

async def generate_risk_code() -> str:
    count = await db.risks.count_documents({})
    return f"RISK-{datetime.now().year}-{count + 1:04d}"

def _parse_iso_datetime(date_str: str) -> datetime:
    parsed = datetime.fromisoformat(date_str)
    if parsed.tzinfo is None:
        parsed = parsed.replace(tzinfo=timezone.utc)
    return parsed

# Email Service
class EmailService:
    @staticmethod
    def send_email(recipients: List[str], subject: str, body: str, html_body: Optional[str] = None) -> bool:
        if not SMTP_PASSWORD:  # Skip if password not set
            print(f"Email would be sent to {recipients}: {subject}")
            return True
            
        try:
            msg = MIMEMultipart("alternative")
            msg["Subject"] = subject
            msg["From"] = f"{SMTP_FROM_NAME} <{SMTP_USER}>"
            msg["To"] = ", ".join(recipients)
            
            text_part = MIMEText(body, "plain")
            msg.attach(text_part)
            
            if html_body:
                html_part = MIMEText(html_body, "html")
                msg.attach(html_part)
            
            context = ssl.create_default_context()
            with smtplib.SMTP(SMTP_SERVER, SMTP_PORT) as server:
                server.starttls(context=context)
                if SMTP_PASSWORD:
                    server.login(SMTP_USER, SMTP_PASSWORD)
                server.sendmail(SMTP_USER, recipients, msg.as_string())
                
            return True
        except Exception as e:
            print(f"Failed to send email: {str(e)}")
            return False

# File Upload Functions
@api_router.post("/upload")
async def upload_file(
    file: UploadFile = File(...),
    module_type: str = Form(...),
    module_id: Optional[str] = Form(None),
    current_user: User = Depends(get_current_user)
):
    # Generate unique filename
    file_id = str(uuid.uuid4())
    file_extension = Path(file.filename).suffix
    filename = f"{file_id}{file_extension}"
    file_path = UPLOAD_DIR / filename
    
    # Save file
    async with aiofiles.open(file_path, 'wb') as f:
        content = await file.read()
        await f.write(content)
    
    # Save file record to database
    file_record = FileUpload(
        filename=filename,
        original_filename=file.filename,
        file_path=str(file_path),
        file_size=len(content),
        mime_type=file.content_type,
        uploaded_by=current_user.id,
        module_type=module_type,
        module_id=module_id
    )
    
    await db.file_uploads.insert_one(file_record.dict())
    
    return {"file_id": file_record.id, "filename": filename, "original_filename": file.filename}

@api_router.get("/download/{file_id}")
async def download_file(file_id: str, current_user: User = Depends(get_current_user)):
    file_record = await db.file_uploads.find_one({"id": file_id})
    if not file_record:
        raise HTTPException(status_code=404, detail="File not found")
    
    file_path = Path(file_record["file_path"])
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="File not found on disk")
    
    return FileResponse(
        path=str(file_path),
        filename=file_record["original_filename"],
        media_type=file_record["mime_type"]
    )

@api_router.delete("/files/{file_id}")
async def delete_file(file_id: str, current_user: User = Depends(get_current_user)):
    file_record = await db.file_uploads.find_one({"id": file_id})
    if not file_record:
        raise HTTPException(status_code=404, detail="File not found")
    
    # Delete from disk
    file_path = Path(file_record["file_path"])
    if file_path.exists():
        file_path.unlink()
    
    # Delete from database
    await db.file_uploads.delete_one({"id": file_id})
    
    return {"message": "File deleted successfully"}

# Authentication Routes
@api_router.post("/auth/login", response_model=Token)
async def login(user_credentials: UserLogin):
    user = await db.users.find_one({"username": user_credentials.username})
    if not user or not verify_password(user_credentials.password, user["password"]):
        raise HTTPException(status_code=401, detail="Invalid username or password")
    
    if not user.get("is_active", True):
        raise HTTPException(status_code=401, detail="User account is disabled")
    
    access_token_expires = timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
    access_token = create_access_token(
        data={"sub": user["username"]}, expires_delta=access_token_expires
    )
    
    user_obj = await build_user_model(user)
    return Token(access_token=access_token, token_type="bearer", user=user_obj)

@api_router.post("/auth/register", response_model=User)
async def register(user_data: UserCreate):
    # Check if user already exists
    existing_user = await db.users.find_one({"$or": [{"username": user_data.username}, {"email": user_data.email}]})
    if existing_user:
        raise HTTPException(status_code=400, detail="Username or email already registered")
    
    # Hash password
    hashed_password = hash_password(user_data.password)
    
    base_roles = user_data.roles or ([] if not user_data.role else [user_data.role])
    roles = _normalize_string_list(base_roles)
    groups = _normalize_string_list(user_data.groups)
    
    # Create user record
    user_record = {
        "id": str(uuid.uuid4()),
        "username": user_data.username,
        "email": user_data.email,
        "password": hashed_password,
        "full_name": user_data.full_name,
        "role": user_data.role.strip() if user_data.role else "",
        "roles": roles,
        "groups": groups,
        "department": user_data.department,
        "is_active": True,
        "created_at": datetime.now(timezone.utc),
    }
    
    await db.users.insert_one(user_record)
    
    return await build_user_model(user_record)

@api_router.get("/auth/me", response_model=User)
async def get_me(current_user: User = Depends(get_current_user)):
    return current_user

# RBAC & User Management Routes


def ensure_admin_permission(user: User, permission: str) -> None:
    ensure_permission(user, permission)


@api_router.get("/roles", response_model=List[UserRole])
async def list_roles(current_user: User = Depends(get_current_user)):
    ensure_admin_permission(current_user, "admin.roles.manage")
    cursor = db.roles.find({}).sort("name", 1)
    roles: List[UserRole] = []
    async for role_doc in cursor:
        roles.append(UserRole(**role_doc))
    return roles


@api_router.post("/roles", response_model=UserRole, status_code=status.HTTP_201_CREATED)
async def create_role(
    payload: RoleCreate,
    current_user: User = Depends(get_current_user),
) -> UserRole:
    ensure_admin_permission(current_user, "admin.roles.manage")
    name = payload.name.strip()
    if not name:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Role name cannot be empty")

    existing = await db.roles.find_one({"name": name})
    if existing:
        raise HTTPException(status_code=status.HTTP_409_CONFLICT, detail="Role already exists")

    now = datetime.now(timezone.utc)
    role_doc = {
        "name": name,
        "description": (payload.description or "").strip() or None,
        "permissions": _normalize_permissions_list(payload.permissions),
        "created_at": now,
        "updated_at": now,
    }
    await db.roles.insert_one(role_doc)
    return UserRole(**role_doc)


@api_router.patch("/roles/{role_name}", response_model=UserRole)
async def update_role(
    role_name: str,
    payload: RoleUpdate,
    current_user: User = Depends(get_current_user),
) -> UserRole:
    ensure_admin_permission(current_user, "admin.roles.manage")
    updates: Dict[str, Any] = {}
    if payload.description is not None:
        updates["description"] = payload.description.strip() or None
    if payload.permissions is not None:
        updates["permissions"] = _normalize_permissions_list(payload.permissions)
    if not updates:
        existing = await db.roles.find_one({"name": role_name})
        if not existing:
            raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Role not found")
        return UserRole(**existing)

    updates["updated_at"] = datetime.now(timezone.utc)
    updated = await db.roles.find_one_and_update(
        {"name": role_name},
        {"$set": updates},
        return_document=True,
    )
    if not updated:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Role not found")
    return UserRole(**updated)


@api_router.delete("/roles/{role_name}", status_code=status.HTTP_204_NO_CONTENT)
async def delete_role(
    role_name: str,
    current_user: User = Depends(get_current_user),
):
    ensure_admin_permission(current_user, "admin.roles.manage")
    result = await db.roles.delete_one({"name": role_name})
    if result.deleted_count == 0:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Role not found")
    # Remove role from users
    await db.users.update_many({}, {"$pull": {"roles": role_name}})
    return Response(status_code=status.HTTP_204_NO_CONTENT)


@api_router.get("/users", response_model=List[UserListItem])
async def list_users(current_user: User = Depends(get_current_user)):
    ensure_admin_permission(current_user, "admin.users.manage")
    cursor = db.users.find({}).sort("created_at", -1)
    items: List[UserListItem] = []
    async for user_doc in cursor:
        items.append(
            UserListItem(
                id=user_doc["id"],
                username=user_doc["username"],
                email=user_doc["email"],
                full_name=user_doc.get("full_name", ""),
                role=user_doc.get("role", ""),
                roles=user_doc.get("roles", []),
                department=user_doc.get("department", ""),
                groups=user_doc.get("groups", []),
                permissions=user_doc.get("permissions", []),
                is_active=user_doc.get("is_active", True),
                created_at=user_doc.get("created_at", datetime.now(timezone.utc)),
            )
        )
    return items


@api_router.patch("/users/{user_id}", response_model=User)
async def update_user(
    user_id: str,
    payload: UserUpdate,
    current_user: User = Depends(get_current_user),
) -> User:
    ensure_admin_permission(current_user, "admin.users.manage")
    user_doc = await db.users.find_one({"id": user_id})
    if not user_doc:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="User not found")

    update_fields: Dict[str, Any] = {}
    if payload.role is not None:
        update_fields["role"] = payload.role.strip()
    if payload.roles is not None:
        update_fields["roles"] = _normalize_string_list(payload.roles)
    if payload.department is not None:
        update_fields["department"] = payload.department.strip()
    if payload.groups is not None:
        update_fields["groups"] = _normalize_string_list(payload.groups)
    if payload.permissions is not None:
        update_fields["permissions"] = _normalize_permissions_list(payload.permissions)
    if payload.is_active is not None:
        update_fields["is_active"] = bool(payload.is_active)

    if not update_fields:
        return await build_user_model(user_doc)

    update_fields["updated_at"] = datetime.now(timezone.utc)
    await db.users.update_one({"id": user_id}, {"$set": update_fields})
    updated_doc = await db.users.find_one({"id": user_id})
    return await build_user_model(updated_doc)

# Dashboard Routes
@api_router.get("/dashboard/stats", response_model=DashboardStats)
async def get_dashboard_stats(current_user: User = Depends(get_current_user)):
    # Get document stats
    total_documents = await db.documents.count_documents({})
    pending_approvals = await db.documents.count_documents({"status": "review"})
    
    # Get complaint stats
    total_complaints = await db.complaints.count_documents({})
    open_complaints = await db.complaints.count_documents({"status": {"$in": ["open", "investigating"]}})
    
    # Get CAPA stats
    total_capas = await db.capas.count_documents({})
    open_capas = await db.capas.count_documents({"status": {"$in": ["open", "investigating", "implementing", "pending_closure"]}})
    
    # Get recent activities
    recent_activities = []
    
    # Recent documents
    recent_docs = await db.documents.find({}).sort("updated_at", -1).limit(5).to_list(5)
    for doc in recent_docs:
        recent_activities.append({
            "type": "document",
            "title": f"Doküman '{doc['title']}' {doc['status']}",
            "timestamp": doc["updated_at"],
            "user": doc["author_id"]
        })
    
    # Recent complaints
    recent_complaints = await db.complaints.find({}).sort("updated_at", -1).limit(3).to_list(3)
    for complaint in recent_complaints:
        recent_activities.append({
            "type": "complaint",
            "title": f"Şikayet #{complaint['complaint_no']} - {complaint['status']}",
            "timestamp": complaint["updated_at"],
            "user": complaint["created_by"]
        })
    
    # Sort activities by timestamp
    recent_activities.sort(key=lambda x: x["timestamp"], reverse=True)
    recent_activities = recent_activities[:10]
    
    # Get notifications count
    notifications_count = await db.notifications.count_documents({"user_id": current_user.id, "is_read": False})
    
    return DashboardStats(
        total_documents=total_documents,
        pending_approvals=pending_approvals,
        total_complaints=total_complaints,
        open_complaints=open_complaints,
        total_capas=total_capas,
        open_capas=open_capas,
        recent_activities=recent_activities,
        notifications_count=notifications_count
    )


@api_router.get("/dashboard/overview", response_model=DashboardOverview)
async def get_dashboard_overview(current_user: User = Depends(get_current_user)):
    # Documents
    documents_total = await db.documents.count_documents({})
    documents_open = await db.documents.count_documents(
        {"status": {"$nin": ["archived", "retired"]}}
    )
    approvals = await list_pending_document_approvals(current_user)
    read_tasks_list = await list_read_tasks(current_user)

    document_items: List[DashboardWorkItem] = []
    for task in approvals[:5]:
        document_items.append(
            DashboardWorkItem(
                id=task.document_id,
                module="document_approval",
                title=f"{task.document_code} - {task.title}",
                code=task.document_code,
                status=f"Approval stage {task.stage}",
                due_date=task.deadline,
                description=f"Asama {task.stage} ({'Hepsi' if task.approval_type == 'all' else 'Herhangi biri'})",
            )
        )
    for task in read_tasks_list[:5]:
        document_items.append(
            DashboardWorkItem(
                id=task.document_id,
                module="document_read",
                title=f"{task.document_code} - {task.title}",
                code=task.document_code,
                status="Zorunlu okuma" if task.required else "Opsiyonel okuma",
                due_date=None,
                description="Okuma onayi bekleniyor",
            )
        )
    documents_section = DashboardSection(
        total=documents_total,
        open=documents_open,
        pending=len(approvals) + len(read_tasks_list),
        items=document_items,
    )

    # Complaints
    complaints_total = await db.complaints.count_documents({})
    complaints_open = await db.complaints.count_documents(
        {"status": {"$in": ["open", "investigating"]}}
    )
    complaint_items: List[DashboardWorkItem] = []
    complaint_cursor = (
        db.complaints.find({"status": {"$in": ["open", "investigating"]}})
        .sort("updated_at", -1)
        .limit(5)
    )
    async for complaint in complaint_cursor:
        complaint_items.append(
            DashboardWorkItem(
                id=complaint["id"],
                module="complaint",
                title=f"{complaint.get('complaint_no', '')} - {complaint.get('customer_name', '')}",
                status=complaint.get("status", ""),
                due_date=complaint.get("complaint_date"),
                description=(complaint.get("description") or "")[:140],
            )
        )
    complaints_section = DashboardSection(
        total=complaints_total,
        open=complaints_open,
        pending=complaints_open,
        items=complaint_items,
    )

    # CAPAs
    capas_total = await db.capas.count_documents({})
    capas_open = await db.capas.count_documents({"status": {"$ne": "closed"}})
    capa_items: List[DashboardWorkItem] = []
    capa_cursor = (
        db.capas.find({"status": {"$ne": "closed"}})
        .sort("target_date", 1)
        .limit(5)
    )
    async for capa in capa_cursor:
        target_date = capa.get("target_date")
        capa_items.append(
            DashboardWorkItem(
                id=capa["id"],
                module="capa",
                title=f"{capa.get('capa_no', '')} - {capa.get('title', '')}",
                status=capa.get("status", ""),
                due_date=target_date,
                description=(capa.get("nonconformity_description") or "")[:140],
            )
        )
    capas_section = DashboardSection(
        total=capas_total,
        open=capas_open,
        pending=capas_open,
        items=capa_items,
    )

    # DÖF tasks
    dof_total = await db.dof_tasks.count_documents({})
    dof_open = await db.dof_tasks.count_documents(
        {"status": {"$in": ["open", "in_progress", "pending_closure"]}}
    )
    dof_items: List[DashboardWorkItem] = []
    dof_cursor = (
        db.dof_tasks.find({"status": {"$in": ["open", "in_progress", "pending_closure"]}})
        .sort("due_date", 1)
        .limit(5)
    )
    async for task in dof_cursor:
        dof_items.append(
            DashboardWorkItem(
                id=task["id"],
                module="dof",
                title=f"{task.get('dof_no', '')} - {task.get('title', '')}",
                status=task.get("status", ""),
                due_date=task.get("due_date"),
                description=(task.get("description") or "")[:140],
            )
        )
    dof_section = DashboardSection(
        total=dof_total,
        open=dof_open,
        pending=dof_open,
        items=dof_items,
    )

    notifications_count = await db.notifications.count_documents(
        {"user_id": current_user.id, "is_read": False}
    )

    return DashboardOverview(
        documents=documents_section,
        complaints=complaints_section,
        capas=capas_section,
        dof=dof_section,
        notifications=notifications_count,
    )

# Document Folder Routes
@api_router.post("/document-folders", response_model=DocumentFolder)
async def create_document_folder(
    payload: DocumentFolderCreate,
    current_user: User = Depends(get_current_user),
):
    ensure_permission(current_user, "doc.folder.manage")
    sanitized_permissions = sanitize_folder_permissions(payload.permissions or [])
    permission_models = [FolderPermission(**perm) for perm in sanitized_permissions]
    folder = DocumentFolder(
        name=payload.name,
        code_prefix=payload.code_prefix,
        department=payload.department,
        description=payload.description,
        parent_id=payload.parent_id,
        auto_code_pattern=payload.auto_code_pattern or DOCUMENT_CODE_PATTERN,
        permissions=permission_models,
    )
    await db.document_folders.insert_one(folder.dict())
    return folder


@api_router.get("/document-folders", response_model=List[DocumentFolder])
async def list_document_folders(
    parent_id: Optional[str] = None,
    current_user: User = Depends(get_current_user),
):
    query: Dict[str, Any] = {}
    if parent_id is not None:
        query["parent_id"] = parent_id
    folders_raw = await db.document_folders.find(query).sort("name", 1).to_list(500)
    visible: List[DocumentFolder] = []
    for data in folders_raw:
        folder = DocumentFolder(**data)
        folder.current_user_capabilities = compute_folder_capabilities_for_user(folder, current_user)
        if user_is_admin(current_user) or "read" in folder.current_user_capabilities:
            visible.append(folder)
    return visible


@api_router.get("/document-folders/{folder_id}", response_model=DocumentFolder)
async def get_document_folder(folder_id: str, current_user: User = Depends(get_current_user)):
    return await ensure_document_folder(folder_id, current_user, capability="read")


@api_router.patch("/document-folders/{folder_id}", response_model=DocumentFolder)
async def update_document_folder(
    folder_id: str,
    payload: DocumentFolderUpdate,
    current_user: User = Depends(get_current_user),
):
    await ensure_document_folder(folder_id, current_user, capability="manage")
    update_data = payload.dict(exclude_none=True)
    if update_data:
        update_data["updated_at"] = datetime.now(timezone.utc)
        await db.document_folders.update_one({"id": folder_id}, {"$set": update_data})
    return await ensure_document_folder(folder_id, current_user, capability="manage")


@api_router.patch("/document-folders/{folder_id}/permissions", response_model=DocumentFolder)
async def set_document_folder_permissions(
    folder_id: str,
    payload: DocumentFolderPermissionsUpdate,
    current_user: User = Depends(get_current_user),
):
    await ensure_document_folder(folder_id, current_user, capability="manage")
    permissions = sanitize_folder_permissions(payload.permissions or [])
    await db.document_folders.update_one(
        {"id": folder_id},
        {"$set": {"permissions": permissions, "updated_at": datetime.now(timezone.utc)}},
    )
    return await ensure_document_folder(folder_id, current_user, capability="manage")


# Document Management Routes
@api_router.get("/documents", response_model=List[Document])
async def list_documents(
    department: Optional[str] = None,
    status_filter: Optional[str] = None,
    folder_id: Optional[str] = None,
    document_type: Optional[str] = None,
    search: Optional[str] = None,
    current_user: User = Depends(get_current_user),
):
    query: Dict[str, Any] = {}
    if department:
        query["department"] = department
    if status_filter:
        query["status"] = status_filter
    if folder_id:
        query["folder_id"] = folder_id
    if document_type:
        query["document_type"] = document_type
    if search:
        query["$or"] = [
            {"title": {"$regex": search, "$options": "i"}},
            {"code": {"$regex": search, "$options": "i"}},
            {"tags": {"$regex": search, "$options": "i"}},
        ]
    allowed_folders = await get_authorized_folder_ids(current_user, "read")
    if allowed_folders is not None:
        if not allowed_folders:
            return []
        folder_filter = query.get("folder_id")
        if folder_filter is None:
            query["folder_id"] = {"$in": list(allowed_folders)}
        elif isinstance(folder_filter, str):
            if folder_filter not in allowed_folders:
                return []
        elif isinstance(folder_filter, dict):
            existing = set(folder_filter.get("$in", []))
            constrained = existing.intersection(allowed_folders) if existing else allowed_folders
            if not constrained:
                return []
            folder_filter["$in"] = list(constrained)
    documents = await db.documents.find(query).sort("updated_at", -1).to_list(500)
    return [hydrate_document(doc) for doc in documents]


@api_router.post("/documents", response_model=Document)
async def create_document(
    payload: DocumentCreate,
    current_user: User = Depends(get_current_user),
):
    folder = await ensure_document_folder(payload.folder_id, current_user, capability="create")
    code = await generate_document_code(folder, payload.document_type)
    distribution = [DocumentDistribution(**item.dict()) for item in payload.distribution_list]
    approval_stages = [DocumentApprovalStage(**stage.dict()) for stage in payload.approval_matrix]
    approval_matrix = normalize_approval_stages(approval_stages)
    read_receipts = build_initial_read_receipts(distribution)

    initial_status = "review" if approval_matrix else "approved"
    now = datetime.now(timezone.utc)

    version_record = DocumentVersion(
        version="1.0",
        changes=payload.version_notes,
        notes=payload.version_notes,
        created_by=current_user.id,
        file_id=payload.file_id,
        status="pending_approval" if approval_matrix else "published",
    )

    status_history: List[DocumentStatus] = [
        DocumentStatus(
            status="draft",
            changed_by=current_user.id,
            comment="Dokuman olusturuldu.",
        )
    ]
    if initial_status != "draft":
        status_history.append(
            DocumentStatus(
                status=initial_status,
                changed_by=current_user.id,
                comment="Onay sureci baslatildi."
                if initial_status == "review"
                else "Dokuman yayinlandi.",
            )
        )

    document = Document(
        folder_id=folder.id,
        code=code,
        title=payload.title,
        description=payload.description,
        document_type=payload.document_type,
        department=payload.department or folder.department or current_user.department,
        status=initial_status,
        author_id=current_user.id,
        version=version_record.version,
        tags=payload.tags or [],
        distribution_list=distribution,
        approval_matrix=approval_matrix,
        read_receipts=read_receipts,
        status_history=status_history,
        version_history=[version_record],
        current_version_id=version_record.id,
        review_date=payload.review_date,
        expiry_date=payload.expiry_date,
        published_at=now if initial_status == "approved" else None,
    )

    await db.documents.insert_one(document.dict())
    if payload.file_id:
        await db.file_uploads.update_one(
            {"id": payload.file_id},
            {"$set": {"module_id": document.id, "module_type": "document"}},
        )

    notification = Notification(
        user_id=current_user.id,
        title="Dokuman olusturuldu",
        message=f"'{document.title}' ({document.code}) dokumani olusturuldu.",
        type="info",
    )
    await db.notifications.insert_one(notification.dict())

    if approval_matrix:
        pending_stage_index = find_pending_stage_index(document)
        if pending_stage_index is not None:
            await notify_document_approvers(document, document.approval_matrix[pending_stage_index])
    elif document.status == "approved" and document.read_receipts:
        await notify_document_readers(
            document,
            note=f"{document.code} / {document.title} yayina alindi. Okuma onayiniz bekleniyor.",
        )

    return document


@api_router.get("/documents/{document_id}", response_model=Document)
async def get_document(document_id: str, current_user: User = Depends(get_current_user)):
    document = await db.documents.find_one({"id": document_id})
    if not document:
        raise HTTPException(status_code=404, detail="Document not found")
    await ensure_document_folder(document["folder_id"], current_user, capability="read")
    return hydrate_document(document)


@api_router.patch("/documents/{document_id}", response_model=Document)
async def update_document(
    document_id: str,
    payload: DocumentUpdate,
    current_user: User = Depends(get_current_user),
):
    document = await db.documents.find_one({"id": document_id})
    if not document:
        raise HTTPException(status_code=404, detail="Document not found")
    await ensure_document_folder(document["folder_id"], current_user, capability="revise")

    update_fields = payload.dict(exclude_none=True)
    set_payload: Dict[str, Any] = {}
    new_read_user_ids: Set[str] = set()

    if "distribution_list" in update_fields:
        distribution = [DocumentDistribution(**item.dict()) for item in update_fields["distribution_list"]]
        set_payload["distribution_list"] = [item.dict() for item in distribution]
        new_receipts = build_initial_read_receipts(distribution)
        set_payload["read_receipts"] = [receipt.dict() for receipt in new_receipts]
        previous_receipts = document.get("read_receipts", [])
        previous_user_ids = {
            str(item.get("user_id"))
            for item in previous_receipts
            if item.get("user_id")
        }
        new_read_user_ids = {
            receipt.user_id
            for receipt in new_receipts
            if receipt.user_id and receipt.user_id not in previous_user_ids
        }
        update_fields.pop("distribution_list")
    if "approval_matrix" in update_fields:
        matrix = [DocumentApprovalStage(**stage.dict()) for stage in update_fields["approval_matrix"]]
        normalized_matrix = normalize_approval_stages(matrix)
        set_payload["approval_matrix"] = [stage.dict() for stage in normalized_matrix]
        update_fields.pop("approval_matrix")

    for key, value in update_fields.items():
        set_payload[key] = value

    if not set_payload:
        return hydrate_document(document)

    set_payload["updated_at"] = datetime.now(timezone.utc)

    await db.documents.update_one({"id": document_id}, {"$set": set_payload})
    updated = await db.documents.find_one({"id": document_id})
    updated_document = hydrate_document(updated)

    if "approval_matrix" in set_payload:
        pending_stage_index = find_pending_stage_index(updated_document)
        if pending_stage_index is not None:
            await notify_document_approvers(
                updated_document,
                updated_document.approval_matrix[pending_stage_index],
            )

    if new_read_user_ids and updated_document.status == "approved":
        await notify_document_readers(
            updated_document,
            new_read_user_ids,
            note=f"{updated_document.code} / {updated_document.title} icin okuma onayi atandi.",
        )

    return updated_document


@api_router.post("/documents/{document_id}/versions", response_model=Document)
async def add_document_version(
    document_id: str,
    payload: DocumentVersionCreate,
    current_user: User = Depends(get_current_user),
):
    document = await db.documents.find_one({"id": document_id})
    if not document:
        raise HTTPException(status_code=404, detail="Document not found")
    await ensure_document_folder(document["folder_id"], current_user, capability="revise")
    hydrated_document = hydrate_document(document)

    version_history_raw = document.get("version_history", [])
    version_label = determine_next_version(version_history_raw, payload.mark_as_published)
    now = datetime.now(timezone.utc)

    version_status = "pending_approval" if hydrated_document.approval_matrix else "published"
    version_record = DocumentVersion(
        version=version_label,
        changes=payload.changes,
        notes=payload.notes,
        created_by=current_user.id,
        created_at=now,
        file_id=payload.file_id,
        status=version_status,
    )

    new_status = "review" if hydrated_document.approval_matrix else "approved"
    status_comment = payload.notes or payload.changes or "Yeni revizyon olusturuldu."
    status_entry = DocumentStatus(
        status=new_status,
        changed_by=current_user.id,
        comment=status_comment,
    )

    set_payload = {
        "current_version_id": version_record.id,
        "status": new_status,
        "version": version_record.version,
        "updated_at": now,
    }
    if new_status == "approved":
        set_payload["published_at"] = now
    if hydrated_document.approval_matrix:
        reset_matrix = reset_approval_stages(hydrated_document.approval_matrix)
        set_payload["approval_matrix"] = [stage.dict() for stage in reset_matrix]

    update_doc: Dict[str, Any] = {
        "$set": set_payload,
        "$push": {
            "version_history": version_record.dict(),
            "status_history": status_entry.dict(),
        },
    }

    await db.documents.update_one({"id": document_id}, update_doc)
    if payload.file_id:
        await db.file_uploads.update_one(
            {"id": payload.file_id},
            {"$set": {"module_id": document_id, "module_type": "document"}},
        )
    updated = await db.documents.find_one({"id": document_id})
    updated_document = hydrate_document(updated)

    if updated_document.approval_matrix:
        pending_stage_index = find_pending_stage_index(updated_document)
        if pending_stage_index is not None:
            await notify_document_approvers(
                updated_document,
                updated_document.approval_matrix[pending_stage_index],
                exclude_user_ids={current_user.id},
            )

    if updated_document.status == "approved":
        await notify_document_readers(
            updated_document,
            note=f"{updated_document.code} / {updated_document.title} icin {updated_document.version} versiyonu yayinda. Okuma onayiniz bekleniyor.",
        )

    return updated_document


@api_router.get("/documents/{document_id}/read-receipts", response_model=List[DocumentReadReceipt])
async def get_document_read_receipts(
    document_id: str,
    current_user: User = Depends(get_current_user),
):
    document = await db.documents.find_one({"id": document_id})
    if not document:
        raise HTTPException(status_code=404, detail="Document not found")
    await ensure_document_folder(document["folder_id"], current_user, capability="read")
    receipts = document.get("read_receipts", [])
    return [DocumentReadReceipt(**receipt) for receipt in receipts]


@api_router.post("/documents/{document_id}/acknowledge", response_model=Document)
async def acknowledge_document_read(
    document_id: str,
    payload: DocumentReadAcknowledge,
    current_user: User = Depends(get_current_user),
):
    document = await db.documents.find_one({"id": document_id})
    if not document:
        raise HTTPException(status_code=404, detail="Document not found")
    await ensure_document_folder(document["folder_id"], current_user, capability="read")

    now = datetime.now(timezone.utc)
    receipts = document.get("read_receipts", [])
    user_id = payload.user_id if user_is_admin(current_user) else current_user.id
    user_id = user_id or current_user.id
    updated = False

    for receipt in receipts:
        if receipt.get("user_id") == user_id:
            receipt["status"] = "read"
            receipt["read_at"] = now
            receipt["note"] = payload.note
            updated = True
            break

    if not updated:
        receipts.append(
            DocumentReadReceipt(user_id=user_id, status="read", read_at=now, note=payload.note).dict()
        )

    await db.documents.update_one(
        {"id": document_id},
        {"$set": {"read_receipts": receipts, "updated_at": now}},
    )
    updated_document = await db.documents.find_one({"id": document_id})
    await db.notifications.update_many(
        {
            "user_id": user_id,
            "is_read": False,
            "message": {"$regex": re.escape(document["code"])},
        },
        {"$set": {"is_read": True}},
    )
    return hydrate_document(updated_document)


@api_router.get("/documents/{document_id}/file", response_model=FileMetadata)
async def get_document_file_metadata(
    document_id: str,
    current_user: User = Depends(get_current_user),
) -> FileMetadata:
    document = await db.documents.find_one({"id": document_id})
    if not document:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Document not found")
    await ensure_document_folder(document["folder_id"], current_user, capability="read")
    file_id = document.get("file_id")
    if not file_id:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Document has no file attached")

    file_record = await db.file_uploads.find_one({"id": file_id})
    if not file_record:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="File not found")

    download_url = app.url_path_for("download_file", file_id=file_id)
    return FileMetadata(
        id=file_record["id"],
        original_filename=file_record["original_filename"],
        mime_type=file_record["mime_type"],
        file_size=file_record["file_size"],
        download_url=download_url,
    )


@api_router.get("/documents/{document_id}/preview", response_model=DocumentPreview)
async def preview_document(
    document_id: str,
    current_user: User = Depends(get_current_user),
) -> DocumentPreview:
    document = await db.documents.find_one({"id": document_id})
    if not document:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Document not found")
    await ensure_document_folder(document["folder_id"], current_user, capability="read")

    file_id = document.get("file_id")
    if not file_id:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Document has no file attached")

    file_record = await db.file_uploads.find_one({"id": file_id})
    if not file_record:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="File not found")

    return _generate_document_preview(document, file_record)


@api_router.get("/documents/approvals/pending", response_model=List[DocumentApprovalTask])
async def list_pending_document_approvals(current_user: User = Depends(get_current_user)):
    allowed_folders = await get_authorized_folder_ids(current_user, "read")
    query: Dict[str, Any] = {"approval_matrix": {"$ne": []}, "status": {"$in": ["review", "draft"]}}
    if allowed_folders is not None:
        if not allowed_folders:
            return []
        query["folder_id"] = {"$in": list(allowed_folders)}

    cursor = db.documents.find(query).sort("updated_at", -1).limit(200)
    tasks: List[DocumentApprovalTask] = []
    async for doc_raw in cursor:
        doc = hydrate_document(doc_raw)
        stage_index = find_pending_stage_index(doc)
        if stage_index is None:
            continue
        stage = doc.approval_matrix[stage_index]
        matching_token = resolve_matching_token(current_user, stage)
        if not matching_token:
            continue
        if any(dec.user_id == current_user.id for dec in stage.decisions):
            continue
        tasks.append(
            DocumentApprovalTask(
                document_id=doc.id,
                document_code=doc.code,
                title=doc.title,
                version=doc.version,
                stage=stage.stage,
                approval_type=stage.approval_type,
                deadline=stage.deadline,
                folder_id=doc.folder_id,
                status=stage.status,
                approvers=stage.approvers,
            )
        )
    return tasks


@api_router.post("/documents/{document_id}/approvals/decision", response_model=Document)
async def submit_document_approval_decision(
    document_id: str,
    payload: DocumentApprovalDecision,
    current_user: User = Depends(get_current_user),
):
    document_raw = await db.documents.find_one({"id": document_id})
    if not document_raw:
        raise HTTPException(status_code=404, detail="Document not found")
    await ensure_document_folder(document_raw["folder_id"], current_user, capability="read")

    document = hydrate_document(document_raw)
    if not document.approval_matrix:
        raise HTTPException(status_code=400, detail="Document does not require approval")

    stage_index = find_pending_stage_index(document)
    if stage_index is None:
        raise HTTPException(status_code=400, detail="No pending approval stage found")

    stage = document.approval_matrix[stage_index]
    matching_token = resolve_matching_token(current_user, stage)
    if not matching_token:
        raise HTTPException(status_code=403, detail="You are not part of the current approval stage")

    if any(dec.user_id == current_user.id for dec in stage.decisions):
        raise HTTPException(status_code=400, detail="Approval decision already recorded")

    now = datetime.now(timezone.utc)
    decision_entry = StageDecision(
        user_id=current_user.id,
        decision=payload.decision,
        comment=payload.comment,
        matched_token=matching_token,
        decided_at=now,
    )
    stage.decisions.append(decision_entry)

    status_entry: Optional[DocumentStatus] = None
    set_payload: Dict[str, Any] = {
        "updated_at": now,
    }

    if payload.decision == "rejected":
        stage.status = "rejected"
        stage.decided_by = current_user.id
        stage.decided_at = now
        stage.comment = payload.comment
        document.status = "draft"
        status_entry = DocumentStatus(
            status="draft",
            changed_by=current_user.id,
            comment=payload.comment or "Onay asamasi reddedildi.",
        )
        version_history = document_raw.get("version_history", [])
        for version in version_history:
            if version.get("id") == document_raw.get("current_version_id"):
                version["status"] = "draft"
                break
        set_payload["status"] = document.status
        set_payload["version_history"] = version_history
    else:
        if stage.approval_type == "any" or approval_stage_is_complete(stage):
            stage.status = "approved"
            stage.decided_by = current_user.id
            stage.decided_at = now
            stage.comment = payload.comment
            remaining_pending = any(s.status == "pending" for s in document.approval_matrix)
            if remaining_pending:
                document.status = "review"
            else:
                document.status = "approved"
                set_payload["published_at"] = now
                version_history = document_raw.get("version_history", [])
                for version in version_history:
                    if version.get("id") == document_raw.get("current_version_id"):
                        version["status"] = "published"
                        break
                set_payload["version_history"] = version_history
            set_payload["status"] = document.status
            status_entry = DocumentStatus(
                status=document.status,
                changed_by=current_user.id,
                comment=payload.comment or f"Asama {stage.stage} onaylandi.",
            )
        else:
            set_payload["status"] = document.status

    set_payload["approval_matrix"] = [stage_model.dict() for stage_model in document.approval_matrix]
    update_doc: Dict[str, Any] = {"$set": set_payload}
    if status_entry:
        update_doc.setdefault("$push", {})["status_history"] = status_entry.dict()

    await db.documents.update_one({"id": document_id}, update_doc)
    updated = await db.documents.find_one({"id": document_id})
    updated_document = hydrate_document(updated)

    if payload.decision == "rejected":
        await notify_users(
            [document.author_id],
            title="Dokuman Onayi Reddedildi",
            message=f"{document.code} / {document.title} {stage.stage}. asamada reddedildi.",
            notif_type="error",
        )
    else:
        next_stage_index = find_pending_stage_index(updated_document)
        if next_stage_index is not None:
            await notify_document_approvers(
                updated_document,
                updated_document.approval_matrix[next_stage_index],
                exclude_user_ids={current_user.id},
            )
        elif updated_document.status == "approved":
            await notify_document_readers(
                updated_document,
                note=f"{updated_document.code} / {updated_document.title} yayina alindi. Okuma onayiniz bekleniyor.",
            )
            await notify_users(
                [document.author_id],
                title="Dokuman Onaylandi",
                message=f"{document.code} / {document.title} yayina alindi.",
                notif_type="success",
            )

    await db.notifications.update_many(
        {
            "user_id": current_user.id,
            "is_read": False,
            "title": {"$regex": "Dokuman Onayi Bekliyor"},
            "message": {"$regex": re.escape(document.code)},
        },
        {"$set": {"is_read": True}},
    )

    return updated_document


@api_router.get("/documents/read-tasks", response_model=List[DocumentReadTask])
async def list_read_tasks(current_user: User = Depends(get_current_user)):
    allowed_folders = await get_authorized_folder_ids(current_user, "read")
    query: Dict[str, Any] = {
        "read_receipts": {
            "$elemMatch": {"user_id": current_user.id, "status": {"$ne": "read"}}
        }
    }
    if allowed_folders is not None:
        if not allowed_folders:
            return []
        query["folder_id"] = {"$in": list(allowed_folders)}

    cursor = db.documents.find(query).sort("updated_at", -1).limit(200)
    tasks: List[DocumentReadTask] = []
    async for doc_raw in cursor:
        doc = hydrate_document(doc_raw)
        receipt = next((r for r in doc.read_receipts if r.user_id == current_user.id and r.status != "read"), None)
        if not receipt:
            continue
        tasks.append(
            DocumentReadTask(
                document_id=doc.id,
                document_code=doc.code,
                title=doc.title,
                version=doc.version,
                required=receipt.required,
                status=receipt.status,
                folder_id=doc.folder_id,
            )
        )
    return tasks


@api_router.get("/documents/report/status", response_model=DocumentStatusReport)
async def get_document_status_report(
    department: Optional[str] = None,
    document_type: Optional[str] = None,
    current_user: User = Depends(get_current_user),
):
    return await build_document_status_report(current_user, department, document_type)


@api_router.get("/documents/report/status/export")
async def export_document_status_report(
    format: str = Query("xlsx", pattern="^(xlsx|pdf)$"),
    department: Optional[str] = None,
    document_type: Optional[str] = None,
    current_user: User = Depends(get_current_user),
):
    report = await build_document_status_report(current_user, department, document_type)
    timestamp = datetime.now(timezone.utc).strftime("%Y%m%d-%H%M%S")
    filename_base = f"document-status-report-{timestamp}"

    if format == "xlsx":
        buffer = BytesIO()
        with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
            status_df = pd.DataFrame(
                [{"Durum": key, "Adet": value} for key, value in report.status_counts.items()]
            )
            department_df = pd.DataFrame(
                [{"Departman": key, "Adet": value} for key, value in report.department_counts.items()]
            )
            type_df = pd.DataFrame(
                [{"Dokuman Turu": key, "Adet": value} for key, value in report.type_counts.items()]
            )
            summary_df = pd.DataFrame([{"Toplam Dokuman": report.total}])

            summary_df.to_excel(writer, sheet_name="Ozet", index=False)
            status_df.to_excel(writer, sheet_name="Durumlar", index=False)
            department_df.to_excel(writer, sheet_name="Departmanlar", index=False)
            type_df.to_excel(writer, sheet_name="DokumanTurleri", index=False)

        buffer.seek(0)
        return StreamingResponse(
            buffer,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            headers={
                "Content-Disposition": f'attachment; filename="{filename_base}.xlsx"'
            },
        )

    pdf_buffer = BytesIO()
    pdf = canvas.Canvas(pdf_buffer, pagesize=letter)
    width, height = letter
    margin = 50
    y = height - margin

    pdf.setFont("Helvetica-Bold", 14)
    pdf.drawString(margin, y, "Dokuman Durum Raporu")
    y -= 20

    pdf.setFont("Helvetica", 10)
    pdf.drawString(margin, y, f"Toplam Dokuman: {report.total}")
    y -= 20

    def draw_section(title: str, items: Dict[str, int], start_y: float) -> float:
        pdf.setFont("Helvetica-Bold", 12)
        pdf.drawString(margin, start_y, title)
        current_y = start_y - 16
        pdf.setFont("Helvetica", 10)
        if not items:
            pdf.drawString(margin, current_y, "Veri bulunmamaktadir.")
            return current_y - 20
        for key, value in items.items():
            pdf.drawString(margin, current_y, f"{key}: {value}")
            current_y -= 14
            if current_y < margin:
                pdf.showPage()
                current_y = height - margin
                pdf.setFont("Helvetica", 10)
        return current_y - 12

    y = draw_section("Durum Dagilimi", report.status_counts, y)
    y = draw_section("Departman Dagilimi", report.department_counts, y)
    draw_section("Dokuman Tur Dagilimi", report.type_counts, y)

    pdf.showPage()
    pdf.save()
    pdf_buffer.seek(0)

    return StreamingResponse(
        pdf_buffer,
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{filename_base}.pdf"'},
    )


# Complaint Routes


async def _get_complaint_or_404(complaint_id: str) -> Dict[str, Any]:
    complaint = await db.complaints.find_one({"id": complaint_id})
    if not complaint:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Complaint not found")
    return complaint


async def _get_complaint_category(category_id: str) -> Optional[ComplaintCategory]:
    if not category_id:
        return None
    category_raw = await db.complaint_categories.find_one({"id": category_id})
    if not category_raw:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Complaint category not found",
        )
    category = ComplaintCategory(**category_raw)
    if not category.is_active:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Complaint category is inactive",
        )
    return category


async def _sync_complaint_category_snapshot(record: Complaint) -> Complaint:
    if record.category_id:
        try:
            category = await _get_complaint_category(record.category_id)
        except HTTPException as exc:
            if exc.status_code == status.HTTP_404_NOT_FOUND:
                record.category_name = None
                return record
            raise
        if category:
            record.category_name = category.name
    else:
        record.category_name = None
    return record


def _merge_unique_values(existing: List[str], incoming: Optional[List[str]]) -> List[str]:
    if not incoming:
        return existing
    filtered_existing = [item for item in existing if item]
    filtered_incoming = [item for item in incoming if item]
    seen = set(filtered_existing)
    merged = list(filtered_existing)
    for item in filtered_incoming:
        if item not in seen:
            merged.append(item)
            seen.add(item)
    return merged


def _sanitize_string_list(values: Optional[Iterable[str]]) -> List[str]:
    if not values:
        return []
    cleaned: List[str] = []
    seen: Set[str] = set()
    for value in values:
        if value is None:
            continue
        text = str(value).strip()
        if not text:
            continue
        if text not in seen:
            cleaned.append(text)
            seen.add(text)
    return cleaned


def _normalize_optional_str(value: Optional[str]) -> Optional[str]:
    if value is None:
        return None
    text = str(value).strip()
    return text or None


def _normalize_required_str(value: str, field_name: str) -> str:
    normalized = _normalize_optional_str(value)
    if not normalized:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"{field_name} is required.",
        )
    return normalized


def _ensure_timezone(value: datetime) -> datetime:
    if value.tzinfo is None:
        return value.replace(tzinfo=timezone.utc)
    return value.astimezone(timezone.utc)


def _normalize_audit_status(value: str) -> str:
    normalized = _normalize_required_str(value, "Audit status").lower()
    if normalized not in ALLOWED_AUDIT_STATUSES:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Invalid audit status: {normalized}. Allowed values: {', '.join(sorted(ALLOWED_AUDIT_STATUSES))}",
        )
    return normalized


def _normalize_finding_type(value: str) -> str:
    normalized = _normalize_required_str(value, "Finding type").lower()
    if normalized not in ALLOWED_AUDIT_FINDING_TYPES:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Invalid finding type: {normalized}. Allowed values: {', '.join(sorted(ALLOWED_AUDIT_FINDING_TYPES))}",
        )
    return normalized


def _normalize_finding_status(value: str) -> str:
    normalized = _normalize_required_str(value, "Finding status").lower()
    if normalized not in ALLOWED_AUDIT_FINDING_STATUSES:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Invalid finding status: {normalized}. Allowed values: {', '.join(sorted(ALLOWED_AUDIT_FINDING_STATUSES))}",
        )
    return normalized


def _normalize_checklist_status(value: Optional[str]) -> str:
    normalized = (value or "pending").strip().lower()
    if normalized not in ALLOWED_AUDIT_CHECKLIST_STATUSES:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Invalid checklist status: {normalized}. Allowed values: {', '.join(sorted(ALLOWED_AUDIT_CHECKLIST_STATUSES))}",
        )
    return normalized


def _sanitize_audit_team_members(members: Optional[Iterable[AuditTeamMember]]) -> List[AuditTeamMember]:
    if not members:
        return []
    sanitized: List[AuditTeamMember] = []
    seen: Set[tuple[str, str]] = set()
    for member in members:
        user_id = _normalize_optional_str(member.user_id)
        role = _normalize_optional_str(member.role)
        if not user_id or not role:
            continue
        full_name = _normalize_optional_str(member.full_name)
        key = (user_id, role)
        if key in seen:
            continue
        seen.add(key)
        sanitized.append(
            AuditTeamMember(
                user_id=user_id,
                role=role,
                full_name=full_name,
            )
        )
    return sanitized


async def _prepare_checklist_items(items: Optional[Iterable[AuditChecklistItem]]) -> List[AuditChecklistItem]:
    if not items:
        return []
    checklist: List[AuditChecklistItem] = []
    pending_question_ids: Set[str] = set()
    original_items: List[AuditChecklistItem] = list(items)
    for item in original_items:
        question_id = _normalize_optional_str(item.question_id)
        question_text = _normalize_optional_str(item.question)
        if not question_text and not question_id:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Checklist item must include a question or question bank reference.",
            )
        if question_id and not question_text:
            pending_question_ids.add(question_id)
    question_lookup: Dict[str, str] = {}
    if pending_question_ids:
        docs = await db.audit_questions.find(
            {"id": {"$in": list(pending_question_ids)}, "is_active": True}
        ).to_list(None)
        question_lookup = {doc["id"]: doc["question"] for doc in docs}
    for item in original_items:
        question_id = _normalize_optional_str(item.question_id)
        question_text = _normalize_optional_str(item.question)
        if not question_text and question_id:
            question_text = question_lookup.get(question_id)
        if not question_text:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Referenced checklist question could not be resolved from the question bank.",
            )
        response = _normalize_optional_str(item.response)
        note = _normalize_optional_str(item.note)
        status_value = _normalize_checklist_status(item.status)
        checklist.append(
            AuditChecklistItem(
                id=item.id or str(uuid.uuid4()),
                question_id=question_id,
                question=question_text,
                response=response,
                status=status_value,
                note=note,
            )
        )
    return checklist


def _sanitize_audit_findings(findings: Optional[Iterable[AuditFinding]]) -> List[AuditFinding]:
    sanitized: List[AuditFinding] = []
    if not findings:
        return sanitized
    for finding in findings:
        finding_type = _normalize_finding_type(finding.finding_type)
        status_value = _normalize_finding_status(finding.status)
        description = _normalize_required_str(finding.description, "Finding description")
        requirement_reference = _normalize_optional_str(finding.requirement_reference)
        related_capa = _normalize_optional_str(finding.related_capa_id)
        corrective_action = _normalize_optional_str(finding.corrective_action)
        created_at = _ensure_timezone(finding.created_at)
        updated_at = _ensure_timezone(finding.updated_at)
        sanitized.append(
            AuditFinding(
                id=finding.id or str(uuid.uuid4()),
                finding_type=finding_type,
                description=description,
                requirement_reference=requirement_reference,
                related_capa_id=related_capa,
                status=status_value,
                corrective_action=corrective_action,
                created_at=created_at,
                updated_at=updated_at,
            )
        )
    return sanitized


def _sanitize_audit_status_history(entries: Optional[Iterable[AuditStatusHistory]]) -> List[AuditStatusHistory]:
    if not entries:
        return []
    sanitized: List[AuditStatusHistory] = []
    for entry in entries:
        status_value = _normalize_audit_status(entry.status)
        changed_by = _normalize_required_str(entry.changed_by, "Status change user")
        note = _normalize_optional_str(entry.note)
        changed_at = _ensure_timezone(entry.changed_at)
        sanitized.append(
            AuditStatusHistory(
                status=status_value,
                changed_by=changed_by,
                note=note,
                changed_at=changed_at,
            )
        )
    return sanitized


async def _get_audit_or_404(audit_id: str) -> Dict[str, Any]:
    audit = await db.audits.find_one({"id": audit_id})
    if not audit:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Audit record not found.")
    return audit


async def _persist_audit_record(record: Audit) -> Audit:
    record.audit_team = _sanitize_audit_team_members(record.audit_team)
    record.checklist = await _prepare_checklist_items(record.checklist)
    record.findings = _sanitize_audit_findings(record.findings)
    record.status_history = _sanitize_audit_status_history(record.status_history)
    record.updated_at = datetime.now(timezone.utc)
    await db.audits.replace_one({"id": record.id}, record.dict(), upsert=True)
    return record


async def _add_finding_link_to_capa(capa_id: str, finding_id: str) -> None:
    capa = await db.capas.find_one({"id": capa_id})
    if not capa:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="CAPA not found.")
    linked = _sanitize_string_list(capa.get("linked_audit_finding_ids") or [])
    if finding_id not in linked:
        linked.append(finding_id)
        await db.capas.update_one(
            {"id": capa_id},
            {
                "$set": {
                    "linked_audit_finding_ids": linked,
                    "updated_at": datetime.now(timezone.utc),
                }
            },
        )


async def _remove_finding_link_from_capa(capa_id: str, finding_id: str) -> None:
    if not capa_id:
        return
    capa = await db.capas.find_one({"id": capa_id})
    if not capa:
        return
    linked = _sanitize_string_list(capa.get("linked_audit_finding_ids") or [])
    if finding_id in linked:
        linked.remove(finding_id)
        await db.capas.update_one(
            {"id": capa_id},
            {
                "$set": {
                    "linked_audit_finding_ids": linked,
                    "updated_at": datetime.now(timezone.utc),
                }
            },
        )


async def _sync_finding_capa_link(finding_id: str, new_capa_id: Optional[str], previous_capa_id: Optional[str]) -> None:
    normalized_new = _normalize_optional_str(new_capa_id)
    normalized_previous = _normalize_optional_str(previous_capa_id)
    if normalized_previous and normalized_previous != normalized_new:
        await _remove_finding_link_from_capa(normalized_previous, finding_id)
    if normalized_new:
        await _add_finding_link_to_capa(normalized_new, finding_id)


def _normalize_risk_status(value: str) -> str:
    normalized = _normalize_required_str(value, "Risk status").lower()
    if normalized not in ALLOWED_RISK_STATUSES:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Invalid risk status: {normalized}. Allowed values: {', '.join(sorted(ALLOWED_RISK_STATUSES))}",
        )
    return normalized


def _sanitize_risk_factors(factors: Optional[Iterable[RiskFactor]]) -> List[RiskFactor]:
    sanitized: List[RiskFactor] = []
    if not factors:
        return sanitized
    for factor in factors:
        name = _normalize_required_str(factor.name, "Risk factor name")
        sanitized.append(
            RiskFactor(
                name=name,
                value=float(factor.value),
                weight=float(factor.weight or 1.0),
            )
        )
    return sanitized


SAFE_FORMULA_GLOBALS = {
    "min": min,
    "max": max,
    "abs": abs,
    "round": round,
    "ceil": math.ceil,
    "floor": math.floor,
    "sqrt": math.sqrt,
}


async def _load_risk_settings() -> RiskSettings:
    doc = await db.risk_settings.find_one({"id": "default"})
    if doc:
        try:
            return RiskSettings(**doc)
        except Exception:
            pass
    return RiskSettings()


async def _persist_risk_settings(settings: RiskSettings) -> RiskSettings:
    settings.updated_at = datetime.now(timezone.utc)
    await db.risk_settings.replace_one({"id": settings.id}, settings.dict(), upsert=True)
    return settings


async def _load_risk_report_template() -> RiskReportTemplate:
    doc = await db.risk_report_templates.find_one({"id": "default"})
    if doc:
        try:
            return RiskReportTemplate(**doc)
        except Exception:
            pass
    return RiskReportTemplate()


async def _persist_risk_report_template(template: RiskReportTemplate) -> RiskReportTemplate:
    template.updated_at = datetime.now(timezone.utc)
    await db.risk_report_templates.replace_one({"id": template.id}, template.dict(), upsert=True)
    return template


def _determine_risk_level(score: float, thresholds: Dict[str, float]) -> str:
    ordered = [
        ("low", thresholds.get("low", 5)),
        ("medium", thresholds.get("medium", 12)),
        ("high", thresholds.get("high", 20)),
        ("critical", thresholds.get("critical", 25)),
    ]
    for level, limit in ordered:
        if limit is None:
            continue
        if score <= float(limit):
            return level
    return "critical"


def _evaluate_formula(expression: str, context: Dict[str, Any]) -> float:
    env = dict(SAFE_FORMULA_GLOBALS)
    env.update(context)
    try:
        result = eval(expression, {"__builtins__": {}}, env)  # pylint: disable=eval-used
    except Exception as exc:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Error evaluating risk formula: {exc}",
        ) from exc
    try:
        return float(result)
    except (TypeError, ValueError) as exc:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Risk formula did not return a numeric value.",
        ) from exc


def _clamp_matrix_index(value: float) -> int:
    return max(1, min(DEFAULT_RISK_MATRIX_SIZE, int(round(value))))


def _calculate_risk_score(
    settings: RiskSettings,
    likelihood: float,
    impact: float,
    detection: Optional[float],
    controls_effectiveness: Optional[float],
    factors: List[RiskFactor],
) -> RiskScore:
    total_factor = sum(f.value * f.weight for f in factors)
    context = {
        "likelihood": float(likelihood),
        "impact": float(impact),
        "detection": float(detection) if detection is not None else 1.0,
        "controls_effectiveness": float(controls_effectiveness) if controls_effectiveness is not None else 0.0,
        "total_factor": total_factor,
    }
    inherent = max(_evaluate_formula(settings.formula, context), 0.0)
    residual_formula = settings.residual_formula or settings.formula
    residual = max(_evaluate_formula(residual_formula, context), 0.0)
    inherent_level = _determine_risk_level(inherent, settings.thresholds)
    residual_level = _determine_risk_level(residual, settings.thresholds)
    row = _clamp_matrix_index(likelihood)
    col = _clamp_matrix_index(impact)
    return RiskScore(
        inherent=inherent,
        residual=residual,
        inherent_level=inherent_level,
        residual_level=residual_level,
        matrix_row=row,
        matrix_col=col,
    )


def _build_risk_matrix_map(settings: RiskSettings) -> Dict[str, RiskMatrixCell]:
    mapping: Dict[str, RiskMatrixCell] = {}
    matrix = settings.matrix
    for row_idx, row in enumerate(matrix, start=1):
        for col_idx, level in enumerate(row, start=1):
            color = settings.palette.get(level, "#cccccc")
            mapping[f"{row_idx}-{col_idx}"] = RiskMatrixCell(
                row=row_idx,
                col=col_idx,
                level=level,
                color=color,
                label=level.title(),
            )
    return mapping


async def _get_risk_or_404(risk_id: str) -> Dict[str, Any]:
    record = await db.risks.find_one({"id": risk_id})
    if not record:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Risk record not found.")
    return record


async def _persist_risk_record(record: RiskAssessment) -> RiskAssessment:
    record.custom_factors = _sanitize_risk_factors(record.custom_factors)
    record.linked_capa_ids = _sanitize_string_list(record.linked_capa_ids)
    record.linked_audit_finding_ids = _sanitize_string_list(record.linked_audit_finding_ids)
    record.updated_at = datetime.now(timezone.utc)
    await db.risks.replace_one({"id": record.id}, record.dict(), upsert=True)
    return record


async def _append_risk_revision(record: RiskAssessment, user_id: str, note: Optional[str] = None) -> RiskAssessment:
    snapshot = record.dict()
    revision_no = (record.revisions[-1].revision_no + 1) if record.revisions else 1
    revision = RiskRevision(
        revision_no=revision_no,
        snapshot=snapshot,
        changed_by=user_id,
        changed_at=datetime.now(timezone.utc),
        note=note,
    )
    record.revisions.append(revision)
    return record


async def _record_risk_trend_point(record: RiskAssessment) -> None:
    point = RiskTrendPoint(
        inherent_score=record.risk_score.inherent,
        residual_score=record.risk_score.residual,
        status=record.status,
    )
    await db.risk_trends.insert_one(
        {
            "risk_id": record.id,
            "recorded_at": point.recorded_at,
            "inherent_score": point.inherent_score,
            "residual_score": point.residual_score,
            "status": point.status,
        }
    )


async def _fetch_risk_trend(risk_id: str, limit: int = 20) -> List[RiskTrendPoint]:
    cursor = (
        db.risk_trends.find({"risk_id": risk_id})
        .sort("recorded_at", -1)
        .limit(limit)
    )
    points: List[RiskTrendPoint] = []
    async for item in cursor:
        points.append(
            RiskTrendPoint(
                recorded_at=item["recorded_at"],
                inherent_score=item["inherent_score"],
                residual_score=item["residual_score"],
                status=item.get("status", "identified"),
            )
        )
    return list(reversed(points))


async def _build_risk_matrix_summary(settings: RiskSettings) -> Dict[str, Any]:
    mapping = _build_risk_matrix_map(settings)
    counts: Dict[str, int] = {key: 0 for key in mapping.keys()}
    async for risk in db.risks.find({}):
        row = risk.get("risk_score", {}).get("matrix_row", 1)
        col = risk.get("risk_score", {}).get("matrix_col", 1)
        key = f"{row}-{col}"
        if key in counts:
            counts[key] += 1
    matrix = []
    for row in range(1, DEFAULT_RISK_MATRIX_SIZE + 1):
        matrix_row = []
        for col in range(1, DEFAULT_RISK_MATRIX_SIZE + 1):
            cell_key = f"{row}-{col}"
            cell = mapping[cell_key]
            matrix_row.append(
                {
                    "row": row,
                    "col": col,
                    "level": cell.level,
                    "color": cell.color,
                    "label": cell.label,
                    "count": counts.get(cell_key, 0),
                }
            )
        matrix.append(matrix_row)
    return {"matrix": matrix, "palette": settings.palette}


def _build_risk_diff(snapshot_a: Dict[str, Any], snapshot_b: Dict[str, Any]) -> Dict[str, Dict[str, Any]]:
    diff: Dict[str, Dict[str, Any]] = {}
    keys = set(snapshot_a.keys()) | set(snapshot_b.keys())
    for key in keys:
        if snapshot_a.get(key) != snapshot_b.get(key):
            diff[key] = {"from": snapshot_a.get(key), "to": snapshot_b.get(key)}
    return diff


def _normalize_device_status(value: str) -> str:
    normalized = _normalize_required_str(value, "Device status").lower()
    if normalized not in CALIBRATION_DEVICE_STATUSES:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Invalid device status: {normalized}. Allowed values: {', '.join(sorted(CALIBRATION_DEVICE_STATUSES))}",
        )
    return normalized


def _normalize_work_order_status(value: str) -> str:
    normalized = _normalize_required_str(value, "Work order status").lower()
    if normalized not in CALIBRATION_WORK_ORDER_STATUSES:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Invalid work order status: {normalized}. Allowed values: {', '.join(sorted(CALIBRATION_WORK_ORDER_STATUSES))}",
        )
    return normalized


def _validate_device_transition(current_status: Optional[str], new_status: str) -> None:
    if current_status is None or current_status == new_status:
        return
    allowed = CALIBRATION_DEVICE_TRANSITIONS.get(current_status, set())
    if new_status not in allowed:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=(
                f"Transition from '{current_status}' to '{new_status}' is not allowed for calibration devices."
            ),
        )


def _validate_work_order_transition(current_status: Optional[str], new_status: str) -> None:
    if current_status is None or current_status == new_status:
        return
    allowed = CALIBRATION_WORK_ORDER_TRANSITIONS.get(current_status, set())
    if new_status not in allowed:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=(
                f"Transition from '{current_status}' to '{new_status}' is not allowed for calibration work orders."
            ),
        )


def _append_device_status_history(
    device: CalibrationDevice,
    new_status: str,
    changed_by: str,
    note: Optional[str] = None,
) -> None:
    normalized_note = _normalize_optional_str(note)
    if device.status_history and device.status_history[-1].status == new_status:
        if normalized_note and device.status_history[-1].note != normalized_note:
            device.status_history.append(
                CalibrationDeviceStatusChange(
                    status=new_status,
                    changed_by=changed_by,
                    note=normalized_note,
                )
            )
        return
    device.status_history.append(
        CalibrationDeviceStatusChange(
            status=new_status,
            changed_by=changed_by,
            note=normalized_note,
        )
    )


def _append_work_order_status_history(
    work_order: CalibrationWorkOrder,
    new_status: str,
    changed_by: str,
    note: Optional[str] = None,
) -> None:
    normalized_note = _normalize_optional_str(note)
    if work_order.status_history and work_order.status_history[-1].status == new_status:
        if normalized_note and work_order.status_history[-1].note != normalized_note:
            work_order.status_history.append(
                CalibrationWorkOrderStatusChange(
                    status=new_status,
                    changed_by=changed_by,
                    note=normalized_note,
                )
            )
        return
    work_order.status_history.append(
        CalibrationWorkOrderStatusChange(
            status=new_status,
            changed_by=changed_by,
            note=normalized_note,
        )
    )


def _normalize_report_module(value: str) -> str:
    module = _normalize_required_str(value, "Report module").lower()
    if module not in REPORT_ALLOWED_MODULES:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Unsupported report module: {module}. Allowed modules: {', '.join(sorted(REPORT_ALLOWED_MODULES))}",
        )
    return module


def _normalize_report_operator(value: Optional[str]) -> str:
    operator = (value or "eq").lower()
    if operator not in REPORT_ALLOWED_OPERATORS:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Unsupported report filter operator: {operator}.",
        )
    return operator


def _get_collection_for_module(module: str):
    normalized = _normalize_report_module(module)
    collection_map = {
        "documents": db.documents,
        "document_versions": db.document_versions,
        "calibration_devices": db.calibration_devices,
        "calibration_work_orders": db.calibration_work_orders,
        "risks": db.risks,
        "capas": db.capas,
        "complaints": db.complaints,
        "dof_tasks": db.dof_tasks,
        "audit_findings": db.audit_findings,
    }
    collection = collection_map.get(normalized)
    if collection is None:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"No collection configured for report module: {normalized}.",
        )
    return collection


def _coerce_report_filter_value(
    definition: ReportFilterDefinition,
    operator: str,
    value: Any,
) -> Any:
    field_type = (definition.field_type or "string").lower()

    def coerce_single(item: Any) -> Any:
        if item is None:
            return None
        if field_type in {"string", "text"}:
            return _normalize_optional_str(str(item))
        if field_type in {"number", "numeric", "float", "int", "integer"}:
            if isinstance(item, (int, float)):
                return float(item)
            text = str(item).strip()
            if not text:
                return None
            try:
                return float(text)
            except ValueError as exc:
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail=f"Value '{text}' is not a valid number for filter '{definition.field}'.",
                ) from exc
        if field_type in {"bool", "boolean"}:
            if isinstance(item, bool):
                return item
            if isinstance(item, (int, float)):
                return bool(item)
            text = str(item).strip().lower()
            if text in {"true", "1", "yes"}:
                return True
            if text in {"false", "0", "no"}:
                return False
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Value '{item}' is not a valid boolean for filter '{definition.field}'.",
            )
        if field_type in {"date", "datetime"}:
            if isinstance(item, datetime):
                return _ensure_datetime(item)
            text = str(item).strip()
            try:
                parsed = _parse_iso_datetime(text)
            except Exception as exc:
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail=f"Value '{item}' is not a valid ISO datetime for filter '{definition.field}'.",
                ) from exc
            return parsed
        if field_type in {"list", "array"}:
            if isinstance(item, list):
                return [_normalize_optional_str(str(v)) for v in item if v is not None]
            return [_normalize_optional_str(str(item))]
        return item

    if operator == "between":
        if not isinstance(value, (list, tuple)) or len(value) != 2:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Filter '{definition.field}' expects a two-value array for 'between' operator.",
            )
        coerced = [coerce_single(value[0]), coerce_single(value[1])]
        return coerced

    if operator in {"in", "nin"} or definition.allow_multiple:
        if not isinstance(value, (list, tuple)):
            value = [value]
        coerced_list = [coerce_single(item) for item in value]
        return [item for item in coerced_list if item is not None]

    if operator == "exists":
        return bool(value)

    return coerce_single(value)


def _merge_report_filters(
    definition: ReportDefinition,
    runtime_filters: Optional[List[ReportFilterValue]],
) -> List[ReportFilterValue]:
    merged: List[ReportFilterValue] = []
    runtime_map: Dict[str, ReportFilterValue] = {}
    for item in runtime_filters or []:
        runtime_map[item.field] = item

    for filter_def in definition.filters:
        runtime_value = runtime_map.pop(filter_def.field, None)
        if runtime_value:
            operator = _normalize_report_operator(runtime_value.operator or filter_def.operator)
            coerced = _coerce_report_filter_value(filter_def, operator, runtime_value.value)
            merged.append(ReportFilterValue(field=filter_def.field, operator=operator, value=coerced))
            continue
        if filter_def.default is not None:
            operator = _normalize_report_operator(filter_def.operator)
            coerced = _coerce_report_filter_value(filter_def, operator, filter_def.default)
            merged.append(ReportFilterValue(field=filter_def.field, operator=operator, value=coerced))
        elif filter_def.required:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Report filter '{filter_def.label or filter_def.field}' is required.",
            )

    if runtime_map:
        unknown = ", ".join(sorted(runtime_map.keys()))
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Unknown report filters provided: {unknown}.",
        )
    return merged


def _apply_query_condition(
    query: Dict[str, Any],
    field: str,
    operator: str,
    value: Any,
) -> None:
    if operator == "eq":
        if field in query and isinstance(query[field], dict):
            query[field]["$eq"] = value
        else:
            query[field] = value
        return

    if field not in query or not isinstance(query[field], dict):
        existing = query.get(field)
        query[field] = {} if isinstance(existing, dict) else {}
        if existing is not None and not isinstance(existing, dict):
            query[field]["$eq"] = existing

    target = query[field]
    if operator == "ne":
        target["$ne"] = value
    elif operator == "in":
        target["$in"] = value or []
    elif operator == "nin":
        target["$nin"] = value or []
    elif operator == "gte":
        target["$gte"] = value
    elif operator == "lte":
        target["$lte"] = value
    elif operator == "gt":
        target["$gt"] = value
    elif operator == "lt":
        target["$lt"] = value
    elif operator == "between":
        target["$gte"] = value[0]
        target["$lte"] = value[1]
    elif operator == "contains":
        if value is None:
            return
        target["$regex"] = re.escape(str(value))
        target["$options"] = "i"
    elif operator == "regex":
        if value is None:
            return
        target["$regex"] = str(value)
    elif operator == "exists":
        target["$exists"] = bool(value)


def _build_report_query(filters: List[ReportFilterValue]) -> Dict[str, Any]:
    query: Dict[str, Any] = {}
    for filter_value in filters:
        operator = _normalize_report_operator(filter_value.operator)
        _apply_query_condition(query, filter_value.field, operator, filter_value.value)
    return query


def _jsonify_report_value(value: Any) -> Any:
    if isinstance(value, datetime):
        return _ensure_datetime(value).isoformat()
    if isinstance(value, uuid.UUID):
        return str(value)
    if isinstance(value, list):
        return [_jsonify_report_value(item) for item in value]
    if isinstance(value, tuple):
        return [_jsonify_report_value(item) for item in value]
    if isinstance(value, dict):
        return {key: _jsonify_report_value(val) for key, val in value.items()}
    return value


def _sanitize_report_rows(
    rows: List[Dict[str, Any]],
    fields: Optional[List[str]],
) -> List[Dict[str, Any]]:
    sanitized: List[Dict[str, Any]] = []
    for row in rows:
        data = {key: value for key, value in row.items() if key != "_id"}
        if fields:
            data = {key: data.get(key) for key in fields if key in data}
        sanitized.append({key: _jsonify_report_value(value) for key, value in data.items()})
    return sanitized


def _collect_metric_values(rows: List[Dict[str, Any]], field: Optional[str]) -> List[Any]:
    if not field:
        return [1 for _ in rows]
    values: List[Any] = []
    for row in rows:
        if field not in row:
            continue
        value = row[field]
        if value is None:
            continue
        values.append(value)
    return values


def _coerce_metric_values(values: List[Any]) -> List[float]:
    numeric: List[float] = []
    for value in values:
        if isinstance(value, bool):
            numeric.append(1.0 if value else 0.0)
        elif isinstance(value, (int, float)):
            numeric.append(float(value))
        elif isinstance(value, str):
            text = value.strip()
            if not text:
                continue
            try:
                numeric.append(float(text))
            except ValueError:
                continue
    return numeric


def _apply_metric(rows: List[Dict[str, Any]], metric: ReportMetricConfig) -> Any:
    operation = (metric.operation or "count").lower()
    field = metric.field
    values = _collect_metric_values(rows, field)

    if operation == "count":
        if field:
            return sum(1 for value in values if value is not None)
        return len(rows)

    numeric_values = _coerce_metric_values(values)
    if not numeric_values:
        return None

    if operation == "sum":
        return sum(numeric_values)
    if operation == "avg":
        return sum(numeric_values) / len(numeric_values)
    if operation == "min":
        return min(numeric_values)
    if operation == "max":
        return max(numeric_values)

    raise HTTPException(
        status_code=status.HTTP_400_BAD_REQUEST,
        detail=f"Unsupported report metric operation: {metric.operation}.",
    )


def _format_metric_value(value: Any, precision: Optional[int]) -> Any:
    if value is None:
        return None
    if isinstance(value, float) and precision is not None:
        return round(value, precision)
    return value


def _build_chart_series(
    grouped_rows: Dict[Any, List[Dict[str, Any]]],
    metrics: List[ReportMetricConfig],
) -> List[Dict[str, Any]]:
    series: List[Dict[str, Any]] = []
    for metric in metrics:
        data_points = []
        for group_key, group_rows in grouped_rows.items():
            metric_value = _apply_metric(group_rows, metric)
            data_points.append(
                {
                    "group": _jsonify_report_value(group_key),
                    "value": _format_metric_value(metric_value, metric.precision),
                }
            )
        series.append(
            {
                "metric": metric.label or metric.name or metric.operation,
                "data": data_points,
            }
        )
    return series


def _build_visualization_payload(
    rows: List[Dict[str, Any]],
    config: ReportVisualizationConfig,
) -> Dict[str, Any]:
    vis_type = (config.type or "table").lower()
    metrics = config.metrics or [ReportMetricConfig(operation="count")]
    payload: Dict[str, Any] = {"type": vis_type, "title": config.title}

    if vis_type == "table":
        fields = config.fields or sorted({key for row in rows for key in row.keys()})
        payload.update(
            {
                "fields": fields,
                "rows": [
                    {field: row.get(field) for field in fields}
                    for row in rows
                ],
            }
        )
        return payload

    if vis_type == "kpi":
        metrics_payload = []
        for metric in metrics:
            value = _apply_metric(rows, metric)
            metrics_payload.append(
                {
                    "metric": metric.label or metric.name or metric.operation,
                    "value": _format_metric_value(value, metric.precision),
                }
            )
        payload["metrics"] = metrics_payload
        return payload

    if config.group_by is None:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Visualization '{config.title or config.type}' requires a group_by field.",
        )

    grouped: Dict[Any, List[Dict[str, Any]]] = defaultdict(list)
    for row in rows:
        key = row.get(config.group_by)
        grouped[key].append(row)

    payload["group_by"] = config.group_by
    payload["series"] = _build_chart_series(grouped, metrics)
    return payload


def _build_report_visualizations(
    rows: List[Dict[str, Any]],
    configs: Optional[List[ReportVisualizationConfig]],
) -> List[Dict[str, Any]]:
    if not configs:
        return []
    return [_build_visualization_payload(rows, config) for config in configs]


def _build_report_summary(
    rows: List[Dict[str, Any]],
    configs: Optional[List[ReportVisualizationConfig]],
) -> Dict[str, Any]:
    summary: Dict[str, Any] = {"total_rows": len(rows)}
    if not configs:
        return summary
    for config in configs:
        if (config.type or "").lower() != "kpi":
            continue
        for metric in config.metrics or []:
            key = (metric.label or metric.name or metric.operation).lower().replace(" ", "_")
            summary[key] = _format_metric_value(_apply_metric(rows, metric), metric.precision)
    return summary


async def _get_report_definition_or_404(definition_id: str) -> Dict[str, Any]:
    record = await db.report_definitions.find_one({"id": definition_id})
    if not record:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Report definition not found.")
    return record


async def _persist_report_definition(definition: ReportDefinition) -> ReportDefinition:
    definition.updated_at = datetime.now(timezone.utc)
    await db.report_definitions.replace_one({"id": definition.id}, definition.dict(), upsert=True)
    return definition


async def _get_report_template_or_404(template_id: str) -> Dict[str, Any]:
    record = await db.report_templates.find_one({"id": template_id})
    if not record:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Report template not found.")
    return record


async def _persist_report_template(template: ReportTemplate) -> ReportTemplate:
    template.updated_at = datetime.now(timezone.utc)
    await db.report_templates.replace_one({"id": template.id}, template.dict(), upsert=True)
    return template


async def _execute_report_run(
    definition: ReportDefinition,
    filters: Optional[List[ReportFilterValue]],
    visualization_overrides: Optional[List[ReportVisualizationConfig]] = None,
) -> ReportRunResult:
    merged_filters = _merge_report_filters(definition, filters)
    query = _build_report_query(merged_filters)
    collection = _get_collection_for_module(definition.module)
    limit = max(1, definition.max_rows or 2000)

    cursor = collection.find(query)
    rows = await cursor.limit(limit).to_list(length=limit)
    sanitized_rows = _sanitize_report_rows(rows, definition.fields)

    visual_configs = visualization_overrides or definition.visualizations
    visualizations = _build_report_visualizations(sanitized_rows, visual_configs)
    summary = _build_report_summary(sanitized_rows, visual_configs or definition.visualizations)
    total_count = await collection.count_documents(query)

    return ReportRunResult(
        definition=definition,
        filters=merged_filters,
        rows=sanitized_rows,
        total=total_count,
        visualizations=visualizations,
        summary=summary,
    )


async def _ensure_report_indexes() -> None:
    logger = logging.getLogger(__name__)
    await db.report_definitions.create_index("module", name="idx_report_def_module")
    await db.report_definitions.create_index("name", unique=True, name="uniq_report_def_name")
    await db.report_definitions.create_index(
        [("default_template_id", 1)],
        name="idx_report_def_default_template",
        sparse=True,
    )
    await db.report_templates.create_index("module", name="idx_report_tpl_module")
    await db.report_templates.create_index("name", unique=True, name="uniq_report_tpl_name")
    await db.report_templates.create_index(
        [("default_definition_id", 1)],
        name="idx_report_tpl_default_definition",
        sparse=True,
    )
    logger.info("Report indexes ensured.")


async def _seed_report_assets() -> None:
    logger = logging.getLogger(__name__)
    now = datetime.now(timezone.utc)

    for template_data in SEED_REPORT_TEMPLATES:
        existing_template = await db.report_templates.find_one({"id": template_data["id"]})
        if existing_template:
            continue
        placeholders = [
            ReportTemplatePlaceholder(**placeholder)
            for placeholder in template_data.get("placeholders", [])
        ]
        template_model = ReportTemplate(
            id=template_data["id"],
            name=template_data["name"],
            module=template_data["module"],
            template_type=template_data["template_type"],
            description=template_data.get("description"),
            file_upload_id=template_data.get("file_upload_id"),
            placeholders=placeholders,
            default_definition_id=template_data.get("default_definition_id"),
            created_by=SYSTEM_USER_ID,
            created_at=now,
            updated_at=now,
        )
        await db.report_templates.insert_one(template_model.dict())
        logger.info("Seeded report template %s", template_model.id)

    for definition_data in SEED_REPORT_DEFINITIONS:
        existing_definition = await db.report_definitions.find_one({"id": definition_data["id"]})
        if existing_definition:
            continue
        filters = [
            ReportFilterDefinition(**filter_definition)
            for filter_definition in definition_data.get("filters", [])
        ]
        visualizations: List[ReportVisualizationConfig] = []
        for visualization in definition_data.get("visualizations", []):
            metrics = [
                ReportMetricConfig(**metric_definition)
                for metric_definition in visualization.get("metrics", [])
            ]
            vis_model = ReportVisualizationConfig(
                type=visualization.get("type", "table"),
                title=visualization.get("title"),
                group_by=visualization.get("group_by"),
                metrics=metrics,
                fields=visualization.get("fields"),
                options=visualization.get("options", {}) or {},
            )
            visualizations.append(vis_model)
        fields = _sanitize_string_list(definition_data.get("fields"))
        definition_model = ReportDefinition(
            id=definition_data["id"],
            name=definition_data["name"],
            module=definition_data["module"],
            description=definition_data.get("description"),
            filters=filters,
            visualizations=visualizations,
            fields=fields or None,
            max_rows=max(1, min(int(definition_data.get("max_rows", 2000)), 20000)),
            default_template_id=definition_data.get("default_template_id"),
            created_by=SYSTEM_USER_ID,
            created_at=now,
            updated_at=now,
        )
        await db.report_definitions.insert_one(definition_model.dict())
        logger.info("Seeded report definition %s", definition_model.id)


def _ensure_datetime(value: Optional[datetime]) -> Optional[datetime]:
    if value is None:
        return None
    if value.tzinfo is None:
        return value.replace(tzinfo=timezone.utc)
    return value.astimezone(timezone.utc)


def _compute_next_due_date(
    last_calibrated_at: Optional[datetime],
    interval_days: int,
    fallback: Optional[datetime] = None,
) -> Optional[datetime]:
    if interval_days <= 0:
        return fallback
    base = last_calibrated_at or fallback
    if not base:
        return fallback
    return base + timedelta(days=interval_days)


async def _get_device_or_404(device_id: str) -> Dict[str, Any]:
    record = await db.calibration_devices.find_one({"id": device_id})
    if not record:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Calibration device not found.")
    return record


async def _persist_device_record(record: CalibrationDevice) -> CalibrationDevice:
    record.linked_capa_ids = _sanitize_string_list(record.linked_capa_ids)
    record.linked_dof_task_ids = _sanitize_string_list(record.linked_dof_task_ids)
    record.file_attachments = _sanitize_string_list(record.file_attachments)
    record.updated_at = datetime.now(timezone.utc)
    await db.calibration_devices.replace_one({"id": record.id}, record.dict(), upsert=True)
    return record


async def _get_work_order_or_404(work_order_id: str) -> Dict[str, Any]:
    record = await db.calibration_work_orders.find_one({"id": work_order_id})
    if not record:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Calibration work order not found.")
    return record


async def _persist_work_order_record(record: CalibrationWorkOrder) -> CalibrationWorkOrder:
    record.linked_capa_ids = _sanitize_string_list(record.linked_capa_ids)
    record.updated_at = datetime.now(timezone.utc)
    await db.calibration_work_orders.replace_one({"id": record.id}, record.dict(), upsert=True)
    return record


async def _generate_device_code() -> str:
    count = await db.calibration_devices.count_documents({})
    return f"DEV-{datetime.now().year}-{count + 1:04d}"


async def _generate_work_order_no() -> str:
    count = await db.calibration_work_orders.count_documents({})
    return f"CWO-{datetime.now().year}-{count + 1:04d}"


def _work_order_cost_total(record: CalibrationWorkOrder) -> float:
    return float(sum(entry.amount for entry in record.cost_entries))


async def _collect_calibration_summary() -> Dict[str, Any]:
    devices = await db.calibration_devices.find({}).to_list(2000)
    work_orders = await db.calibration_work_orders.find({}).to_list(5000)
    now = datetime.now(timezone.utc)
    overdue_devices = 0
    upcoming_devices = 0
    total_cost = 0.0
    for device in devices:
        next_due = device.get("next_due_date")
        interval = device.get("calibration_interval_days", 0)
        notice_days = device.get("notice_days", 14)
        if next_due:
            if isinstance(next_due, datetime):
                due_dt = next_due
            else:
                due_dt = _parse_iso_datetime(str(next_due))
            if due_dt < now:
                overdue_devices += 1
            elif due_dt <= now + timedelta(days=notice_days or 14):
                upcoming_devices += 1
    for wo in work_orders:
        for cost in wo.get("cost_entries", []):
            total_cost += float(cost.get("amount") or 0.0)

    return {
        "total_devices": len(devices),
        "active_devices": sum(1 for device in devices if (device.get("status") or "").lower() == "active"),
        "overdue_devices": overdue_devices,
        "upcoming_devices": upcoming_devices,
        "total_work_orders": len(work_orders),
        "open_work_orders": sum(
            1 for wo in work_orders if (wo.get("status") or "").lower() in {"planned", "in_progress"}
        ),
        "cost_total": total_cost,
    }


async def _create_dof_for_work_order(
    device: CalibrationDevice,
    work_order: CalibrationWorkOrder,
    request: CalibrationDofRequest,
    current_user: User,
) -> DofTask:
    department = _normalize_optional_str(request.department) or device.department or "Kalibrasyon"
    responsible = (
        _normalize_optional_str(request.responsible_person)
        or device.responsible_person
        or current_user.full_name
        or current_user.username
    )
    if not responsible:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Responsible person is required to create a DOF task.",
        )
    due_date = request.due_date or (datetime.now(timezone.utc) + timedelta(days=7))
    if due_date.tzinfo is None:
        due_date = due_date.replace(tzinfo=timezone.utc)

    dof_no = await generate_dof_no()
    now = datetime.now(timezone.utc)
    task = DofTask(
        dof_no=dof_no,
        title=request.title or f"Calibration action for {device.name}",
        description=(
            f"Calibration work order {work_order.work_order_no} for device {device.device_code}."
        ),
        department=department,
        responsible_person=responsible,
        due_date=due_date,
        created_by=current_user.id,
        status_history=[
            DofStatusHistory(status="open", changed_by=current_user.id, changed_at=now)
        ],
    )
    await db.dof_tasks.insert_one(task.dict())
    return task


async def _persist_complaint_record(record: Complaint) -> Complaint:
    record = await _sync_complaint_category_snapshot(record)
    record.updated_at = datetime.now(timezone.utc)
    result = await db.complaints.replace_one({"id": record.id}, record.dict())
    if result.matched_count == 0:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Complaint not found")
    return record


async def _get_capa_or_404(capa_id: str) -> Dict[str, Any]:
    capa = await db.capas.find_one({"id": capa_id})
    if not capa:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="CAPA not found")
    return capa


async def _persist_capa_record(record: Capa) -> Capa:
    record.team_members = _sanitize_string_list(record.team_members)
    record.linked_risk_ids = _sanitize_string_list(record.linked_risk_ids)
    record.linked_equipment_ids = _sanitize_string_list(record.linked_equipment_ids)
    record.linked_audit_finding_ids = _sanitize_string_list(record.linked_audit_finding_ids)
    record.updated_at = datetime.now(timezone.utc)
    await db.capas.replace_one({"id": record.id}, record.dict())
    return record


@api_router.get("/complaint-categories", response_model=List[ComplaintCategory])
async def list_complaint_categories() -> List[ComplaintCategory]:
    categories = (
        await db.complaint_categories.find({})
        .sort("name", 1)
        .to_list(200)
    )
    return [ComplaintCategory(**category) for category in categories]


@api_router.post("/complaint-categories", response_model=ComplaintCategory, status_code=status.HTTP_201_CREATED)
async def create_complaint_category(payload: ComplaintCategoryCreate) -> ComplaintCategory:
    name = (payload.name or "").strip()
    if not name:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Category name is required",
        )
    existing = await db.complaint_categories.find_one(
        {"name": {"$regex": f"^{re.escape(name)}$", "$options": "i"}}
    )
    if existing:
        raise HTTPException(
            status_code=status.HTTP_409_CONFLICT,
            detail="Category with same name already exists",
        )
    category = ComplaintCategory(
        name=name,
        description=(payload.description or "").strip() or None,
    )
    await db.complaint_categories.insert_one(category.dict())
    return category


@api_router.patch("/complaint-categories/{category_id}", response_model=ComplaintCategory)
async def update_complaint_category(
    category_id: str,
    payload: ComplaintCategoryUpdate,
) -> ComplaintCategory:
    category_raw = await db.complaint_categories.find_one({"id": category_id})
    if not category_raw:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Complaint category not found")
    update_fields: Dict[str, Any] = {}
    if payload.name is not None:
        name = payload.name.strip()
        if not name:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Category name cannot be empty",
            )
        existing = await db.complaint_categories.find_one(
            {
                "id": {"$ne": category_id},
                "name": {"$regex": f"^{re.escape(name)}$", "$options": "i"},
            }
        )
        if existing:
            raise HTTPException(
                status_code=status.HTTP_409_CONFLICT,
                detail="Another category with same name exists",
            )
        update_fields["name"] = name
    if payload.description is not None:
        update_fields["description"] = payload.description.strip() or None
    if payload.is_active is not None:
        update_fields["is_active"] = payload.is_active
    if not update_fields:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="No fields provided for update")
    update_fields["updated_at"] = datetime.now(timezone.utc)
    await db.complaint_categories.update_one({"id": category_id}, {"$set": update_fields})
    updated_raw = await db.complaint_categories.find_one({"id": category_id})
    return ComplaintCategory(**updated_raw)


@api_router.delete("/complaint-categories/{category_id}", status_code=status.HTTP_204_NO_CONTENT)
async def delete_complaint_category(category_id: str) -> Response:
    now = datetime.now(timezone.utc)
    result = await db.complaint_categories.update_one(
        {"id": category_id},
        {"$set": {"is_active": False, "updated_at": now}},
    )
    if result.matched_count == 0:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Complaint category not found")
    return Response(status_code=status.HTTP_204_NO_CONTENT)


@api_router.get("/complaints", response_model=List[Complaint])
async def get_complaints(current_user: User = Depends(get_current_user)):
    complaints = await db.complaints.find({}).sort("updated_at", -1).to_list(200)
    category_ids: Set[str] = {
        complaint.get("category_id")
        for complaint in complaints
        if complaint.get("category_id")
    }
    category_map: Dict[str, str] = {}
    if category_ids:
        categories = await db.complaint_categories.find(
            {"id": {"$in": list(category_ids)}}
        ).to_list(len(category_ids))
        category_map = {
            category["id"]: category.get("name", "")
            for category in categories
        }
    enriched: List[Complaint] = []
    for complaint in complaints:
        category_id = complaint.get("category_id")
        if category_id and not complaint.get("category_name"):
            complaint["category_name"] = category_map.get(category_id)
        enriched.append(Complaint(**complaint))
    return enriched

@api_router.post("/complaints", response_model=Complaint)
async def create_complaint(complaint_data: ComplaintCreate, current_user: User = Depends(get_current_user)):
    complaint_no = await generate_complaint_no()
    
    complaint_dict = complaint_data.dict(exclude_unset=True)
    complaint_dict["complaint_no"] = complaint_no
    complaint_dict["created_by"] = current_user.id
    complaint_dict.setdefault("solution_team", [])
    complaint_dict.setdefault("file_attachments", [])
    complaint_dict.setdefault("related_task_ids", [])
    complaint_dict.setdefault("related_capa_ids", [])
    category_id = complaint_dict.pop("category_id", None)
    if category_id:
        category = await _get_complaint_category(category_id)
        complaint_dict["category_id"] = category.id
        complaint_dict["category_name"] = category.name
    
    complaint_obj = Complaint(**complaint_dict)
    complaint_obj.status_history.append(
        ComplaintStatus(
            status=complaint_obj.status,
            changed_by=current_user.id,
            changed_at=complaint_obj.created_at,
            comment="Complaint created",
        )
    )
    
    # Save to database
    await db.complaints.insert_one(complaint_obj.dict())
    
    # Create notification
    notification = Notification(
        user_id=current_user.id,
        title="Yeni Şikayet",
        message=f"#{complaint_obj.complaint_no} numaralı şikayet oluşturuldu",
        type="warning"
    )
    await db.notifications.insert_one(notification.dict())
    
    return complaint_obj

@api_router.get("/complaints/{complaint_id}", response_model=Complaint)
async def get_complaint(complaint_id: str, current_user: User = Depends(get_current_user)):
    complaint = await _get_complaint_or_404(complaint_id)
    category_id = complaint.get("category_id")
    if category_id and not complaint.get("category_name"):
        category = await db.complaint_categories.find_one({"id": category_id})
        if category:
            complaint["category_name"] = category.get("name")
    return Complaint(**complaint)

@api_router.put("/complaints/{complaint_id}/status")
async def update_complaint_status(
    complaint_id: str,
    status: str,
    comment: Optional[str] = None,
    current_user: User = Depends(get_current_user)
):
    complaint = await _get_complaint_or_404(complaint_id)

    status_change = ComplaintStatus(
        status=status,
        changed_by=current_user.id,
        comment=comment,
    )

    update_fields = {
        "status": status,
        "updated_at": datetime.now(timezone.utc),
    }

    if status == "resolved":
        update_fields["resolution_date"] = datetime.now(timezone.utc)

    await db.complaints.update_one(
        {"id": complaint_id},
        {
            "$set": update_fields,
            "$push": {"status_history": status_change.dict()},
        },
    )

    return {"message": "Complaint status updated successfully"}


@api_router.patch("/complaints/{complaint_id}/assignment", response_model=Complaint)
async def update_complaint_assignment(
    complaint_id: str,
    payload: ComplaintAssignmentUpdate,
    current_user: User = Depends(get_current_user),
) -> Complaint:
    complaint_raw = await _get_complaint_or_404(complaint_id)
    complaint_model = Complaint(**complaint_raw)

    changed = False
    if payload.assigned_to is not None:
        complaint_model.assigned_to = payload.assigned_to
        changed = True
    if payload.team_leader is not None:
        complaint_model.team_leader = payload.team_leader
        changed = True
    if payload.solution_team is not None:
        complaint_model.solution_team = [member for member in payload.solution_team if member]
        changed = True
    if payload.initial_response is not None:
        complaint_model.initial_response = payload.initial_response
        changed = True

    if not changed:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="No assignment fields provided")

    if complaint_model.status == "open":
        complaint_model.status = "investigating"
        complaint_model.status_history.append(
            ComplaintStatus(
                status="investigating",
                changed_by=current_user.id,
                comment="Complaint assigned to investigation team",
            )
        )

    return await _persist_complaint_record(complaint_model)


@api_router.patch("/complaints/{complaint_id}/investigation", response_model=Complaint)
async def update_complaint_investigation(
    complaint_id: str,
    payload: ComplaintInvestigationUpdate,
    current_user: User = Depends(get_current_user),
) -> Complaint:
    complaint_raw = await _get_complaint_or_404(complaint_id)
    complaint_model = Complaint(**complaint_raw)

    complaint_model.investigation_report = payload.investigation_report
    complaint_model.investigation_reported_at = payload.investigation_reported_at or datetime.now(timezone.utc)
    complaint_model.investigation_reported_by = current_user.id

    complaint_model.file_attachments = _merge_unique_values(
        complaint_model.file_attachments,
        payload.file_attachments,
    )
    complaint_model.related_task_ids = _merge_unique_values(
        complaint_model.related_task_ids,
        payload.related_task_ids,
    )

    if complaint_model.status == "open":
        complaint_model.status = "investigating"
        complaint_model.status_history.append(
            ComplaintStatus(
                status="investigating",
                changed_by=current_user.id,
                comment="Investigation report submitted",
            )
        )

    return await _persist_complaint_record(complaint_model)


@api_router.patch("/complaints/{complaint_id}/finalize", response_model=Complaint)
async def finalize_complaint(
    complaint_id: str,
    payload: ComplaintFinalizationUpdate,
    current_user: User = Depends(get_current_user),
) -> Complaint:
    complaint_raw = await _get_complaint_or_404(complaint_id)
    complaint_model = Complaint(**complaint_raw)

    now = datetime.now(timezone.utc)
    complaint_model.final_report = payload.final_report
    complaint_model.final_reported_at = now
    complaint_model.final_reported_by = current_user.id

    if payload.final_response is not None:
        complaint_model.final_response = payload.final_response

    complaint_model.file_attachments = _merge_unique_values(
        complaint_model.file_attachments,
        payload.file_attachments,
    )

    if payload.mark_resolved and complaint_model.status != "resolved":
        complaint_model.status = "resolved"
        complaint_model.resolution_date = now
        complaint_model.status_history.append(
            ComplaintStatus(
                status="resolved",
                changed_by=current_user.id,
                comment="Final report approved",
            )
        )

    return await _persist_complaint_record(complaint_model)


@api_router.patch("/complaints/{complaint_id}/metadata", response_model=Complaint)
async def update_complaint_metadata(
    complaint_id: str,
    payload: ComplaintMetadataUpdate,
    current_user: User = Depends(get_current_user),
) -> Complaint:
    complaint_raw = await _get_complaint_or_404(complaint_id)
    complaint_model = Complaint(**complaint_raw)

    changed = False
    if payload.complaint_type is not None:
        complaint_model.complaint_type = payload.complaint_type
        changed = True
    if payload.priority is not None:
        complaint_model.priority = payload.priority
        changed = True
    if payload.department is not None:
        complaint_model.department = payload.department
        changed = True
    if payload.category_id is not None:
        if not payload.category_id:
            complaint_model.category_id = None
            complaint_model.category_name = None
        else:
            category = await _get_complaint_category(payload.category_id)
            complaint_model.category_id = category.id
            complaint_model.category_name = category.name
        changed = True

    if not changed:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="No metadata fields provided",
        )

    return await _persist_complaint_record(complaint_model)


@api_router.get("/complaints/report/categories")
async def get_complaint_category_report(current_user: User = Depends(get_current_user)) -> Dict[str, Any]:
    complaints = await db.complaints.find({}).to_list(1000)
    categories = await db.complaint_categories.find({}).to_list(500)
    category_map = {category["id"]: category.get("name", "") for category in categories}

    summary: Dict[str, Dict[str, Any]] = {}
    for complaint in complaints:
        category_id = complaint.get("category_id") or "uncategorized"
        entry = summary.setdefault(
            category_id,
            {
                "category_id": None if category_id == "uncategorized" else category_id,
                "category_name": "Kategorize Edilmedi",
                "count": 0,
                "status_counts": {},
                "priority_counts": {},
            },
        )
        if category_id != "uncategorized":
            entry["category_name"] = (
                complaint.get("category_name")
                or category_map.get(category_id)
                or "Kategori"
            )
        entry["count"] += 1
        status_value = complaint.get("status") or "unknown"
        entry["status_counts"][status_value] = entry["status_counts"].get(status_value, 0) + 1
        priority_value = complaint.get("priority") or "unknown"
        entry["priority_counts"][priority_value] = entry["priority_counts"].get(priority_value, 0) + 1

    categories_summary = sorted(
        summary.values(),
        key=lambda item: item["count"],
        reverse=True,
    )

    return {
        "total": len(complaints),
        "categories": categories_summary,
    }


@api_router.post("/complaints/{complaint_id}/attachments", response_model=Complaint)
async def add_complaint_attachments(
    complaint_id: str,
    payload: ComplaintAttachmentAdd,
    current_user: User = Depends(get_current_user),
) -> Complaint:
    if not payload.file_ids:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="file_ids listesi boş olamaz.",
        )
    complaint_raw = await _get_complaint_or_404(complaint_id)
    complaint_model = Complaint(**complaint_raw)

    complaint_model.file_attachments = _merge_unique_values(
        complaint_model.file_attachments,
        payload.file_ids,
    )

    return await _persist_complaint_record(complaint_model)


@api_router.post(
    "/complaints/{complaint_id}/dof",
    response_model=DofTask,
    status_code=status.HTTP_201_CREATED,
)
async def create_dof_for_complaint(
    complaint_id: str,
    payload: ComplaintDofCreate,
    current_user: User = Depends(get_current_user),
) -> DofTask:
    complaint_raw = await _get_complaint_or_404(complaint_id)
    complaint_model = Complaint(**complaint_raw)

    now = datetime.now(timezone.utc)
    title = payload.title or f"Şikayet {complaint_model.complaint_no}"
    description = payload.description or complaint_model.description
    department = payload.department or complaint_model.department
    if not department:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Department is required to create a DÖF kaydı.",
        )

    responsible = (
        (payload.responsible_person or "").strip()
        or (complaint_model.assigned_to or "").strip()
        or (complaint_model.team_leader or "").strip()
        or (current_user.full_name or "").strip()
        or (current_user.username or "").strip()
    )
    if not responsible:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Sorumlu kişi belirtmelisiniz.",
        )

    due_date = payload.due_date or (now + timedelta(days=7))
    if due_date.tzinfo is None:
        due_date = due_date.replace(tzinfo=timezone.utc)

    team_members = _sanitize_string_list(payload.team_members or [])
    initial_report_date = payload.initial_improvement_report_date
    if initial_report_date and initial_report_date.tzinfo is None:
        initial_report_date = initial_report_date.replace(tzinfo=timezone.utc)

    dof_no = await generate_dof_no()
    task = DofTask(
        dof_no=dof_no,
        title=title,
        description=description,
        department=department,
        responsible_person=responsible,
        due_date=due_date,
        team_members=team_members,
        initial_improvement_report_date=initial_report_date,
        created_by=current_user.id,
        status_history=[
            DofStatusHistory(status="open", changed_by=current_user.id, changed_at=now)
        ],
    )

    await db.dof_tasks.insert_one(task.dict())

    complaint_model.related_task_ids = _merge_unique_values(
        complaint_model.related_task_ids,
        [task.id],
    )
    await _persist_complaint_record(complaint_model)

    notification = Notification(
        user_id=current_user.id,
        title="Yeni DÖF",
        message=f"{task.dof_no} numaralı DÖF, {complaint_model.complaint_no} şikayetiyle ilişkilendirildi.",
        type="info",
    )
    await db.notifications.insert_one(notification.dict())

    return task


@api_router.post(
    "/complaints/{complaint_id}/capas",
    response_model=Capa,
    status_code=status.HTTP_201_CREATED,
)
async def create_capa_for_complaint(
    complaint_id: str,
    payload: ComplaintCapaCreate,
    current_user: User = Depends(get_current_user),
) -> Capa:
    complaint_raw = await _get_complaint_or_404(complaint_id)
    complaint_model = Complaint(**complaint_raw)

    now = datetime.now(timezone.utc)
    title = payload.title or f"Şikayet {complaint_model.complaint_no} CAPA"
    department = payload.department or complaint_model.department
    if not department:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Department is required to create a CAPA kaydı.",
        )

    team_leader = (
        (payload.team_leader or "").strip()
        or (complaint_model.team_leader or "").strip()
        or (complaint_model.assigned_to or "").strip()
        or (current_user.full_name or "").strip()
        or (current_user.username or "").strip()
    )
    if not team_leader:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="CAPA team leader belirtilmelidir.",
        )

    source = payload.source or "customer_complaint"
    nonconformity_description = (
        payload.nonconformity_description or complaint_model.description
    )

    target_date = payload.target_date
    if target_date and target_date.tzinfo is None:
        target_date = target_date.replace(tzinfo=timezone.utc)

    initial_report_date = payload.initial_improvement_report_date
    if initial_report_date and initial_report_date.tzinfo is None:
        initial_report_date = initial_report_date.replace(tzinfo=timezone.utc)

    attachments = _sanitize_string_list(payload.file_attachments or [])
    team_members = _sanitize_string_list(payload.team_members or [])
    linked_risks = _sanitize_string_list(payload.linked_risk_ids or [])
    linked_equipment = _sanitize_string_list(payload.linked_equipment_ids or [])
    linked_audit_findings = _sanitize_string_list(payload.linked_audit_finding_ids or [])

    capa_no = await generate_capa_no()
    capa_obj = Capa(
        capa_no=capa_no,
        title=title,
        source=source,
        department=department,
        initiated_by=current_user.id,
        team_leader=team_leader,
        target_date=target_date,
        nonconformity_description=nonconformity_description,
        file_attachments=attachments,
        team_members=team_members,
        initial_improvement_report_date=initial_report_date,
        linked_risk_ids=linked_risks,
        linked_equipment_ids=linked_equipment,
        linked_audit_finding_ids=linked_audit_findings,
    )

    await db.capas.insert_one(capa_obj.dict())

    complaint_model.related_capa_ids = _merge_unique_values(
        complaint_model.related_capa_ids,
        [capa_obj.id],
    )
    await _persist_complaint_record(complaint_model)

    notification = Notification(
        user_id=current_user.id,
        title="Yeni CAPA",
        message=f"{capa_obj.capa_no} numaralı CAPA, {complaint_model.complaint_no} şikayetiyle ilişkilendirildi.",
        type="info",
    )
    await db.notifications.insert_one(notification.dict())

    return capa_obj

# CAPA Routes
@api_router.get("/capas", response_model=List[Capa])
async def get_capas(current_user: User = Depends(get_current_user)):
    capas = await db.capas.find({}).sort("updated_at", -1).to_list(100)
    return [Capa(**capa) for capa in capas]

@api_router.post("/capas", response_model=Capa)
async def create_capa(capa_data: CapaCreate, current_user: User = Depends(get_current_user)):
    capa_no = await generate_capa_no()
    
    capa_dict = capa_data.dict()
    capa_dict["capa_no"] = capa_no
    capa_dict["initiated_by"] = current_user.id
    if capa_dict.get("target_date") and capa_dict["target_date"].tzinfo is None:
        capa_dict["target_date"] = capa_dict["target_date"].replace(tzinfo=timezone.utc)
    initial_report_date = capa_dict.pop("initial_improvement_report_date", None)
    if initial_report_date and initial_report_date.tzinfo is None:
        initial_report_date = initial_report_date.replace(tzinfo=timezone.utc)
    team_members = _sanitize_string_list(capa_dict.pop("team_members", []))
    linked_risks = _sanitize_string_list(capa_dict.pop("linked_risk_ids", []))
    linked_equipment = _sanitize_string_list(capa_dict.pop("linked_equipment_ids", []))
    linked_audit_findings = _sanitize_string_list(capa_dict.pop("linked_audit_finding_ids", []))
    
    capa_obj = Capa(**capa_dict)
    capa_obj.team_members = team_members
    capa_obj.initial_improvement_report_date = initial_report_date
    capa_obj.linked_risk_ids = linked_risks
    capa_obj.linked_equipment_ids = linked_equipment
    capa_obj.linked_audit_finding_ids = linked_audit_findings
    capa_obj.closure_requested_at = None
    capa_obj.closure_requested_by = None
    capa_obj.closure_request_note = None
    capa_obj.closure_approved_at = None
    capa_obj.closure_approved_by = None
    capa_obj.closure_decision_note = None
    
    # Save to database
    await db.capas.insert_one(capa_obj.dict())
    
    # Create notification
    notification = Notification(
        user_id=current_user.id,
        title="Yeni DÖF/CAPA",
        message=f"#{capa_obj.capa_no} numaralı DÖF/CAPA oluşturuldu",
        type="info"
    )
    await db.notifications.insert_one(notification.dict())
    
    return capa_obj

@api_router.get("/capas/{capa_id}", response_model=Capa)
async def get_capa(capa_id: str, current_user: User = Depends(get_current_user)):
    capa = await _get_capa_or_404(capa_id)
    return Capa(**capa)


@api_router.patch("/capas/{capa_id}", response_model=Capa)
async def update_capa_record(
    capa_id: str,
    payload: CapaUpdate,
    current_user: User = Depends(get_current_user),
) -> Capa:
    capa_raw = await _get_capa_or_404(capa_id)
    capa_model = Capa(**capa_raw)
    update_fields = payload.dict(exclude_unset=True)
    if not update_fields:
        return capa_model

    if "target_date" in update_fields and update_fields["target_date"] is not None:
        target = update_fields["target_date"]
        if target.tzinfo is None:
            target = target.replace(tzinfo=timezone.utc)
        capa_model.target_date = target
    if "initial_improvement_report_date" in update_fields:
        initial_date = update_fields.pop("initial_improvement_report_date")
        if initial_date and initial_date.tzinfo is None:
            initial_date = initial_date.replace(tzinfo=timezone.utc)
        capa_model.initial_improvement_report_date = initial_date
    if "team_members" in update_fields:
        capa_model.team_members = _sanitize_string_list(update_fields.pop("team_members"))
    if "linked_risk_ids" in update_fields:
        capa_model.linked_risk_ids = _sanitize_string_list(update_fields.pop("linked_risk_ids"))
    if "linked_equipment_ids" in update_fields:
        capa_model.linked_equipment_ids = _sanitize_string_list(
            update_fields.pop("linked_equipment_ids")
        )
    if "linked_audit_finding_ids" in update_fields:
        capa_model.linked_audit_finding_ids = _sanitize_string_list(
            update_fields.pop("linked_audit_finding_ids")
        )
    if "status" in update_fields:
        new_status = update_fields["status"]
        if new_status not in ALLOWED_CAPA_STATUSES:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Gecersiz CAPA durumu. Kullanilabilir degerler: {', '.join(sorted(ALLOWED_CAPA_STATUSES))}",
            )
        capa_model.status = new_status
        update_fields.pop("status")
        if new_status != "pending_closure":
            capa_model.closure_requested_at = None
            capa_model.closure_requested_by = None
            capa_model.closure_request_note = None
    for key in ["title", "source", "department", "team_leader",
                "nonconformity_description", "root_cause_analysis",
                "immediate_action", "effectiveness_review"]:
        if key in update_fields:
            setattr(capa_model, key, update_fields[key])
    if "file_attachments" in update_fields and update_fields["file_attachments"] is not None:
        capa_model.file_attachments = [item for item in update_fields["file_attachments"] if item]

    return await _persist_capa_record(capa_model)


@api_router.post("/capas/{capa_id}/closure/request", response_model=Capa)
async def request_capa_closure(
    capa_id: str,
    payload: CapaClosureRequest,
    current_user: User = Depends(get_current_user),
) -> Capa:
    capa_raw = await _get_capa_or_404(capa_id)
    capa_model = Capa(**capa_raw)
    if capa_model.status in {"closed", "cancelled"}:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Kapanmis veya iptal edilmis kayit icin onay istenemez.")
    now = datetime.now(timezone.utc)
    capa_model.status = "pending_closure"
    capa_model.closure_requested_at = now
    capa_model.closure_requested_by = current_user.id
    capa_model.closure_request_note = payload.note
    capa_model.closure_approved_at = None
    capa_model.closure_approved_by = None
    capa_model.closure_decision_note = None
    return await _persist_capa_record(capa_model)


@api_router.post("/capas/{capa_id}/closure/decision", response_model=Capa)
async def decide_capa_closure(
    capa_id: str,
    payload: CapaClosureDecision,
    current_user: User = Depends(get_current_user),
) -> Capa:
    capa_raw = await _get_capa_or_404(capa_id)
    capa_model = Capa(**capa_raw)
    if capa_model.status != "pending_closure":
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Onay bekleyen bir kapanis bulunmuyor.")
    now = datetime.now(timezone.utc)
    capa_model.closure_decision_note = payload.note
    if payload.approve:
        capa_model.status = "closed"
        capa_model.closure_approved_at = now
        capa_model.closure_approved_by = current_user.id
        capa_model.closed_at = now
    else:
        capa_model.status = "implementing"
        capa_model.closure_requested_at = None
        capa_model.closure_requested_by = None
        capa_model.closure_request_note = None
        capa_model.closure_approved_at = None
        capa_model.closure_approved_by = None
    return await _persist_capa_record(capa_model)


def _find_capa_action(capa: Capa, action_id: str) -> tuple[Optional[CapaAction], Optional[str]]:
    for action in capa.corrective_actions:
        if action.id == action_id:
            return action, "corrective"
    for action in capa.preventive_actions:
        if action.id == action_id:
            return action, "preventive"
    return None, None


@api_router.post("/capas/{capa_id}/actions", response_model=Capa, status_code=status.HTTP_201_CREATED)
async def add_capa_action(
    capa_id: str,
    payload: CapaActionCreate,
    current_user: User = Depends(get_current_user),
) -> Capa:
    capa_raw = await _get_capa_or_404(capa_id)
    capa_model = Capa(**capa_raw)

    due_date = payload.due_date
    if due_date.tzinfo is None:
        due_date = due_date.replace(tzinfo=timezone.utc)

    action = CapaAction(
        action_description=payload.action_description,
        responsible_person=payload.responsible_person,
        due_date=due_date,
        status=payload.status or "open",
        evidence=payload.evidence,
    )
    if action.status == "completed":
        action.completion_date = datetime.now(timezone.utc)

    if payload.action_type == "corrective":
        capa_model.corrective_actions.append(action)
    else:
        capa_model.preventive_actions.append(action)

    return await _persist_capa_record(capa_model)


@api_router.patch("/capas/{capa_id}/actions/{action_id}", response_model=Capa)
async def update_capa_action(
    capa_id: str,
    action_id: str,
    payload: CapaActionUpdate,
    current_user: User = Depends(get_current_user),
) -> Capa:
    capa_raw = await _get_capa_or_404(capa_id)
    capa_model = Capa(**capa_raw)
    action, action_type = _find_capa_action(capa_model, action_id)
    if action is None:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="CAPA action not found")

    update_fields = payload.dict(exclude_unset=True)
    if not update_fields:
        return capa_model

    if "action_description" in update_fields:
        action.action_description = update_fields["action_description"]
    if "responsible_person" in update_fields:
        action.responsible_person = update_fields["responsible_person"]
    if "due_date" in update_fields and update_fields["due_date"] is not None:
        due_date = update_fields["due_date"]
        if due_date.tzinfo is None:
            due_date = due_date.replace(tzinfo=timezone.utc)
        action.due_date = due_date
    if "status" in update_fields and update_fields["status"]:
        action.status = update_fields["status"]
        if action.status == "completed" and not update_fields.get("completion_date"):
            action.completion_date = datetime.now(timezone.utc)
    if "completion_date" in update_fields:
        completion_date = update_fields["completion_date"]
        if completion_date:
            if completion_date.tzinfo is None:
                completion_date = completion_date.replace(tzinfo=timezone.utc)
            action.completion_date = completion_date
        else:
            action.completion_date = None
    if "evidence" in update_fields:
        action.evidence = update_fields["evidence"]

    # Ensure lists are updated with new action data
    if action_type == "corrective":
        capa_model.corrective_actions = [
            action if existing.id == action.id else existing
            for existing in capa_model.corrective_actions
        ]
    else:
        capa_model.preventive_actions = [
            action if existing.id == action.id else existing
            for existing in capa_model.preventive_actions
        ]

    return await _persist_capa_record(capa_model)

# DÖF (Corrective & Preventive Task) Routes

# DÖF (Corrective & Preventive Task) Routes
# Audit Question Routes
@api_router.get("/audit-questions", response_model=List[AuditQuestion])
async def list_audit_questions(
    search: Optional[str] = None,
    category: Optional[str] = None,
    include_inactive: bool = False,
    current_user: User = Depends(get_current_user),
) -> List[AuditQuestion]:
    query: Dict[str, Any] = {}
    if not include_inactive:
        query["is_active"] = True
    if category:
        query["category"] = category.strip()
    if search:
        pattern = {"$regex": re.escape(search), "$options": "i"}
        query["$or"] = [
            {"question": pattern},
            {"requirement_reference": pattern},
            {"tags": pattern},
        ]
    cursor = (
        db.audit_questions.find(query)
        .sort("updated_at", -1)
        .limit(200)
    )
    return [AuditQuestion(**doc) async for doc in cursor]


@api_router.post("/audit-questions", response_model=AuditQuestion, status_code=status.HTTP_201_CREATED)
async def create_audit_question(
    payload: AuditQuestionCreate,
    current_user: User = Depends(get_current_user),
) -> AuditQuestion:
    question_text = _normalize_required_str(payload.question, "Question text")
    category = _normalize_optional_str(payload.category)
    requirement = _normalize_optional_str(payload.requirement_reference)
    tags = _sanitize_string_list(payload.tags)
    is_active = True if payload.is_active is None else bool(payload.is_active)

    question = AuditQuestion(
        question=question_text,
        category=category,
        requirement_reference=requirement,
        tags=tags,
        is_active=is_active,
        created_by=current_user.id,
    )
    await db.audit_questions.insert_one(question.dict())
    return question


@api_router.patch("/audit-questions/{question_id}", response_model=AuditQuestion)
async def update_audit_question(
    question_id: str,
    payload: AuditQuestionUpdate,
    current_user: User = Depends(get_current_user),
) -> AuditQuestion:
    existing = await db.audit_questions.find_one({"id": question_id})
    if not existing:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Question not found.")

    update_fields = payload.dict(exclude_unset=True)
    updates: Dict[str, Any] = {}
    if "question" in update_fields:
        updates["question"] = _normalize_required_str(update_fields["question"], "Question text")
    if "category" in update_fields:
        updates["category"] = _normalize_optional_str(update_fields["category"])
    if "requirement_reference" in update_fields:
        updates["requirement_reference"] = _normalize_optional_str(update_fields["requirement_reference"])
    if "tags" in update_fields:
        updates["tags"] = _sanitize_string_list(update_fields["tags"])
    if "is_active" in update_fields:
        updates["is_active"] = bool(update_fields["is_active"])
    if not updates:
        return AuditQuestion(**existing)

    updates["updated_at"] = datetime.now(timezone.utc)
    await db.audit_questions.update_one(
        {"id": question_id},
        {"$set": updates},
    )
    updated = await db.audit_questions.find_one({"id": question_id})
    return AuditQuestion(**updated)


@api_router.delete("/audit-questions/{question_id}", status_code=status.HTTP_204_NO_CONTENT)
async def delete_audit_question(
    question_id: str,
    current_user: User = Depends(get_current_user),
) -> Response:
    result = await db.audit_questions.delete_one({"id": question_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Question not found.")
    return Response(status_code=status.HTTP_204_NO_CONTENT)


# Audit Routes
@api_router.get("/audits", response_model=List[Audit])
async def list_audits(
    status_filter: Optional[str] = Query(default=None, alias="status"),
    department: Optional[str] = None,
    audit_type: Optional[str] = None,
    search: Optional[str] = None,
    start_date: Optional[str] = None,
    end_date: Optional[str] = None,
    current_user: User = Depends(get_current_user),
) -> List[Audit]:
    query: Dict[str, Any] = {}
    if status_filter:
        query["status"] = _normalize_audit_status(status_filter)
    if department:
        query["department"] = department.strip()
    if audit_type:
        query["audit_type"] = audit_type.strip()
    if search:
        pattern = {"$regex": re.escape(search), "$options": "i"}
        query["$or"] = [
            {"audit_code": pattern},
            {"scope": pattern},
            {"lead_auditor": pattern},
        ]
    date_filter: Dict[str, Any] = {}
    if start_date:
        try:
            date_filter["$gte"] = _parse_iso_datetime(start_date)
        except ValueError:
            raise HTTPException(status_code=400, detail="start_date must be ISO 8601 format.")
    if end_date:
        try:
            date_filter["$lte"] = _parse_iso_datetime(end_date)
        except ValueError:
            raise HTTPException(status_code=400, detail="end_date must be ISO 8601 format.")
    if date_filter:
        query["start_date"] = date_filter

    cursor = db.audits.find(query).sort("updated_at", -1)
    return [Audit(**doc) async for doc in cursor]


@api_router.post("/audits", response_model=Audit, status_code=status.HTTP_201_CREATED)
async def create_audit(
    payload: AuditCreate,
    current_user: User = Depends(get_current_user),
) -> Audit:
    audit_code = payload.audit_code or await generate_audit_code()
    audit_type = _normalize_required_str(payload.audit_type, "Audit type")
    scope = _normalize_required_str(payload.scope, "Audit scope")
    department = _normalize_required_str(payload.department, "Department")
    lead_auditor = _normalize_required_str(payload.lead_auditor, "Lead auditor")

    start_date = _ensure_timezone(payload.start_date)
    end_date = _ensure_timezone(payload.end_date)
    if end_date < start_date:
        raise HTTPException(status_code=400, detail="End date cannot be earlier than start date.")

    team = _sanitize_audit_team_members(payload.audit_team or [])
    checklist = await _prepare_checklist_items(payload.checklist or [])
    status_history = [
        AuditStatusHistory(
            status="planned",
            changed_by=current_user.id,
            note="Audit created",
        )
    ]

    audit = Audit(
        audit_code=audit_code,
        audit_type=audit_type,
        scope=scope,
        department=department,
        start_date=start_date,
        end_date=end_date,
        status="planned",
        lead_auditor=lead_auditor,
        audit_team=team,
        auditee_representative=_normalize_optional_str(payload.auditee_representative),
        objectives=_normalize_optional_str(payload.objectives),
        checklist=checklist,
        findings=[],
        status_history=status_history,
        created_by=current_user.id,
    )
    await db.audits.insert_one(audit.dict())
    return audit


@api_router.get("/audits/{audit_id}", response_model=Audit)
async def get_audit(
    audit_id: str,
    current_user: User = Depends(get_current_user),
) -> Audit:
    audit_doc = await _get_audit_or_404(audit_id)
    return Audit(**audit_doc)


@api_router.patch("/audits/{audit_id}", response_model=Audit)
async def update_audit_record(
    audit_id: str,
    payload: AuditUpdate,
    current_user: User = Depends(get_current_user),
) -> Audit:
    audit_doc = await _get_audit_or_404(audit_id)
    audit_model = Audit(**audit_doc)
    data = payload.dict(exclude_unset=True)
    if not data:
        return audit_model

    if "audit_type" in data:
        audit_model.audit_type = _normalize_required_str(data["audit_type"], "Audit type")
    if "scope" in data:
        audit_model.scope = _normalize_required_str(data["scope"], "Audit scope")
    if "department" in data:
        audit_model.department = _normalize_required_str(data["department"], "Department")
    if "lead_auditor" in data:
        audit_model.lead_auditor = _normalize_required_str(data["lead_auditor"], "Lead auditor")
    if "auditee_representative" in data:
        audit_model.auditee_representative = _normalize_optional_str(data["auditee_representative"])
    if "objectives" in data:
        audit_model.objectives = _normalize_optional_str(data["objectives"])
    if "start_date" in data:
        audit_model.start_date = _ensure_timezone(data["start_date"])
    if "end_date" in data:
        audit_model.end_date = _ensure_timezone(data["end_date"])
    if audit_model.end_date < audit_model.start_date:
        raise HTTPException(status_code=400, detail="End date cannot be earlier than start date.")
    if "audit_team" in data and data["audit_team"] is not None:
        audit_model.audit_team = _sanitize_audit_team_members(data["audit_team"])
    if "checklist" in data and data["checklist"] is not None:
        audit_model.checklist = await _prepare_checklist_items(data["checklist"])
    if "status" in data and data["status"] is not None:
        new_status = _normalize_audit_status(data["status"])
        if new_status != audit_model.status:
            audit_model.status = new_status
            audit_model.status_history.append(
                AuditStatusHistory(
                    status=new_status,
                    changed_by=current_user.id,
                    note="Audit status updated",
                )
            )

    audit_model = await _persist_audit_record(audit_model)
    return audit_model


@api_router.post("/audits/{audit_id}/status", response_model=Audit)
async def update_audit_status(
    audit_id: str,
    payload: AuditStatusUpdate,
    current_user: User = Depends(get_current_user),
) -> Audit:
    audit_doc = await _get_audit_or_404(audit_id)
    audit_model = Audit(**audit_doc)
    new_status = _normalize_audit_status(payload.status)
    note = _normalize_optional_str(payload.note)

    audit_model.status = new_status
    audit_model.status_history.append(
        AuditStatusHistory(
            status=new_status,
            changed_by=current_user.id,
            note=note,
        )
    )
    audit_model = await _persist_audit_record(audit_model)
    return audit_model


@api_router.post("/audits/{audit_id}/findings", response_model=Audit)
async def add_audit_finding(
    audit_id: str,
    payload: AuditFindingCreate,
    current_user: User = Depends(get_current_user),
) -> Audit:
    audit_doc = await _get_audit_or_404(audit_id)
    audit_model = Audit(**audit_doc)

    finding = AuditFinding(
        finding_type=_normalize_finding_type(payload.finding_type),
        description=_normalize_required_str(payload.description, "Finding description"),
        requirement_reference=_normalize_optional_str(payload.requirement_reference),
        related_capa_id=_normalize_optional_str(payload.related_capa_id),
        status="open",
        corrective_action=_normalize_optional_str(payload.corrective_action),
        created_at=datetime.now(timezone.utc),
        updated_at=datetime.now(timezone.utc),
    )
    audit_model.findings.append(finding)
    audit_model = await _persist_audit_record(audit_model)
    await _sync_finding_capa_link(finding.id, finding.related_capa_id, None)
    return audit_model


@api_router.patch("/audits/{audit_id}/findings/{finding_id}", response_model=Audit)
async def update_audit_finding(
    audit_id: str,
    finding_id: str,
    payload: AuditFindingUpdate,
    current_user: User = Depends(get_current_user),
) -> Audit:
    audit_doc = await _get_audit_or_404(audit_id)
    audit_model = Audit(**audit_doc)
    finding = next((item for item in audit_model.findings if item.id == finding_id), None)
    if not finding:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Audit finding not found.")

    data = payload.dict(exclude_unset=True)
    previous_capa = finding.related_capa_id
    if "finding_type" in data:
        finding.finding_type = _normalize_finding_type(data["finding_type"])
    if "description" in data:
        finding.description = _normalize_required_str(data["description"], "Finding description")
    if "requirement_reference" in data:
        finding.requirement_reference = _normalize_optional_str(data["requirement_reference"])
    if "related_capa_id" in data:
        finding.related_capa_id = _normalize_optional_str(data["related_capa_id"])
    if "status" in data:
        finding.status = _normalize_finding_status(data["status"])
    if "corrective_action" in data:
        finding.corrective_action = _normalize_optional_str(data["corrective_action"])
    finding.updated_at = datetime.now(timezone.utc)

    audit_model = await _persist_audit_record(audit_model)
    await _sync_finding_capa_link(finding.id, finding.related_capa_id, previous_capa)
    return audit_model

# Calibration Routes
@api_router.get("/calibration/devices", response_model=List[CalibrationDevice])
async def list_calibration_devices(
    status_filter: Optional[str] = Query(default=None, alias="status"),
    category: Optional[str] = None,
    search: Optional[str] = None,
    current_user: User = Depends(get_current_user),
) -> List[CalibrationDevice]:
    query: Dict[str, Any] = {}
    if status_filter:
        query["status"] = _normalize_device_status(status_filter)
    if category:
        query["category"] = category.strip()
    if search:
        pattern = {"$regex": re.escape(search), "$options": "i"}
        query["$or"] = [
            {"device_code": pattern},
            {"name": pattern},
            {"serial_number": pattern},
            {"location": pattern},
        ]
    cursor = (
        db.calibration_devices.find(query)
        .sort([("next_due_date", 1), ("name", 1)])
        .limit(500)
    )
    return [CalibrationDevice(**doc) async for doc in cursor]


@api_router.post("/calibration/devices", response_model=CalibrationDevice, status_code=status.HTTP_201_CREATED)
async def create_calibration_device(
    payload: CalibrationDeviceCreate,
    current_user: User = Depends(get_current_user),
) -> CalibrationDevice:
    status_value = _normalize_device_status(payload.status or "active")
    last_calibrated = _ensure_datetime(payload.last_calibrated_at)
    next_due = _ensure_datetime(payload.next_due_date)
    interval = payload.calibration_interval_days or 1
    notice_days = payload.notice_days if payload.notice_days is not None else 14
    if not next_due:
        next_due = _compute_next_due_date(
            last_calibrated,
            interval,
            datetime.now(timezone.utc),
        )
    device = CalibrationDevice(
        device_code=await _generate_device_code(),
        name=_normalize_required_str(payload.name, "Device name"),
        category=_normalize_required_str(payload.category, "Device category"),
        location=_normalize_optional_str(payload.location),
        manufacturer=_normalize_optional_str(payload.manufacturer),
        model=_normalize_optional_str(payload.model),
        serial_number=_normalize_optional_str(payload.serial_number),
        department=_normalize_optional_str(payload.department),
        responsible_person=_normalize_optional_str(payload.responsible_person),
        status=status_value,
        calibration_interval_days=max(int(interval), 1),
        notice_days=max(int(notice_days), 0),
        last_calibrated_at=last_calibrated,
        next_due_date=next_due,
        linked_capa_ids=_sanitize_string_list(payload.linked_capa_ids),
        linked_dof_task_ids=_sanitize_string_list(payload.linked_dof_task_ids),
        file_attachments=_sanitize_string_list(payload.file_attachments),
        notes=_normalize_optional_str(payload.notes),
        created_by=current_user.id,
    )
    _append_device_status_history(device, status_value, current_user.id, "Device created")
    device = await _persist_device_record(device)
    return device


@api_router.get("/calibration/devices/{device_id}", response_model=CalibrationDevice)
async def get_calibration_device(
    device_id: str,
    current_user: User = Depends(get_current_user),
) -> CalibrationDevice:
    device_doc = await _get_device_or_404(device_id)
    return CalibrationDevice(**device_doc)


@api_router.patch("/calibration/devices/{device_id}", response_model=CalibrationDevice)
async def update_calibration_device(
    device_id: str,
    payload: CalibrationDeviceUpdate,
    current_user: User = Depends(get_current_user),
) -> CalibrationDevice:
    device_doc = await _get_device_or_404(device_id)
    device_model = CalibrationDevice(**device_doc)
    data = payload.dict(exclude_unset=True)
    status_note = _normalize_optional_str(data.pop("status_note", None))
    raw_status = data.pop("status", None)
    status_changed = False
    if "name" in data:
        device_model.name = _normalize_required_str(data["name"], "Device name")
    if "category" in data:
        device_model.category = _normalize_required_str(data["category"], "Device category")
    if "location" in data:
        device_model.location = _normalize_optional_str(data["location"])
    if "manufacturer" in data:
        device_model.manufacturer = _normalize_optional_str(data["manufacturer"])
    if "model" in data:
        device_model.model = _normalize_optional_str(data["model"])
    if "serial_number" in data:
        device_model.serial_number = _normalize_optional_str(data["serial_number"])
    if "department" in data:
        device_model.department = _normalize_optional_str(data["department"])
    if "responsible_person" in data:
        device_model.responsible_person = _normalize_optional_str(data["responsible_person"])
    if "calibration_interval_days" in data and data["calibration_interval_days"] is not None:
        device_model.calibration_interval_days = max(int(data["calibration_interval_days"]), 1)
    if "notice_days" in data and data["notice_days"] is not None:
        device_model.notice_days = max(int(data["notice_days"]), 0)
    if "notes" in data:
        device_model.notes = _normalize_optional_str(data["notes"])
    if "file_attachments" in data and data["file_attachments"] is not None:
        device_model.file_attachments = _sanitize_string_list(data["file_attachments"])
    if "linked_capa_ids" in data and data["linked_capa_ids"] is not None:
        device_model.linked_capa_ids = _sanitize_string_list(data["linked_capa_ids"])
    if "linked_dof_task_ids" in data and data["linked_dof_task_ids"] is not None:
        device_model.linked_dof_task_ids = _sanitize_string_list(data["linked_dof_task_ids"])
    if "last_calibrated_at" in data:
        device_model.last_calibrated_at = _ensure_datetime(data["last_calibrated_at"])
    if "next_due_date" in data:
        device_model.next_due_date = _ensure_datetime(data["next_due_date"])
    if raw_status is not None:
        new_status = _normalize_device_status(raw_status)
        _validate_device_transition(device_model.status, new_status)
        status_changed = new_status != device_model.status
        device_model.status = new_status

    if device_model.next_due_date is None:
        device_model.next_due_date = _compute_next_due_date(
            device_model.last_calibrated_at,
            device_model.calibration_interval_days,
            datetime.now(timezone.utc),
        )

    if status_changed or status_note:
        _append_device_status_history(
            device_model,
            device_model.status,
            current_user.id,
            status_note,
        )

    device_model = await _persist_device_record(device_model)
    return device_model


@api_router.get("/calibration/work-orders", response_model=List[CalibrationWorkOrder])
async def list_calibration_work_orders(
    status_filter: Optional[str] = Query(default=None, alias="status"),
    device_id: Optional[str] = None,
    due_before: Optional[str] = None,
    due_after: Optional[str] = None,
    current_user: User = Depends(get_current_user),
) -> List[CalibrationWorkOrder]:
    query: Dict[str, Any] = {}
    if status_filter:
        query["status"] = _normalize_work_order_status(status_filter)
    if device_id:
        query["device_id"] = device_id
    if due_before or due_after:
        date_filter: Dict[str, Any] = {}
        if due_before:
            try:
                date_filter["$lte"] = _parse_iso_datetime(due_before)
            except ValueError:
                raise HTTPException(status_code=400, detail="due_before must be ISO 8601 formatted.")
        if due_after:
            try:
                date_filter["$gte"] = _parse_iso_datetime(due_after)
            except ValueError:
                raise HTTPException(status_code=400, detail="due_after must be ISO 8601 formatted.")
        query["due_date"] = date_filter
    cursor = (
        db.calibration_work_orders.find(query)
        .sort([("due_date", 1), ("created_at", -1)])
        .limit(500)
    )
    return [CalibrationWorkOrder(**doc) async for doc in cursor]


@api_router.post("/calibration/devices/{device_id}/work-orders", response_model=CalibrationWorkOrder, status_code=status.HTTP_201_CREATED)
async def create_calibration_work_order(
    device_id: str,
    payload: CalibrationWorkOrderCreate,
    current_user: User = Depends(get_current_user),
) -> CalibrationWorkOrder:
    device_doc = await _get_device_or_404(device_id)
    device_model = CalibrationDevice(**device_doc)
    planned = _ensure_datetime(payload.planned_date)
    due = _ensure_datetime(payload.due_date)
    if due is None:
        due = planned or datetime.now(timezone.utc)
    work_order = CalibrationWorkOrder(
        work_order_no=await _generate_work_order_no(),
        device_id=device_id,
        planned_date=planned or due,
        due_date=due,
        status="planned",
        assigned_to=_normalize_optional_str(payload.assigned_to),
        notes=_normalize_optional_str(payload.notes),
        created_by=current_user.id,
    )
    _append_work_order_status_history(work_order, work_order.status, current_user.id, "Work order created")
    work_order = await _persist_work_order_record(work_order)
    return work_order


@api_router.post("/calibration/work-orders/scheduler/run")
async def run_calibration_scheduler(
    horizon_days: int = 30,
    current_user: User = Depends(get_current_user),
) -> Dict[str, Any]:
    now = datetime.now(timezone.utc)
    horizon = now + timedelta(days=max(horizon_days, 1))
    created_orders = []
    cursor = db.calibration_devices.find({"status": {"$ne": "inactive"}})
    async for doc in cursor:
        device = CalibrationDevice(**doc)
        next_due = device.next_due_date or _compute_next_due_date(
            device.last_calibrated_at,
            device.calibration_interval_days,
            now,
        )
        if not next_due:
            continue
        if next_due > horizon:
            continue
        existing = await db.calibration_work_orders.count_documents(
            {
                "device_id": device.id,
                "status": {"$in": ["planned", "in_progress"]},
                "due_date": {"$gte": now - timedelta(days=30)},
            }
        )
        if existing:
            continue
        work_order = CalibrationWorkOrder(
            work_order_no=await _generate_work_order_no(),
            device_id=device.id,
            planned_date=now,
            due_date=next_due,
            status="planned",
            assigned_to=device.responsible_person,
            notes="Automatically generated calibration work order.",
            created_by=current_user.id,
        )
        _append_work_order_status_history(
            work_order,
            work_order.status,
            current_user.id,
            "Scheduler generated work order",
        )
        await _persist_work_order_record(work_order)
        created_orders.append(work_order.work_order_no)
    return {"created": len(created_orders), "work_orders": created_orders}


@api_router.get("/calibration/work-orders/{work_order_id}", response_model=CalibrationWorkOrder)
async def get_calibration_work_order(
    work_order_id: str,
    current_user: User = Depends(get_current_user),
) -> CalibrationWorkOrder:
    record = await _get_work_order_or_404(work_order_id)
    return CalibrationWorkOrder(**record)


@api_router.patch("/calibration/work-orders/{work_order_id}", response_model=CalibrationWorkOrder)
async def update_calibration_work_order(
    work_order_id: str,
    payload: CalibrationWorkOrderUpdate,
    current_user: User = Depends(get_current_user),
) -> CalibrationWorkOrder:
    record = await _get_work_order_or_404(work_order_id)
    work_order = CalibrationWorkOrder(**record)
    data = payload.dict(exclude_unset=True)
    status_note = _normalize_optional_str(data.pop("status_note", None))
    raw_status = data.pop("status", None)
    status_changed = False
    if "planned_date" in data and data["planned_date"] is not None:
        work_order.planned_date = _ensure_datetime(data["planned_date"])
    if "due_date" in data and data["due_date"] is not None:
        work_order.due_date = _ensure_datetime(data["due_date"])
    if "assigned_to" in data:
        work_order.assigned_to = _normalize_optional_str(data["assigned_to"])
    if "notes" in data:
        work_order.notes = _normalize_optional_str(data["notes"])
    if "linked_dof_task_id" in data:
        work_order.linked_dof_task_id = _normalize_optional_str(data["linked_dof_task_id"])
    if raw_status is not None:
        new_status = _normalize_work_order_status(raw_status)
        _validate_work_order_transition(work_order.status, new_status)
        status_changed = new_status != work_order.status
        work_order.status = new_status
        if new_status == "completed":
            work_order.completed_at = work_order.completed_at or datetime.now(timezone.utc)
            if work_order.result is None:
                if any(not record.pass_status for record in work_order.measurement_records):
                    work_order.result = "fail"
                else:
                    work_order.result = "pass"
        elif status_changed:
            work_order.completed_at = None
    if status_changed or status_note:
        _append_work_order_status_history(
            work_order,
            work_order.status,
            current_user.id,
            status_note,
        )
    work_order = await _persist_work_order_record(work_order)
    if work_order.status == "completed":
        device_doc = await _get_device_or_404(work_order.device_id)
        device_model = CalibrationDevice(**device_doc)
        if work_order.result != "fail":
            device_model.last_calibrated_at = work_order.completed_at or datetime.now(timezone.utc)
            device_model.next_due_date = _compute_next_due_date(
                device_model.last_calibrated_at,
                device_model.calibration_interval_days,
                device_model.last_calibrated_at,
            )
        await _persist_device_record(device_model)
    return work_order


@api_router.post("/calibration/work-orders/{work_order_id}/measurements", response_model=CalibrationWorkOrder)
async def add_work_order_measurement(
    work_order_id: str,
    payload: CalibrationMeasurementCreate,
    current_user: User = Depends(get_current_user),
) -> CalibrationWorkOrder:
    record = await _get_work_order_or_404(work_order_id)
    work_order = CalibrationWorkOrder(**record)
    measurement = CalibrationMeasurementRecord(
        parameter=_normalize_required_str(payload.parameter, "Measurement parameter"),
        nominal=float(payload.nominal),
        tolerance=float(payload.tolerance) if payload.tolerance is not None else None,
        observed=float(payload.observed),
        pass_status=bool(payload.pass_status),
        recorded_at=_ensure_datetime(payload.recorded_at) or datetime.now(timezone.utc),
        recorded_by=_normalize_optional_str(payload.recorded_by) or current_user.full_name or current_user.username,
        note=_normalize_optional_str(payload.note),
    )
    work_order.measurement_records.append(measurement)
    if not measurement.pass_status:
        work_order.result = "fail"
    elif work_order.result is None and all(rec.pass_status for rec in work_order.measurement_records):
        work_order.result = "pass"
    if work_order.status == "planned":
        _validate_work_order_transition(work_order.status, "in_progress")
        work_order.status = "in_progress"
        _append_work_order_status_history(
            work_order,
            work_order.status,
            current_user.id,
            "First measurement recorded",
        )
    work_order = await _persist_work_order_record(work_order)
    return work_order


@api_router.post("/calibration/work-orders/{work_order_id}/costs", response_model=CalibrationWorkOrder)
async def add_work_order_cost(
    work_order_id: str,
    payload: CalibrationCostCreate,
    current_user: User = Depends(get_current_user),
) -> CalibrationWorkOrder:
    record = await _get_work_order_or_404(work_order_id)
    work_order = CalibrationWorkOrder(**record)
    entry = CalibrationCostEntry(
        description=_normalize_required_str(payload.description, "Cost description"),
        amount=float(payload.amount),
        currency=payload.currency or "TRY",
        incurred_at=_ensure_datetime(payload.incurred_at) or datetime.now(timezone.utc),
        supplier=_normalize_optional_str(payload.supplier),
        reference=_normalize_optional_str(payload.reference),
        note=_normalize_optional_str(payload.note),
    )
    work_order.cost_entries.append(entry)
    work_order = await _persist_work_order_record(work_order)
    return work_order


@api_router.post("/calibration/work-orders/{work_order_id}/dof", response_model=CalibrationWorkOrder)
async def create_dof_for_calibration_work_order(
    work_order_id: str,
    payload: CalibrationDofRequest,
    current_user: User = Depends(get_current_user),
) -> CalibrationWorkOrder:
    work_order_doc = await _get_work_order_or_404(work_order_id)
    work_order = CalibrationWorkOrder(**work_order_doc)
    device_doc = await _get_device_or_404(work_order.device_id)
    device_model = CalibrationDevice(**device_doc)
    if work_order.linked_dof_task_id:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="DOF task already linked.")
    dof_task = await _create_dof_for_work_order(device_model, work_order, payload, current_user)
    work_order.linked_dof_task_id = dof_task.id
    work_order = await _persist_work_order_record(work_order)
    if dof_task.id not in device_model.linked_dof_task_ids:
        device_model.linked_dof_task_ids.append(dof_task.id)
        await _persist_device_record(device_model)
    return work_order


@api_router.get("/calibration/reports/summary", response_model=CalibrationReportSummary)
async def get_calibration_summary_report(
    current_user: User = Depends(get_current_user),
) -> CalibrationReportSummary:
    summary = await _collect_calibration_summary()
    return CalibrationReportSummary(
        total_devices=summary["total_devices"],
        active_devices=summary["active_devices"],
        overdue_devices=summary["overdue_devices"],
        upcoming_devices=summary["upcoming_devices"],
        total_work_orders=summary["total_work_orders"],
        open_work_orders=summary["open_work_orders"],
        cost_total=summary["cost_total"],
    )


# Reporting Routes
@api_router.get("/reports/modules", response_model=List[str])
async def list_report_modules(current_user: User = Depends(get_current_user)) -> List[str]:
    return sorted(REPORT_ALLOWED_MODULES)


@api_router.get("/reports/definitions", response_model=List[ReportDefinition])
async def list_report_definitions(current_user: User = Depends(get_current_user)) -> List[ReportDefinition]:
    cursor = db.report_definitions.find({}).sort("updated_at", -1)
    results: List[ReportDefinition] = []
    async for doc in cursor:
        results.append(ReportDefinition(**doc))
    return results


@api_router.post("/reports/definitions", response_model=ReportDefinition, status_code=status.HTTP_201_CREATED)
async def create_report_definition(
    payload: ReportDefinitionCreate,
    current_user: User = Depends(get_current_user),
) -> ReportDefinition:
    module = _normalize_report_module(payload.module)
    max_rows = payload.max_rows or 2000
    max_rows = max(1, min(int(max_rows), 20000))
    definition = ReportDefinition(
        name=_normalize_required_str(payload.name, "Report name"),
        module=module,
        description=_normalize_optional_str(payload.description),
        filters=payload.filters or [],
        visualizations=payload.visualizations or [],
        fields=_sanitize_string_list(payload.fields) or None,
        max_rows=max_rows,
        default_template_id=_normalize_optional_str(payload.default_template_id),
        created_by=current_user.id,
    )
    return await _persist_report_definition(definition)


@api_router.get("/reports/definitions/{definition_id}", response_model=ReportDefinition)
async def get_report_definition(
    definition_id: str,
    current_user: User = Depends(get_current_user),
) -> ReportDefinition:
    record = await _get_report_definition_or_404(definition_id)
    return ReportDefinition(**record)


@api_router.patch("/reports/definitions/{definition_id}", response_model=ReportDefinition)
async def update_report_definition(
    definition_id: str,
    payload: ReportDefinitionUpdate,
    current_user: User = Depends(get_current_user),
) -> ReportDefinition:
    record = await _get_report_definition_or_404(definition_id)
    definition = ReportDefinition(**record)
    data = payload.dict(exclude_unset=True)
    if "name" in data and data["name"] is not None:
        definition.name = _normalize_required_str(data["name"], "Report name")
    if "description" in data:
        definition.description = _normalize_optional_str(data["description"])
    if "filters" in data and data["filters"] is not None:
        definition.filters = data["filters"]
    if "visualizations" in data and data["visualizations"] is not None:
        definition.visualizations = data["visualizations"]
    if "fields" in data:
        sanitized_fields = _sanitize_string_list(data["fields"])
        definition.fields = sanitized_fields or None
    if "max_rows" in data and data["max_rows"] is not None:
        definition.max_rows = max(1, min(int(data["max_rows"]), 20000))
    if "default_template_id" in data:
        definition.default_template_id = _normalize_optional_str(data["default_template_id"])
    return await _persist_report_definition(definition)


@api_router.delete("/reports/definitions/{definition_id}", status_code=status.HTTP_204_NO_CONTENT)
async def delete_report_definition(
    definition_id: str,
    current_user: User = Depends(get_current_user),
) -> Response:
    await db.report_definitions.delete_one({"id": definition_id})
    return Response(status_code=status.HTTP_204_NO_CONTENT)


@api_router.post("/reports/run", response_model=ReportRunResult)
async def run_report(
    payload: ReportRunRequest,
    current_user: User = Depends(get_current_user),
) -> ReportRunResult:
    record = await _get_report_definition_or_404(payload.definition_id)
    definition = ReportDefinition(**record)
    result = await _execute_report_run(definition, payload.filters, payload.visualization_overrides)
    return result


@api_router.get("/reports/templates", response_model=List[ReportTemplate])
async def list_report_templates(current_user: User = Depends(get_current_user)) -> List[ReportTemplate]:
    cursor = db.report_templates.find({}).sort("updated_at", -1)
    templates: List[ReportTemplate] = []
    async for doc in cursor:
        templates.append(ReportTemplate(**doc))
    return templates


@api_router.post("/reports/templates", response_model=ReportTemplate, status_code=status.HTTP_201_CREATED)
async def create_report_template(
    payload: ReportTemplateCreate,
    current_user: User = Depends(get_current_user),
) -> ReportTemplate:
    module = _normalize_report_module(payload.module)
    template = ReportTemplate(
        name=_normalize_required_str(payload.name, "Template name"),
        module=module,
        template_type=payload.template_type,
        description=_normalize_optional_str(payload.description),
        file_upload_id=_normalize_optional_str(payload.file_upload_id),
        placeholders=payload.placeholders or [],
        default_definition_id=_normalize_optional_str(payload.default_definition_id),
        created_by=current_user.id,
    )
    return await _persist_report_template(template)


@api_router.get("/reports/templates/{template_id}", response_model=ReportTemplate)
async def get_report_template(
    template_id: str,
    current_user: User = Depends(get_current_user),
) -> ReportTemplate:
    record = await _get_report_template_or_404(template_id)
    return ReportTemplate(**record)


@api_router.patch("/reports/templates/{template_id}", response_model=ReportTemplate)
async def update_report_template(
    template_id: str,
    payload: ReportTemplateUpdate,
    current_user: User = Depends(get_current_user),
) -> ReportTemplate:
    record = await _get_report_template_or_404(template_id)
    template = ReportTemplate(**record)
    data = payload.dict(exclude_unset=True)
    if "name" in data and data["name"] is not None:
        template.name = _normalize_required_str(data["name"], "Template name")
    if "description" in data:
        template.description = _normalize_optional_str(data["description"])
    if "file_upload_id" in data:
        template.file_upload_id = _normalize_optional_str(data["file_upload_id"])
    if "template_type" in data and data["template_type"] is not None:
        template.template_type = data["template_type"]
    if "placeholders" in data and data["placeholders"] is not None:
        template.placeholders = data["placeholders"]
    if "default_definition_id" in data:
        template.default_definition_id = _normalize_optional_str(data["default_definition_id"])
    return await _persist_report_template(template)


@api_router.delete("/reports/templates/{template_id}", status_code=status.HTTP_204_NO_CONTENT)
async def delete_report_template(
    template_id: str,
    current_user: User = Depends(get_current_user),
) -> Response:
    await db.report_templates.delete_one({"id": template_id})
    return Response(status_code=status.HTTP_204_NO_CONTENT)


@api_router.post("/reports/templates/{template_id}/render", response_model=ReportTemplateRenderResponse)
async def render_report_template(
    template_id: str,
    payload: ReportTemplateRenderRequest,
    current_user: User = Depends(get_current_user),
) -> ReportTemplateRenderResponse:
    template_record = await _get_report_template_or_404(template_id)
    template = ReportTemplate(**template_record)
    definition_id = payload.definition_id or template.default_definition_id
    run_result: Optional[ReportRunResult] = None
    definition_model: Optional[ReportDefinition] = None
    if definition_id:
        definition_record = await _get_report_definition_or_404(definition_id)
        definition_model = ReportDefinition(**definition_record)
        run_result = await _execute_report_run(definition_model, payload.filters, None)
        if run_result:
            if not payload.include_rows:
                run_result.rows = []
            if not payload.include_summary:
                run_result.summary = {}
            if not payload.include_visualizations:
                run_result.visualizations = []
    return ReportTemplateRenderResponse(
        template=template,
        definition=definition_model,
        run_result=run_result,
        placeholders=template.placeholders,
    )


# Risk Routes
@api_router.get("/risk/settings", response_model=RiskSettings)
async def get_risk_settings(current_user: User = Depends(get_current_user)) -> RiskSettings:
    return await _load_risk_settings()


@api_router.put("/risk/settings", response_model=RiskSettings)
async def update_risk_settings(
    payload: RiskSettings,
    current_user: User = Depends(get_current_user),
) -> RiskSettings:
    payload.updated_by = current_user.id
    return await _persist_risk_settings(payload)


@api_router.get("/risk/report-template", response_model=RiskReportTemplate)
async def get_risk_report_template(current_user: User = Depends(get_current_user)) -> RiskReportTemplate:
    return await _load_risk_report_template()


@api_router.put("/risk/report-template", response_model=RiskReportTemplate)
async def update_risk_report_template(
    payload: RiskReportTemplate,
    current_user: User = Depends(get_current_user),
) -> RiskReportTemplate:
    payload.updated_by = current_user.id
    return await _persist_risk_report_template(payload)


@api_router.get("/risks", response_model=List[RiskAssessment])
async def list_risks(
    status_filter: Optional[str] = Query(default=None, alias="status"),
    category: Optional[str] = None,
    owner: Optional[str] = None,
    search: Optional[str] = None,
    current_user: User = Depends(get_current_user),
) -> List[RiskAssessment]:
    query: Dict[str, Any] = {}
    if status_filter:
        query["status"] = _normalize_risk_status(status_filter)
    if category:
        query["category"] = category.strip()
    if owner:
        query["owner"] = owner.strip()
    if search:
        pattern = {"$regex": re.escape(search), "$options": "i"}
        query["$or"] = [
            {"risk_code": pattern},
            {"title": pattern},
            {"process": pattern},
        ]
    cursor = db.risks.find(query).sort("updated_at", -1)
    return [RiskAssessment(**doc) async for doc in cursor]


@api_router.post("/risks", response_model=RiskAssessment, status_code=status.HTTP_201_CREATED)
async def create_risk_record(
    payload: RiskAssessmentCreate,
    current_user: User = Depends(get_current_user),
) -> RiskAssessment:
    settings = await _load_risk_settings()
    status_value = _normalize_risk_status(payload.status or "identified")
    factors = _sanitize_risk_factors(payload.custom_factors or [])
    risk_score = _calculate_risk_score(
        settings=settings,
        likelihood=payload.likelihood,
        impact=payload.impact,
        detection=payload.detection,
        controls_effectiveness=payload.controls_effectiveness,
        factors=factors,
    )
    risk = RiskAssessment(
        risk_code=await generate_risk_code(),
        title=_normalize_required_str(payload.title, "Risk title"),
        category=_normalize_required_str(payload.category, "Risk category"),
        process=_normalize_optional_str(payload.process),
        owner=_normalize_required_str(payload.owner, "Risk owner"),
        description=_normalize_optional_str(payload.description),
        status=status_value,
        likelihood=payload.likelihood,
        impact=payload.impact,
        detection=payload.detection,
        controls_effectiveness=payload.controls_effectiveness,
        custom_factors=factors,
        risk_score=risk_score,
        linked_capa_ids=_sanitize_string_list(payload.linked_capa_ids),
        linked_audit_finding_ids=_sanitize_string_list(payload.linked_audit_finding_ids),
        next_review_date=_ensure_timezone(payload.next_review_date) if payload.next_review_date else None,
        created_by=current_user.id,
    )
    risk = await _persist_risk_record(risk)
    await _record_risk_trend_point(risk)
    return risk


@api_router.get("/risks/{risk_id}", response_model=RiskAssessment)
async def get_risk(
    risk_id: str,
    include_trend: bool = False,
    current_user: User = Depends(get_current_user),
) -> RiskAssessment:
    risk_doc = await _get_risk_or_404(risk_id)
    risk_model = RiskAssessment(**risk_doc)
    if include_trend:
        risk_model.trend = await _fetch_risk_trend(risk_id)
    return risk_model


@api_router.patch("/risks/{risk_id}", response_model=RiskAssessment)
async def update_risk(
    risk_id: str,
    payload: RiskAssessmentUpdate,
    current_user: User = Depends(get_current_user),
) -> RiskAssessment:
    risk_doc = await _get_risk_or_404(risk_id)
    risk_model = RiskAssessment(**risk_doc)
    await _append_risk_revision(risk_model, current_user.id, payload.revision_note)

    data = payload.dict(exclude_unset=True)
    settings = await _load_risk_settings()
    if "title" in data:
        risk_model.title = _normalize_required_str(data["title"], "Risk title")
    if "category" in data:
        risk_model.category = _normalize_required_str(data["category"], "Risk category")
    if "process" in data:
        risk_model.process = _normalize_optional_str(data["process"])
    if "owner" in data:
        risk_model.owner = _normalize_required_str(data["owner"], "Risk owner")
    if "description" in data:
        risk_model.description = _normalize_optional_str(data["description"])
    if "status" in data:
        risk_model.status = _normalize_risk_status(data["status"])
    if "likelihood" in data:
        risk_model.likelihood = float(data["likelihood"])
    if "impact" in data:
        risk_model.impact = float(data["impact"])
    if "detection" in data:
        risk_model.detection = float(data["detection"]) if data["detection"] is not None else None
    if "controls_effectiveness" in data:
        risk_model.controls_effectiveness = (
            float(data["controls_effectiveness"]) if data["controls_effectiveness"] is not None else None
        )
    if "custom_factors" in data and data["custom_factors"] is not None:
        risk_model.custom_factors = _sanitize_risk_factors(data["custom_factors"])
    if "linked_capa_ids" in data and data["linked_capa_ids"] is not None:
        risk_model.linked_capa_ids = _sanitize_string_list(data["linked_capa_ids"])
    if "linked_audit_finding_ids" in data and data["linked_audit_finding_ids"] is not None:
        risk_model.linked_audit_finding_ids = _sanitize_string_list(data["linked_audit_finding_ids"])
    if "next_review_date" in data:
        risk_model.next_review_date = (
            _ensure_timezone(data["next_review_date"]) if data["next_review_date"] else None
        )

    risk_model.risk_score = _calculate_risk_score(
        settings,
        likelihood=risk_model.likelihood,
        impact=risk_model.impact,
        detection=risk_model.detection,
        controls_effectiveness=risk_model.controls_effectiveness,
        factors=risk_model.custom_factors,
    )
    risk_model.last_reviewed_at = datetime.now(timezone.utc)
    risk_model = await _persist_risk_record(risk_model)
    await _record_risk_trend_point(risk_model)
    return risk_model


@api_router.get("/risks/{risk_id}/revisions", response_model=List[RiskRevision])
async def list_risk_revisions(
    risk_id: str,
    current_user: User = Depends(get_current_user),
) -> List[RiskRevision]:
    risk_doc = await _get_risk_or_404(risk_id)
    risk_model = RiskAssessment(**risk_doc)
    return sorted(risk_model.revisions, key=lambda rev: rev.revision_no)


@api_router.get("/risks/{risk_id}/compare")
async def compare_risk_revisions(
    risk_id: str,
    rev_a: int = Query(..., description="Base revision number"),
    rev_b: int = Query(..., description="Target revision number"),
    current_user: User = Depends(get_current_user),
) -> Dict[str, Any]:
    risk_doc = await _get_risk_or_404(risk_id)
    risk_model = RiskAssessment(**risk_doc)
    revisions = {rev.revision_no: rev for rev in risk_model.revisions}
    if rev_a not in revisions or rev_b not in revisions:
        raise HTTPException(status_code=404, detail="Revision not found.")
    diff = _build_risk_diff(revisions[rev_a].snapshot, revisions[rev_b].snapshot)
    return {
        "risk_id": risk_id,
        "base_revision": rev_a,
        "target_revision": rev_b,
        "diff": diff,
    }


@api_router.get("/risks/matrix", response_model=RiskMatrixSummary)
async def get_risk_matrix_summary(current_user: User = Depends(get_current_user)) -> RiskMatrixSummary:
    settings = await _load_risk_settings()
    summary = await _build_risk_matrix_summary(settings)
    return RiskMatrixSummary(**summary)


@api_router.get("/risks/trends", response_model=RiskTrendResponse)
async def get_risk_trends(
    risk_id: Optional[str] = None,
    limit: int = 20,
    current_user: User = Depends(get_current_user),
) -> RiskTrendResponse:
    if risk_id:
        points = await _fetch_risk_trend(risk_id, limit)
    else:
        cursor = (
            db.risk_trends.find({}).sort("recorded_at", -1).limit(limit)
        )
        points = [
            RiskTrendPoint(
                recorded_at=item["recorded_at"],
                inherent_score=item["inherent_score"],
                residual_score=item["residual_score"],
                status=item.get("status", "identified"),
            )
            async for item in cursor
        ]
        points.reverse()
    return RiskTrendResponse(points=points)


@api_router.get("/risks/reports/custom")
async def generate_risk_report(
    risk_id: str,
    current_user: User = Depends(get_current_user),
) -> Dict[str, Any]:
    risk_doc = await _get_risk_or_404(risk_id)
    risk_model = RiskAssessment(**risk_doc)
    template = await _load_risk_report_template()
    context = {
        "title": risk_model.title,
        "status": risk_model.status,
        "inherent_score": risk_model.risk_score.inherent,
        "residual_score": risk_model.risk_score.residual,
        "owner": risk_model.owner,
        "controls": risk_model.controls_effectiveness,
        "category": risk_model.category,
        "risk_code": risk_model.risk_code,
    }
    body = template.body
    for key, value in context.items():
        body = body.replace(f"{{{{{key}}}}}", str(value if value is not None else "-"))
    return {
        "template_name": template.name,
        "generated_at": datetime.now(timezone.utc),
        "risk_id": risk_model.id,
        "title": template.name,
        "body": body,
    }
# D�F (Corrective & Preventive Task) Routes
ALLOWED_DOF_STATUSES = {"open", "in_progress", "pending_closure", "closed", "cancelled"}
DEFAULT_DOF_DEPARTMENTS = [
    "Resepsiyon",
    "Ön Büro",
    "Kat Hizmetleri",
    "Housekeeping",
    "Yiyecek & İçecek",
    "Mutfak",
    "Satış ve Pazarlama",
    "Rezervasyon",
    "Misafir İlişkileri",
    "İnsan Kaynakları",
    "Muhasebe",
    "Finans",
    "Satın Alma",
    "Teknik Servis",
    "Bakım",
    "Güvenlik",
    "Bilgi Teknolojileri",
    "Spa ve Wellness",
    "Eğlence ve Animasyon",
    "Kalite Yönetimi",
]

def _build_dof_filters(
    department: Optional[str],
    status: Optional[str],
    search: Optional[str],
    start_date: Optional[str],
    end_date: Optional[str],
) -> Dict[str, Any]:
    filters: Dict[str, Any] = {}

    if department:
        filters["department"] = department
    if status:
        filters["status"] = status
    if start_date or end_date:
        date_filter: Dict[str, Any] = {}
        if start_date:
            try:
                date_filter["$gte"] = _parse_iso_datetime(start_date)
            except ValueError:
                raise HTTPException(status_code=400, detail="Invalid start_date format. Use ISO 8601.")
        if end_date:
            try:
                date_filter["$lte"] = _parse_iso_datetime(end_date)
            except ValueError:
                raise HTTPException(status_code=400, detail="Invalid end_date format. Use ISO 8601.")
        if date_filter:
            filters["created_at"] = date_filter

    if search:
        regex = {"$regex": re.escape(search), "$options": "i"}
        filters["$or"] = [
            {"title": regex},
            {"description": regex},
            {"responsible_person": regex},
            {"dof_no": regex},
        ]

    return filters

@api_router.get("/dof-tasks", response_model=DofTaskListResponse)
async def list_dof_tasks(
    department: Optional[str] = None,
    status: Optional[str] = None,
    search: Optional[str] = None,
    start_date: Optional[str] = None,
    end_date: Optional[str] = None,
    page: int = 1,
    page_size: int = 20,
    current_user: User = Depends(get_current_user),
):
    if page < 1 or page_size < 1:
        raise HTTPException(status_code=400, detail="page and page_size must be positive integers.")

    filters = _build_dof_filters(department, status, search, start_date, end_date)

    skip = (page - 1) * page_size
    total = await db.dof_tasks.count_documents(filters)
    cursor = (
        db.dof_tasks.find(filters)
        .sort("updated_at", -1)
        .skip(skip)
        .limit(page_size)
    )
    tasks = [DofTask(**task) async for task in cursor]

    return DofTaskListResponse(
        items=tasks,
        total=total,
        page=page,
        page_size=page_size,
    )

@api_router.post("/dof-tasks", response_model=DofTask, status_code=status.HTTP_201_CREATED)
async def create_dof_task(
    dof_data: DofTaskCreate,
    current_user: User = Depends(get_current_user),
):
    dof_no = await generate_dof_no()
    now = datetime.now(timezone.utc)
    initial_report_date = dof_data.initial_improvement_report_date
    if initial_report_date and initial_report_date.tzinfo is None:
        initial_report_date = initial_report_date.replace(tzinfo=timezone.utc)
    team_members = _sanitize_string_list(dof_data.team_members or [])

    task = DofTask(
        dof_no=dof_no,
        title=dof_data.title,
        description=dof_data.description,
        department=dof_data.department,
        responsible_person=dof_data.responsible_person,
        due_date=dof_data.due_date,
        team_members=team_members,
        initial_improvement_report_date=initial_report_date,
        created_by=current_user.id,
        status_history=[
            DofStatusHistory(status="open", changed_by=current_user.id, changed_at=now)
        ],
    )

    await db.dof_tasks.insert_one(task.dict())
    return task

@api_router.get("/dof-tasks/{dof_id}", response_model=DofTask)
async def get_dof_task(
    dof_id: str,
    current_user: User = Depends(get_current_user),
):
    task = await db.dof_tasks.find_one({"id": dof_id})
    if not task:
        raise HTTPException(status_code=404, detail="DÖF kaydı bulunamadı.")
    return DofTask(**task)

@api_router.put("/dof-tasks/{dof_id}", response_model=DofTask)
async def update_dof_task(
    dof_id: str,
    update_data: DofTaskUpdate,
    current_user: User = Depends(get_current_user),
):
    task = await db.dof_tasks.find_one({"id": dof_id})
    if not task:
        raise HTTPException(status_code=404, detail="DÖF kaydı bulunamadı.")

    update_fields = {k: v for k, v in update_data.dict(exclude_unset=True).items()}
    if not update_fields:
        return DofTask(**task)

    if "initial_improvement_report_date" in update_fields:
        report_date = update_fields["initial_improvement_report_date"]
        if report_date and report_date.tzinfo is None:
            report_date = report_date.replace(tzinfo=timezone.utc)
        update_fields["initial_improvement_report_date"] = report_date
    if "team_members" in update_fields:
        update_fields["team_members"] = _sanitize_string_list(update_fields["team_members"])

    update_fields["updated_at"] = datetime.now(timezone.utc)

    await db.dof_tasks.update_one({"id": dof_id}, {"$set": update_fields})
    updated = await db.dof_tasks.find_one({"id": dof_id})
    return DofTask(**updated)

@api_router.post("/dof-tasks/{dof_id}/closure/request", response_model=DofTask)
async def request_dof_closure(
    dof_id: str,
    payload: DofClosureRequest,
    current_user: User = Depends(get_current_user),
) -> DofTask:
    task = await db.dof_tasks.find_one({"id": dof_id})
    if not task:
        raise HTTPException(status_code=404, detail="DÖF kaydı bulunamadı.")
    if task.get("status") in {"closed", "cancelled"}:
        raise HTTPException(status_code=400, detail="Bu kayıt için kapanış isteği gönderilemez.")

    now = datetime.now(timezone.utc)
    status_entry = DofStatusHistory(
        status="pending_closure",
        changed_by=current_user.id,
        changed_at=now,
        note=payload.note,
    )
    update_doc: Dict[str, Any] = {
        "status": "pending_closure",
        "updated_at": now,
        "closure_requested_at": now,
        "closure_requested_by": current_user.id,
        "closure_request_note": payload.note,
        "closure_approved_at": None,
        "closure_approved_by": None,
        "closure_decision_note": None,
    }

    await db.dof_tasks.update_one(
        {"id": dof_id},
        {
            "$set": update_doc,
            "$push": {"status_history": status_entry.dict()},
        },
    )
    updated = await db.dof_tasks.find_one({"id": dof_id})
    return DofTask(**updated)


@api_router.post("/dof-tasks/{dof_id}/closure/decision", response_model=DofTask)
async def decide_dof_closure(
    dof_id: str,
    payload: DofClosureDecision,
    current_user: User = Depends(get_current_user),
) -> DofTask:
    task = await db.dof_tasks.find_one({"id": dof_id})
    if not task:
        raise HTTPException(status_code=404, detail="DÖF kaydı bulunamadı.")
    if task.get("status") != "pending_closure":
        raise HTTPException(status_code=400, detail="Onay bekleyen kapanış isteği bulunamadı.")

    now = datetime.now(timezone.utc)
    if payload.approve:
        new_status = "closed"
        update_doc = {
            "status": new_status,
            "updated_at": now,
            "closed_at": now,
            "closure_approved_at": now,
            "closure_approved_by": current_user.id,
            "closure_decision_note": payload.note,
        }
    else:
        new_status = "in_progress"
        update_doc = {
            "status": new_status,
            "updated_at": now,
            "closure_requested_at": None,
            "closure_requested_by": None,
            "closure_request_note": None,
            "closure_approved_at": None,
            "closure_approved_by": None,
            "closure_decision_note": payload.note,
            "closed_at": None,
        }

    status_entry = DofStatusHistory(
        status=new_status,
        changed_by=current_user.id,
        changed_at=now,
        note=payload.note,
    )

    await db.dof_tasks.update_one(
        {"id": dof_id},
        {
            "$set": update_doc,
            "$push": {"status_history": status_entry.dict()},
        },
    )
    updated = await db.dof_tasks.find_one({"id": dof_id})
    return DofTask(**updated)

@api_router.patch("/dof-tasks/{dof_id}/status", response_model=DofTask)
async def update_dof_status(
    dof_id: str,
    status_update: DofStatusUpdate,
    current_user: User = Depends(get_current_user),
):
    if status_update.status not in ALLOWED_DOF_STATUSES:
        raise HTTPException(
            status_code=400,
            detail=f"Geçersiz durum. Kullanılabilir durumlar: {', '.join(sorted(ALLOWED_DOF_STATUSES))}",
        )

    task = await db.dof_tasks.find_one({"id": dof_id})
    if not task:
        raise HTTPException(status_code=404, detail="DÖF kaydı bulunamadı.")

    now = datetime.now(timezone.utc)
    update_doc: Dict[str, Any] = {
        "status": status_update.status,
        "updated_at": now,
    }

    if status_update.status == "pending_closure":
        update_doc["closure_requested_at"] = now
        update_doc["closure_requested_by"] = current_user.id
        update_doc["closure_request_note"] = status_update.note
        update_doc["closure_approved_at"] = None
        update_doc["closure_approved_by"] = None
        update_doc["closure_decision_note"] = None
    elif status_update.status == "closed":
        update_doc["closed_at"] = now
        update_doc["closure_approved_at"] = now
        update_doc["closure_approved_by"] = current_user.id
        update_doc["closure_decision_note"] = status_update.note
    else:
        if task.get("closed_at"):
            update_doc["closed_at"] = None
        update_doc["closure_requested_at"] = None
        update_doc["closure_requested_by"] = None
        update_doc["closure_request_note"] = None
        update_doc["closure_approved_at"] = None
        update_doc["closure_approved_by"] = None
        update_doc["closure_decision_note"] = None

    status_entry = DofStatusHistory(
        status=status_update.status,
        changed_by=current_user.id,
        changed_at=now,
        note=status_update.note,
    )

    await db.dof_tasks.update_one(
        {"id": dof_id},
        {
            "$set": update_doc,
            "$push": {"status_history": status_entry.dict()},
        },
    )

    updated = await db.dof_tasks.find_one({"id": dof_id})
    return DofTask(**updated)

@api_router.get("/dof-tasks/report/summary")
async def get_dof_summary_report(
    department: Optional[str] = None,
    status: Optional[str] = None,
    start_date: Optional[str] = None,
    end_date: Optional[str] = None,
    current_user: User = Depends(get_current_user),
):
    filters = _build_dof_filters(department, status, None, start_date, end_date)
    now = datetime.now(timezone.utc)

    total = await db.dof_tasks.count_documents(filters)

    status_pipeline = [
        {"$match": filters},
        {"$group": {"_id": "$status", "count": {"$sum": 1}}},
    ]
    status_counts_raw = await db.dof_tasks.aggregate(status_pipeline).to_list(None)
    status_counts = {item["_id"]: item["count"] for item in status_counts_raw}

    department_pipeline = [
        {"$match": filters},
        {"$group": {"_id": "$department", "count": {"$sum": 1}}},
        {"$sort": {"count": -1}},
    ]
    department_counts = await db.dof_tasks.aggregate(department_pipeline).to_list(None)

    overdue_filter = dict(filters)
    if "status" not in overdue_filter:
        overdue_filter["status"] = {"$in": ["open", "in_progress", "pending_closure"]}
    overdue_filter["due_date"] = {"$lt": now}
    overdue = await db.dof_tasks.count_documents(overdue_filter)

    upcoming_filter = dict(filters)
    if "status" not in upcoming_filter:
        upcoming_filter["status"] = {"$in": ["open", "in_progress", "pending_closure"]}
    upcoming_filter["due_date"] = {"$gte": now}
    upcoming_cursor = (
        db.dof_tasks.find(upcoming_filter)
        .sort("due_date", 1)
        .limit(5)
    )
    upcoming_deadlines = [
        {
            "id": item["id"],
            "dof_no": item["dof_no"],
            "title": item["title"],
            "due_date": item["due_date"],
            "department": item["department"],
            "responsible_person": item["responsible_person"],
        }
        async for item in upcoming_cursor
    ]

    monthly_pipeline = [
        {"$match": filters},
        {
            "$group": {
                "_id": {
                    "year": {"$year": "$created_at"},
                    "month": {"$month": "$created_at"},
                },
                "count": {"$sum": 1},
            }
        },
        {"$sort": {"_id.year": 1, "_id.month": 1}},
    ]
    monthly_raw = await db.dof_tasks.aggregate(monthly_pipeline).to_list(None)
    monthly_trends = [
        {
            "label": f"{item['_id']['year']}-{item['_id']['month']:02d}",
            "count": item["count"],
        }
        for item in monthly_raw
    ]

    return {
        "total": total,
        "status_counts": status_counts,
        "department_counts": department_counts,
        "overdue": overdue,
        "upcoming_deadlines": upcoming_deadlines,
        "monthly_trends": monthly_trends,
    }

@api_router.get("/dof-tasks/departments", response_model=List[str])
async def get_dof_departments(
    current_user: User = Depends(get_current_user),
):
    user_departments = await db.users.distinct("department")
    task_departments = await db.dof_tasks.distinct("department")
    unique_departments = sorted(
        {
            dep
            for dep in [
                *DEFAULT_DOF_DEPARTMENTS,
                *user_departments,
                *task_departments,
            ]
            if dep
        }
    )
    return unique_departments

# Notification Routes
@api_router.get("/notifications", response_model=List[Notification])
async def get_notifications(current_user: User = Depends(get_current_user)):
    notifications = await db.notifications.find({"user_id": current_user.id}).sort("created_at", -1).to_list(50)
    return [Notification(**notif) for notif in notifications]

@api_router.put("/notifications/{notification_id}/read")
async def mark_notification_read(notification_id: str, current_user: User = Depends(get_current_user)):
    await db.notifications.update_one(
        {"id": notification_id, "user_id": current_user.id},
        {"$set": {"is_read": True}}
    )
    return {"message": "Notification marked as read"}

# Email Route
@api_router.post("/send-email")
async def send_email(
    recipients: List[EmailStr],
    subject: str,
    message: str,
    background_tasks: BackgroundTasks,
    current_user: User = Depends(get_current_user)
):
    background_tasks.add_task(
        EmailService.send_email,
        recipients=recipients,
        subject=subject,
        body=message,
        html_body=f"<p>{message}</p>"
    )
    return {"message": "Email queued for delivery"}

# Initialize default admin user
@api_router.post("/init/admin")
async def init_admin():
    admin_exists = await db.users.find_one({"username": "admin"})
    if admin_exists:
        return {"message": "Admin user already exists"}
    
    admin_data = {
        "id": str(uuid.uuid4()),
        "username": "admin",
        "email": "admin@calista.com.tr",
        "password": hash_password("admin123"),
        "full_name": "System Administrator",
        "role": "Admin",
        "roles": ["Admin"],
        "permissions": ["*"],
        "department": "IT",
        "is_active": True,
        "created_at": datetime.now(timezone.utc)
    }
    
    await db.users.insert_one(admin_data)
    return {"message": "Admin user created successfully", "username": "admin", "password": "admin123"}

# Include router
app.include_router(api_router)

app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=os.environ.get('CORS_ORIGINS', '*').split(','),
    allow_methods=["*"],
    allow_headers=["*"],
)

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

@app.on_event("startup")
async def startup_event():
    await _ensure_report_indexes()
    await _seed_report_assets()

@app.on_event("shutdown")
async def shutdown_db_client():
    client.close()






